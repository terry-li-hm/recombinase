"""Inspect a pptx template and print structural metadata.

Prints only structural information (shape names, types, placeholder info,
text character counts). Does NOT print actual text content, so it's safe
to share the output publicly even if the source template contains draft
sample data.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


@dataclass
class ShapeInfo:
    """Structural summary of a single shape (no text content)."""

    name: str
    shape_type: str
    is_placeholder: bool
    placeholder_type: str | None
    placeholder_idx: int | None
    text_chars: int | None
    paragraph_count: int | None
    run_count: int | None
    depth: int = 0
    """Nesting depth inside group shapes. 0 = top-level, 1 = inside a group, etc."""

    preset_geom: str | None = None
    """The `<a:prstGeom prst="...">` value if the shape uses a preset geometry
    (`ellipse`, `roundRect`, `rect`, etc.). `None` if the shape has no preset
    geometry or uses a custom geometry. Useful for detecting circle-cropped
    profile pictures — if a Picture shape has `preset_geom="ellipse"` the
    underlying image bitmap is still a square, it's just being rendered with
    a circle mask."""

    has_table: bool = False
    """True if the shape is a table shape (has a `shape.table` attribute).
    Used by `recombinase validate` to catch config errors where a text
    placeholder is pointed at a table shape or vice versa."""

    table_cells: list[CellInfo] | None = None
    """For table shapes, per-cell structural metadata: (row, col) with
    paragraph/run/br counts for the first paragraph. Used to diagnose
    multi-run-br cell idioms (CV "Role\\n(duration)" pattern) — cells
    matching that pattern must have paras=1, runs>=2, brs>=1 on cell
    (row>=1, col=0). `None` for non-table shapes."""


@dataclass
class CellInfo:
    """Structural summary of a single table cell (no text content)."""

    row: int
    col: int
    paragraph_count: int
    """Total number of `<a:p>` paragraphs in the cell's text frame."""

    first_paragraph_run_count: int
    """Number of `<a:r>` runs in the cell's FIRST paragraph — what
    `_write_multirun_br` inspects. A value >= 2 is one of the three
    conditions for the dual-run-br idiom."""

    first_paragraph_br_count: int
    """Number of `<a:br/>` soft-break elements in the cell's first
    paragraph. A value >= 1 is the third condition for the idiom."""

    text_chars: int
    """Total character count across the cell's text."""

    @property
    def is_multirun_br(self) -> bool:
        """True if the first paragraph matches the dual-run-br idiom:
        single paragraph containing ≥2 runs separated by ≥1 soft break."""
        return (
            self.paragraph_count == 1
            and self.first_paragraph_run_count >= 2
            and self.first_paragraph_br_count >= 1
        )


@dataclass
class SlideInfo:
    """Structural summary of a single slide."""

    index: int
    layout_name: str
    shapes: list[ShapeInfo]


@dataclass
class TemplateInfo:
    """Structural summary of an entire pptx template."""

    path: Path
    slide_count: int
    slides: list[SlideInfo]
    layout_names: list[str]


def _detect_preset_geom(shape: Any) -> str | None:
    """Return the value of `<a:prstGeom prst="...">` for this shape, or None.

    Walks the shape's XML element looking for the first `prstGeom` descendant.
    Returns the preset name string (e.g. `"ellipse"`, `"roundRect"`, `"rect"`).
    Used by the inspector to flag circle-cropped profile pictures so users
    can see at a glance which shapes apply a mask to the underlying bitmap.
    """
    element = getattr(shape, "_element", None)
    if element is None:
        return None
    prst_geom = element.find(f".//{qn('a:prstGeom')}")
    if prst_geom is None:
        return None
    return prst_geom.get("prst")


def _cell_info(cell: Any, row: int, col: int) -> CellInfo:
    """Capture structural metadata for a single table cell."""
    paragraphs = list(cell.text_frame.paragraphs)
    first_runs = 0
    first_brs = 0
    if paragraphs:
        first_p_xml = paragraphs[0]._p
        first_runs = len(first_p_xml.findall(qn("a:r")))
        first_brs = len(first_p_xml.findall(qn("a:br")))
    return CellInfo(
        row=row,
        col=col,
        paragraph_count=len(paragraphs),
        first_paragraph_run_count=first_runs,
        first_paragraph_br_count=first_brs,
        text_chars=len(cell.text_frame.text or ""),
    )


def _collect_table_cells(shape: Any) -> list[CellInfo]:
    """Walk a table shape's cells and collect per-cell structural info.

    Skips spanned (merged non-origin) cells since they have no accessible
    text frame. Used to diagnose multi-run-br cell idioms in CV templates.
    """
    result: list[CellInfo] = []
    if not getattr(shape, "has_table", False):
        return result
    table = shape.table
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            if getattr(cell, "is_spanned", False):
                continue
            result.append(_cell_info(cell, row_index, col_index))
    return result


def _shape_info(shape: Any, depth: int = 0) -> ShapeInfo:
    placeholder_type: str | None = None
    placeholder_idx: int | None = None
    if getattr(shape, "is_placeholder", False):
        ph = shape.placeholder_format
        placeholder_type = str(ph.type) if ph.type is not None else None
        placeholder_idx = ph.idx

    text_chars: int | None = None
    paragraph_count: int | None = None
    run_count: int | None = None
    if getattr(shape, "has_text_frame", False):
        tf = shape.text_frame
        text_chars = len(tf.text or "")
        paragraph_count = len(tf.paragraphs)
        run_count = sum(len(p.runs) for p in tf.paragraphs)

    has_table = bool(getattr(shape, "has_table", False))
    table_cells = _collect_table_cells(shape) if has_table else None

    return ShapeInfo(
        name=shape.name,
        shape_type=str(shape.shape_type),
        is_placeholder=bool(getattr(shape, "is_placeholder", False)),
        placeholder_type=placeholder_type,
        placeholder_idx=placeholder_idx,
        text_chars=text_chars,
        paragraph_count=paragraph_count,
        run_count=run_count,
        depth=depth,
        preset_geom=_detect_preset_geom(shape),
        has_table=has_table,
        table_cells=table_cells,
    )


def _collect_shapes(shapes: Any, depth: int = 0) -> list[ShapeInfo]:
    """Walk a shape collection recursively, descending into groups.

    Each shape is emitted once; group shapes are emitted before their
    children, with `depth` incremented for nested shapes so formatters
    can indent them.
    """
    result: list[ShapeInfo] = []
    for shape in shapes:
        result.append(_shape_info(shape, depth))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.extend(_collect_shapes(shape.shapes, depth + 1))
    return result


def inspect_template(path: Path | str) -> TemplateInfo:
    """Inspect a pptx/pptm template and return its structural metadata."""
    path = Path(path).expanduser().resolve()
    prs = Presentation(str(path))

    slides: list[SlideInfo] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        shape_infos = _collect_shapes(slide.shapes)
        slides.append(
            SlideInfo(
                index=slide_index,
                layout_name=slide.slide_layout.name,
                shapes=shape_infos,
            )
        )

    return TemplateInfo(
        path=path,
        slide_count=len(prs.slides),
        slides=slides,
        layout_names=[layout.name for layout in prs.slide_layouts],
    )


def format_template_info(info: TemplateInfo) -> str:
    """Format a TemplateInfo for human-readable printing."""
    lines: list[str] = [
        f"File: {info.path}",
        f"Slide count: {info.slide_count}",
        "",
    ]

    for slide in info.slides:
        lines.append(f"=== Slide {slide.index} (layout: {slide.layout_name!r}) ===")
        if not slide.shapes:
            lines.append("  (no shapes)")
        for shape in slide.shapes:
            indent = "  " + ("  " * shape.depth)
            bits = [repr(shape.name), f"type={shape.shape_type}"]
            if shape.is_placeholder:
                bits.append(
                    f"placeholder(type={shape.placeholder_type}, idx={shape.placeholder_idx})"
                )
            if shape.text_chars is not None:
                bits.append(f"text_chars={shape.text_chars}")
                bits.append(f"paras={shape.paragraph_count}, runs={shape.run_count}")
            if shape.preset_geom is not None:
                bits.append(f"geom={shape.preset_geom}")
            if shape.depth > 0:
                bits.append(f"depth={shape.depth}")
            lines.append(indent + "- " + " | ".join(bits))
            # Drill into table cells: structural metadata per cell so the
            # user can see which cells match the multi-run-br idiom that
            # `populate_table` preserves on write.
            if shape.table_cells:
                cell_indent = indent + "    "
                lines.append(
                    cell_indent
                    + "cells (row,col) paras / first-para runs+brs / chars / dual-run-br:"
                )
                for cell in shape.table_cells:
                    marker = " ★" if cell.is_multirun_br else ""
                    lines.append(
                        cell_indent
                        + f"  ({cell.row},{cell.col}) "
                        + f"paras={cell.paragraph_count} "
                        + f"runs={cell.first_paragraph_run_count} "
                        + f"brs={cell.first_paragraph_br_count} "
                        + f"chars={cell.text_chars}{marker}"
                    )
        lines.append("")

    lines.append("=== Slide Layouts available on master ===")
    for layout_index, layout_name in enumerate(info.layout_names):
        lines.append(f"  [{layout_index}] {layout_name!r}")

    return "\n".join(lines)


def shape_names_from_slide(info: TemplateInfo, slide_index: int) -> list[str]:
    """Return the list of shape names on a specific slide (1-based index).

    Includes shapes nested inside groups so scaffold configs written from
    this list cover every addressable shape on the slide.
    """
    for slide in info.slides:
        if slide.index == slide_index:
            return [shape.name for shape in slide.shapes]
    return []


def shape_types_from_slide(info: TemplateInfo, slide_index: int) -> dict[str, bool]:
    """Return a dict mapping shape name -> is_table for a specific slide.

    Used by ``recombinase validate`` to check that placeholders aren't
    pointed at table shapes and vice versa. Later shapes with the same
    name win — normally shapes have unique names on a slide anyway.
    """
    for slide in info.slides:
        if slide.index == slide_index:
            return {shape.name: shape.has_table for shape in slide.shapes}
    return {}
