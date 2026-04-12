"""Generate a populated pptx deck from a template + structured YAML data.

Pattern: duplicate a well-formatted source slide in the template once per
record, populate the duplicated slide's shapes by name, save the output.
Relies on the `source slide` in the template being already styled correctly
(the 'filled example' approach) so that visual fidelity is inherited for
free via slide duplication.
"""

from __future__ import annotations

import warnings
from collections.abc import Iterator
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Any

import yaml
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.slide import Slide

from recombinase.config import (
    SectionConfig,
    TableConfig,
    TemplateConfig,
    _check_duplicate_yaml_keys,
)

# OOXML namespace used by r:id / r:embed / r:link attributes inside shape XML.
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def load_records(
    data_dir: Path | str,
    sort_by: str | None = None,
) -> list[dict[str, Any]]:
    """Load all YAML files from a directory as a list of records.

    Each .yaml or .yml file becomes one record. Files are loaded in
    sorted filename order by default.

    If *sort_by* is given, records are re-sorted by the value of that
    field in each record after loading. Numeric values sort numerically;
    string values sort lexicographically. Records missing the field
    sort last. Records with equal sort values preserve their filename
    order (stable sort).

    Each loaded record has `_recombinase_record_dir` stamped with the
    absolute path to its containing directory, so downstream code (in
    particular `set_picture`) can resolve relative image paths against the
    YAML file's directory rather than against the caller's CWD.
    """
    data_dir = Path(data_dir).expanduser().resolve()
    if not data_dir.exists():
        raise FileNotFoundError(f"Data directory not found: {data_dir}")
    if not data_dir.is_dir():
        raise NotADirectoryError(f"Not a directory: {data_dir}")

    records: list[dict[str, Any]] = []
    yaml_files = sorted([*data_dir.glob("*.yaml"), *data_dir.glob("*.yml")])
    for yaml_file in yaml_files:
        file_size = yaml_file.stat().st_size
        if file_size > 10 * 1024 * 1024:
            raise ValueError(
                f"{yaml_file}: file is {file_size / 1024 / 1024:.1f} MB, exceeding the 10 MB safety limit"
            )
        with yaml_file.open("r", encoding="utf-8") as fh:
            content = fh.read()
        _check_duplicate_yaml_keys(yaml_file, content)
        data = yaml.safe_load(content)
        if data is None:
            warnings.warn(f"{yaml_file}: YAML file is empty; skipping", stacklevel=2)
            continue
        if not isinstance(data, dict):
            raise ValueError(f"{yaml_file}: expected top-level mapping, got {type(data).__name__}")
        # Stamp the record's source directory so set_picture can resolve
        # relative image paths against the YAML file's own directory.
        data.setdefault("_recombinase_record_dir", str(yaml_file.parent))
        records.append(data)

    if sort_by is not None and records:
        # Sentinel that sorts after any real value.
        _missing = (1, "")  # (1, ...) sorts after (0, ...) below

        def _sort_key(record: dict[str, Any]) -> tuple[int, int | float | str]:
            value = record.get(sort_by)
            if value is None:
                return _missing
            if isinstance(value, (int, float)) and not isinstance(value, bool):
                return (0, value)
            return (0, str(value))

        records.sort(key=_sort_key)

    return records


def _walk_shapes(shapes: Any) -> Iterator[Any]:
    """Walk a shape collection, yielding every shape including group members.

    PowerPoint allows shapes to be nested inside GroupShape elements
    (`<p:grpSp>`), and iterating `slide.shapes` directly only yields the
    top-level children. This walker recurses into groups so every named
    shape is reachable by `find_shape_by_name` and visible to `inspect`.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)


def duplicate_slide(presentation: Any, source_slide: Slide) -> Slide:
    """Duplicate a slide, preserving shapes, formatting, and relationships.

    python-pptx has no native `slide.duplicate()` — this is the canonical
    workaround extended to also copy the source slide's relationships and
    rewrite `r:id` / `r:embed` / `r:link` references in the copied XML.
    Without the rel copy, deep-copied shapes referencing pictures, hyperlinks,
    or embedded charts produce dangling references on the new slide and
    those elements render as broken or empty.

    Ref: https://github.com/scanny/python-pptx/issues/132
    """
    blank_layout = source_slide.slide_layout
    new_slide = presentation.slides.add_slide(blank_layout)

    # Remove any shapes the layout added automatically — we want only the
    # source slide's shapes, not the layout's defaults.
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy relationships from the source slide onto the new slide, building an
    # old-rId -> new-rId map so that copied shape XML can be rewritten. Skip
    # notesSlide relationships — those belong to the notes subsystem and
    # copying them would attach the source's notes to the new slide.
    rel_id_map: dict[str, str] = {}
    for old_rid, rel in source_slide.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
        rel_id_map[old_rid] = new_rid

    # Copy each top-level shape XML element from the source slide into the new
    # slide, rewriting any r: references so they resolve against the new rels.
    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        for descendant in new_el.iter():
            for attr in ("id", "embed", "link"):
                qname = f"{{{_R_NS}}}{attr}"
                val = descendant.get(qname)
                if val and val in rel_id_map and val != rel_id_map[val]:
                    descendant.set(qname, rel_id_map[val])
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def find_shape_by_name(slide: Slide, name: str) -> Any | None:
    """Find a shape on a slide by its `.name` property, recursing into groups.

    Returns the first matching shape at any nesting depth, or None if no
    shape with that name exists on the slide. CV templates commonly group
    Name + Role + Headshot as a single unit, and without group recursion
    those named shapes would be unreachable here.
    """
    for shape in _walk_shapes(slide.shapes):
        if shape.name == name:
            return shape
    return None


def set_shape_value(shape: Any, value: Any) -> None:
    """Write a value into a shape's text frame.

    Handles:
    - None or empty string -> clear the text frame
    - list -> each item becomes a separate paragraph (bullet points inherit
      from the placeholder's paragraph-level formatting in the template)
    - scalar (str / int / float / any str-convertible) -> a single text
      value; strings containing `\\n` are split into paragraphs

    All non-empty values route through ``_write_paragraphs`` so the template's
    paragraph-level formatting (pPr) AND run-level formatting (rPr — font,
    size, weight, color) are preserved. A prior version took a "fast path"
    via ``tf.text = text`` for single-line scalars, but python-pptx's
    ``.text`` setter wipes both pPr and rPr, which silently flattened
    run-level font on every scalar-valued shape (name fields, role titles,
    summary headers). The one-item list path preserves them by capturing
    and re-injecting before/after clear().
    """
    if not getattr(shape, "has_text_frame", False):
        return

    tf = shape.text_frame

    if value is None or value == "":
        tf.clear()
        return

    if isinstance(value, list):
        items = [str(item) for item in value if item is not None and item != ""]
        _write_paragraphs(tf, items)
        return

    # Scalar: stringify and route through the right writer. Multi-line scalars
    # that target a multi-run-br template paragraph (run + <a:br/> + run with
    # distinct rPr) get special handling to preserve each run's formatting —
    # that's the CV idiom of "bold Role, Firm\n(italic duration)". Otherwise
    # fall back to _write_paragraphs so rPr/pPr are preserved.
    text = str(value)
    if "\n" in text and _is_multirun_br_first_paragraph(tf):
        _write_multirun_br(tf, text.split("\n"))
    elif "\n" in text:
        _write_paragraphs(tf, text.split("\n"))
    else:
        _write_paragraphs(tf, [text])


def _write_paragraphs(text_frame: Any, items: list[str]) -> None:
    """Write a list of strings as paragraphs into a text frame, one per item.

    Preserves paragraph-level formatting (bullet style, indent, alignment)
    by capturing the first existing paragraph's `<a:pPr>` before clearing
    the text frame, then re-injecting that pPr into every new paragraph.

    Without this preservation step, `text_frame.clear()` wipes all pPr and
    `add_paragraph()` produces bare paragraphs that lose the template's
    bullet styling — so a 3-item list ends up with one bullet on the first
    item and none on the rest. This is the v0.1.10 fix.

    Empty input clears the text frame. Non-empty input keeps the first
    existing paragraph's style profile for every output paragraph.
    """
    if not items:
        text_frame.clear()
        return

    # Capture the first existing paragraph's pPr + first run's rPr BEFORE
    # clear() wipes them. These encode the template's bullet/indent/font style.
    first_pPr_copy = None
    first_rPr_copy = None
    source_multirun_template: Any = None
    existing_paras = text_frame.paragraphs
    if existing_paras:
        first_p_xml = existing_paras[0]._p
        first_pPr = first_p_xml.find(qn("a:pPr"))
        if first_pPr is not None:
            first_pPr_copy = deepcopy(first_pPr)
        # Also capture the first run's rPr (font size, color, weight, etc.)
        first_r = first_p_xml.find(qn("a:r"))
        if first_r is not None:
            first_rPr = first_r.find(qn("a:rPr"))
            if first_rPr is not None:
                first_rPr_copy = deepcopy(first_rPr)
        # If the source paragraph itself is multi-run-br (run + <a:br/> + run
        # with distinct rPrs), capture the whole <a:p> as a prototype so we
        # can clone it per item and preserve per-run formatting when items
        # contain '\n'. Unblocks the CV "background" list idiom:
        #   • Bold Role, Firm
        #     (italic N years)
        # where each bullet is a multi-run-br paragraph, not a plain run.
        source_runs = first_p_xml.findall(qn("a:r"))
        source_brs = first_p_xml.findall(qn("a:br"))
        if len(source_runs) >= 2 and len(source_brs) >= 1:
            source_multirun_template = deepcopy(first_p_xml)

    # Route: if any item contains '\n' AND the source is a multi-run-br
    # paragraph template, clone the template per item (preserving per-run
    # rPr + br) and inject split parts into each run. Otherwise fall back
    # to the classic single-run-per-paragraph path.
    if source_multirun_template is not None and any("\n" in item for item in items):
        _write_paragraphs_cloning_multirun_br(text_frame, items, source_multirun_template)
        return

    text_frame.clear()

    # After clear(), text_frame has exactly one empty paragraph to reuse.
    first_p = text_frame.paragraphs[0]
    first_p.text = items[0]
    _apply_preserved_format(first_p, first_pPr_copy, first_rPr_copy)

    for item in items[1:]:
        new_p = text_frame.add_paragraph()
        new_p.text = item
        _apply_preserved_format(new_p, first_pPr_copy, first_rPr_copy)


def _write_paragraphs_cloning_multirun_br(
    text_frame: Any, items: list[str], source_p_template: Any
) -> None:
    """Emit one paragraph per item by deep-cloning a multi-run-br prototype.

    Used when the template's first paragraph is a multi-run-br idiom
    (e.g. bold `Role, Firm` + `<a:br/>` + italic `(N years)`) and the
    caller wants every list item rendered with the same structure.
    Unlike `_write_paragraphs`, which creates fresh paragraphs and
    reapplies a single rPr profile to the whole line, this function
    clones the source paragraph XML (pPr, runs, brs, per-run rPr) and
    only overwrites the `<a:t>` text content of each run.

    Per-item rules:
    - If item contains '\\n': split into parts, walk cloned runs in
      order, replace each `<a:t>` text with the corresponding part.
      Padding / merging follows the same rules as `_write_multirun_br`.
    - If item has no '\\n': replace the first run's `<a:t>` with the
      full item text, clear trailing runs' `<a:t>`, and strip `<a:br/>`
      children so the cloned paragraph collapses to a single line.
    """
    # Locate the text body via the existing first paragraph so we don't
    # have to poke at TextFrame's private `_txBody` handle directly.
    if not text_frame.paragraphs:
        return
    first_p_xml = text_frame.paragraphs[0]._p
    txBody = first_p_xml.getparent()
    if txBody is None:
        return

    # Drop every existing paragraph before appending clones — the
    # template prototype was captured before this wipe, so it's safe.
    for existing_p in txBody.findall(qn("a:p")):
        txBody.remove(existing_p)

    for item in items:
        cloned_p = deepcopy(source_p_template)
        cloned_segments = _segment_runs_by_br(cloned_p)
        if not cloned_segments:
            txBody.append(cloned_p)
            continue

        if "\n" in item:
            parts = item.split("\n")
            if len(parts) < len(cloned_segments):
                aligned: list[str] = list(parts) + [""] * (len(cloned_segments) - len(parts))
            elif len(parts) > len(cloned_segments):
                aligned = list(parts[: len(cloned_segments) - 1])
                aligned.append(" ".join(parts[len(cloned_segments) - 1 :]))
            else:
                aligned = list(parts)
            # Per-segment write: first run of each segment gets the text,
            # trailing runs in the same segment are cleared. Same semantics
            # as `_write_multirun_br` for the scalar case.
            for segment, part in zip(cloned_segments, aligned, strict=True):
                if not segment:
                    continue
                first_t = segment[0].find(qn("a:t"))
                if first_t is None:
                    first_t = etree.SubElement(segment[0], qn("a:t"))
                first_t.text = part
                for extra_run in segment[1:]:
                    extra_t = extra_run.find(qn("a:t"))
                    if extra_t is not None:
                        extra_t.text = ""
        else:
            # Single-line item against a multi-segment prototype: put the
            # text in segment 0's first run, clear everything else, strip
            # the `<a:br/>` children so the paragraph collapses to one
            # flat visual line.
            first_run = cloned_segments[0][0]
            first_t = first_run.find(qn("a:t"))
            if first_t is None:
                first_t = etree.SubElement(first_run, qn("a:t"))
            first_t.text = item
            # Clear everything else across all segments
            for segment in cloned_segments:
                for run_el in segment if segment is not cloned_segments[0] else segment[1:]:
                    extra_t = run_el.find(qn("a:t"))
                    if extra_t is not None:
                        extra_t.text = ""
            for br_el in cloned_p.findall(qn("a:br")):
                cloned_p.remove(br_el)

        txBody.append(cloned_p)

    # OOXML requires at least one <a:p> in a txBody.
    if not txBody.findall(qn("a:p")):
        etree.SubElement(txBody, qn("a:p"))


def _apply_preserved_format(paragraph: Any, pPr_copy: Any, rPr_copy: Any) -> None:
    """Inject a preserved pPr and the first run's rPr into a paragraph.

    pPr is inserted at the start of the paragraph element (OOXML requires
    pPr to precede text runs). rPr is injected at the start of the
    paragraph's first run.
    """
    if pPr_copy is not None:
        # Remove any existing pPr (added by add_paragraph) then insert ours
        existing_pPr = paragraph._p.find(qn("a:pPr"))
        if existing_pPr is not None:
            paragraph._p.remove(existing_pPr)
        paragraph._p.insert(0, deepcopy(pPr_copy))

    if rPr_copy is not None:
        first_r = paragraph._p.find(qn("a:r"))
        if first_r is not None:
            existing_rPr = first_r.find(qn("a:rPr"))
            if existing_rPr is not None:
                first_r.remove(existing_rPr)
            first_r.insert(0, deepcopy(rPr_copy))


def _is_multirun_br_first_paragraph(text_frame: Any) -> bool:
    """Return True if the text frame's first paragraph is a "dual-run-br" idiom.

    The pattern is a single paragraph containing two or more `<a:r>` runs
    separated by at least one `<a:br/>` soft break — each run carrying its
    own `<a:rPr>` so the runs can render with distinct formatting (bold
    primary text, italic bracketed secondary, etc.). This is the CV
    template idiom for cells like::

        Management Principal, client        (run 0, bold)
        (3 years)                          (run 1, italic)

    where the line break is Shift+Enter in PowerPoint — a `<a:br/>`, not
    a new paragraph. Used by `set_shape_value` and `populate_table` to
    decide whether to route a multi-line scalar value through
    `_write_multirun_br` (which preserves per-run rPr and the br) or
    through `_write_paragraphs` (which treats each line as a new
    paragraph and flattens to one rPr profile).
    """
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        return False
    p_xml = paragraphs[0]._p
    runs = p_xml.findall(qn("a:r"))
    brs = p_xml.findall(qn("a:br"))
    return len(runs) >= 2 and len(brs) >= 1


def _segment_runs_by_br(paragraph_xml: Any) -> list[list[Any]]:
    """Walk a paragraph's direct children in document order and group
    `<a:r>` runs into segments separated by `<a:br/>` elements.

    Returns a list of segments, each a list of run elements. A paragraph
    with no runs returns `[]`. A paragraph with runs but no `<a:br/>`
    returns a single segment containing all runs. This lets the writer
    preserve the VISUAL line structure of a multi-run-br cell even when
    the template has extra inline runs inside a single visual line
    (e.g., 3 runs + 1 br where the first visual line is authored as two
    adjacent bold runs — a common side effect of editing in PowerPoint).
    """
    segments: list[list[Any]] = [[]]
    r_tag = qn("a:r")
    br_tag = qn("a:br")
    for child in paragraph_xml:
        if child.tag == r_tag:
            segments[-1].append(child)
        elif child.tag == br_tag:
            segments.append([])
    # Drop any trailing empty segment (paragraph ends with <a:br/>)
    if segments and not segments[-1]:
        segments.pop()
    # Drop any leading empty segment (paragraph starts with <a:br/>)
    if segments and not segments[0]:
        segments.pop(0)
    return segments


def _write_multirun_br(text_frame: Any, parts: list[str]) -> None:
    """Overwrite the first paragraph's visual lines with `parts`.

    Segments runs by `<a:br/>` boundary, then assigns one `part` per
    segment. Within each segment, the FIRST run's `<a:t>` gets the
    assigned text and all trailing runs in that segment are cleared
    (their `<a:rPr>` is preserved so the rendered font/weight stays
    authored, but their text is empty so they don't leak template
    content). This matches how a reader sees the cell: one line per
    visual segment, whatever authoring runs the segment happens to
    contain.

    Preserves per-run `<a:rPr>` (font, weight, italics, colour) and the
    `<a:br/>` elements themselves. The paragraph structure stays exactly
    as the template authored it; only `<a:t>` text content changes. Any
    paragraphs after the first are dropped so the output matches the
    template's single-paragraph layout.

    Segment / part count handling:
    - len(parts) == len(segments): one-to-one.
    - len(parts) <  len(segments): trailing segments are cleared (first
      run of each trailing segment gets empty text, rest stay cleared).
    - len(parts) >  len(segments): excess parts are joined with a single
      space and merged into the last segment so content is never
      silently dropped.

    Empty input clears the text frame and returns.
    """
    if not parts:
        text_frame.clear()
        return

    paragraphs = text_frame.paragraphs
    if not paragraphs:
        return
    first_p_xml = paragraphs[0]._p
    segments = _segment_runs_by_br(first_p_xml)
    if not segments:
        # No runs to overwrite — fall back to paragraph-style write so
        # the caller's content still lands somewhere visible.
        _write_paragraphs(text_frame, parts)
        return

    # Pad-or-merge parts to match segment count.
    if len(parts) < len(segments):
        aligned: list[str] = list(parts) + [""] * (len(segments) - len(parts))
    elif len(parts) > len(segments):
        aligned = list(parts[: len(segments) - 1])
        aligned.append(" ".join(parts[len(segments) - 1 :]))
    else:
        aligned = list(parts)

    for segment, part in zip(segments, aligned, strict=True):
        if not segment:
            continue
        # First run of the segment gets the assigned text; trailing runs
        # in the same segment are cleared (rPr preserved so the visual
        # authoring stays, but content is replaced cleanly).
        first_run = segment[0]
        first_t = first_run.find(qn("a:t"))
        if first_t is None:
            first_t = etree.SubElement(first_run, qn("a:t"))
        first_t.text = part
        for extra_run in segment[1:]:
            extra_t = extra_run.find(qn("a:t"))
            if extra_t is not None:
                extra_t.text = ""

    # Drop any paragraphs beyond the first — this idiom is
    # single-paragraph by definition, and leaving trailing paragraphs
    # from the template example would leak stale content.
    txBody = first_p_xml.getparent()
    all_paras = txBody.findall(qn("a:p"))
    for extra_p in all_paras[1:]:
        txBody.remove(extra_p)


def _capture_paragraph_profile(paragraphs: Any, index: int) -> tuple[Any | None, Any | None]:
    """Capture a (pPr, first-run rPr) style profile from a specific paragraph.

    Returns deep-copied XML elements so the caller can safely reuse them
    after the containing text frame is cleared. Either element may be
    None if the source paragraph doesn't have a pPr (uses layout defaults)
    or has no runs (empty paragraph). Index out of range also returns
    (None, None) — the caller should fall back or warn.

    This is the machinery behind `populate_sections`: a sectioned-list
    shape has two distinct visual styles (header, bullet) encoded in
    different paragraphs of the template example; capturing both before
    `tf.clear()` lets us re-emit alternating profiles as we populate.
    """
    if index < 0 or index >= len(paragraphs):
        return None, None
    target_p_xml = paragraphs[index]._p
    pPr = target_p_xml.find(qn("a:pPr"))
    pPr_copy = deepcopy(pPr) if pPr is not None else None
    first_r = target_p_xml.find(qn("a:r"))
    rPr_copy = None
    if first_r is not None:
        rPr = first_r.find(qn("a:rPr"))
        if rPr is not None:
            rPr_copy = deepcopy(rPr)
    return pPr_copy, rPr_copy


def populate_sections(
    shape: Any, section_config: SectionConfig, sections_data: list[Any]
) -> list[str]:
    """Populate a sectioned-list shape from a list of {header, items} dicts.

    The template shape must be a text frame whose example paragraphs
    encode the intended visual profile for headers (at
    ``section_config.header_index``) and bullets (at
    ``section_config.bullet_index``). This function captures both
    profiles, clears the text frame, then re-emits paragraphs in the
    order prescribed by `sections_data`:

        for each section in sections_data:
            emit one header paragraph (header profile)
            emit N bullet paragraphs (bullet profile), one per item

    Returns a list of warning strings (empty on clean populate). Common
    warnings:
    - shape has no text frame (misrouted config)
    - template has fewer paragraphs than the configured indices reach
    - a section entry is not a dict or is missing header/items
    - header profile or bullet profile could not be captured (template
      example paragraph is empty)

    On a completely empty `sections_data` the text frame is cleared.
    """
    warnings: list[str] = []

    if not getattr(shape, "has_text_frame", False):
        warnings.append(
            f"shape {shape.name!r} is configured as a sectioned list but has "
            "no text frame; skipping"
        )
        return warnings

    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)

    if section_config.header_index >= len(paragraphs):
        warnings.append(
            f"shape {shape.name!r}: header_index={section_config.header_index} "
            f"but template has only {len(paragraphs)} paragraph(s); cannot "
            "capture header style profile"
        )
        return warnings
    if section_config.bullet_index >= len(paragraphs):
        warnings.append(
            f"shape {shape.name!r}: bullet_index={section_config.bullet_index} "
            f"but template has only {len(paragraphs)} paragraph(s); cannot "
            "capture bullet style profile"
        )
        return warnings

    header_pPr, header_rPr = _capture_paragraph_profile(paragraphs, section_config.header_index)
    bullet_pPr, bullet_rPr = _capture_paragraph_profile(paragraphs, section_config.bullet_index)
    if header_pPr is None and header_rPr is None:
        warnings.append(
            f"shape {shape.name!r}: template paragraph at header_index="
            f"{section_config.header_index} has no pPr or runs; header style "
            "will fall back to defaults"
        )
    if bullet_pPr is None and bullet_rPr is None:
        warnings.append(
            f"shape {shape.name!r}: template paragraph at bullet_index="
            f"{section_config.bullet_index} has no pPr or runs; bullet style "
            "will fall back to defaults"
        )

    # Build the flat emission plan: list of (text, pPr, rPr) tuples.
    # Validation is lenient (warn and skip malformed sections) so one bad
    # section doesn't nuke the whole placeholder.
    plan: list[tuple[str, Any, Any]] = []
    for section_index, section in enumerate(sections_data):
        if not isinstance(section, dict):
            warnings.append(
                f"shape {shape.name!r}: section {section_index} is "
                f"{type(section).__name__}, expected a dict with 'header' and "
                "'items'; skipping"
            )
            continue
        header_text = section.get("header")
        items_value = section.get("items", [])
        if not isinstance(header_text, str) or header_text == "":
            warnings.append(
                f"shape {shape.name!r}: section {section_index} is missing "
                "a non-empty 'header' string; skipping"
            )
            continue
        if not isinstance(items_value, list):
            warnings.append(
                f"shape {shape.name!r}: section {section_index} ({header_text!r}): "
                f"'items' must be a list, got {type(items_value).__name__}; "
                "rendering header alone"
            )
            items_value = []
        plan.append((header_text, header_pPr, header_rPr))
        for item in items_value:
            if item is None or item == "":
                continue
            plan.append((str(item), bullet_pPr, bullet_rPr))

    if not plan:
        tf.clear()
        return warnings

    tf.clear()
    # After clear() there is exactly one empty paragraph to reuse.
    first_p = tf.paragraphs[0]
    first_text, first_pPr_copy, first_rPr_copy = plan[0]
    first_p.text = first_text
    _apply_preserved_format(first_p, first_pPr_copy, first_rPr_copy)

    for text, pPr_copy, rPr_copy in plan[1:]:
        new_p = tf.add_paragraph()
        new_p.text = text
        _apply_preserved_format(new_p, pPr_copy, rPr_copy)

    return warnings


def is_picture_placeholder(shape: Any) -> bool:
    """Return True if the shape is a python-pptx PicturePlaceholder.

    PicturePlaceholders have an `insert_picture(path)` method to embed an
    image into the placeholder slot, keeping the shape's crop geometry and
    position. Regular Picture shapes (inserted via `add_picture`) are a
    different class and need different handling.
    """
    return isinstance(shape, PicturePlaceholder)


def set_picture(shape: Any, value: Any, base_dir: Path | None = None) -> None:
    """Insert an image file into a PicturePlaceholder shape.

    `value` is interpreted as a path to an image file. If it's a relative
    path and `base_dir` is provided, it resolves against `base_dir` (the
    directory of the YAML record file, typically). Otherwise the path is
    used as-is.

    Empty string or None is treated as "no picture" and silently no-ops —
    the placeholder retains whatever it already had (from the source slide
    inheritance), so leaving a `photo:` field empty keeps the example
    headshot as a fallback.
    """
    if value is None or value == "":
        return

    image_path = Path(str(value)).expanduser()
    if not image_path.is_absolute() and base_dir is not None:
        image_path = (base_dir / image_path).resolve()

    if not image_path.exists():
        raise FileNotFoundError(f"picture file not found: {image_path}")

    shape.insert_picture(str(image_path))


def populate_table(shape: Any, table_config: TableConfig, rows: list[Any]) -> list[str]:
    """Populate a table shape from a list of row dicts.

    Each element of `rows` is a dict whose keys should match the names in
    `table_config.columns`. Column values can be scalars or lists; lists
    are joined with `table_config.list_joiner` to produce multi-line cell
    text that inherits the template's paragraph-level bullet formatting.

    Returns a list of warning strings (empty if everything populated cleanly).
    Common warnings:
    - data has more rows than the template table can hold (truncation)
    - data has fewer rows than the template's data rows (excess rows cleared)
    - a row dict is missing a column key (that cell is cleared + warned)
    - the shape isn't actually a table (recombinase is called on the wrong shape)

    Stale-data invariant: because the duplicate-and-populate pattern brings
    the source slide's example text forward, any cell we don't explicitly
    overwrite would leak example content into the output. `populate_table`
    therefore clears every data cell that is not replaced with a real value.
    """
    warnings: list[str] = []

    if not getattr(shape, "has_table", False):
        warnings.append(
            f"shape {shape.name!r} is configured as a table but is not a "
            "table shape on the template; skipping"
        )
        return warnings

    table = shape.table
    all_rows = list(table.rows)
    start_row = 1 if table_config.header_row else 0
    end_row = len(all_rows) - table_config.footer_rows
    if end_row <= start_row:
        warnings.append(
            f"table {shape.name!r} has header_row={table_config.header_row} and "
            f"footer_rows={table_config.footer_rows} but only {len(all_rows)} "
            "rows total; no data rows left to populate"
        )
        return warnings
    data_rows = all_rows[start_row:end_row]

    if len(rows) > len(data_rows):
        warnings.append(
            f"table {shape.name!r} has {len(data_rows)} data rows but "
            f"{len(rows)} records were provided; truncating to {len(data_rows)}"
        )

    for row_index, row_dict in enumerate(rows[: len(data_rows)]):
        target_row = data_rows[row_index]
        is_dict = isinstance(row_dict, dict)
        if not is_dict:
            warnings.append(
                f"table {shape.name!r} row {row_index}: expected a "
                f"dict, got {type(row_dict).__name__}; clearing row"
            )
        for col_index, column_name in enumerate(table_config.columns):
            if col_index >= len(target_row.cells):
                warnings.append(
                    f"table {shape.name!r} row {row_index} has only "
                    f"{len(target_row.cells)} cells; column {column_name!r} "
                    "exceeds the template width, skipping"
                )
                continue
            cell = target_row.cells[col_index]
            if not is_dict:
                _clear_cell(cell)
                continue
            if column_name not in row_dict:
                warnings.append(
                    f"table {shape.name!r} row {row_index}: missing column "
                    f"{column_name!r}; cell cleared"
                )
                _clear_cell(cell)
                continue
            value = row_dict[column_name]
            if value is None or value == "":
                _clear_cell(cell)
                continue
            if isinstance(value, list):
                items = [str(item) for item in value if item is not None and item != ""]
                if not items:
                    _clear_cell(cell)
                    continue
                text = table_config.list_joiner.join(items)
            elif isinstance(value, dict):
                warnings.append(
                    f"table {shape.name!r} row {row_index} column "
                    f"{column_name!r}: value is a dict; expected scalar or list. "
                    "Cell cleared."
                )
                _clear_cell(cell)
                continue
            else:
                text = str(value)
            # Route the write based on the source cell's structure:
            # - Multi-run-br idiom (one para, >=2 runs, >=1 br) with a
            #   newline-containing value: preserve per-run rPr and the br
            #   so the CV "bold Role\n(italic duration)" idiom survives.
            # - Otherwise route through _write_paragraphs so the cell's pPr
            #   AND run-level rPr (font size, colour, weight) are preserved.
            #   A prior fast path via `cell.text_frame.text = text` flattened
            #   run-level font on single-line writes — the v0.1.13 fix.
            cell_tf = cell.text_frame
            if "\n" in text and _is_multirun_br_first_paragraph(cell_tf):
                _write_multirun_br(cell_tf, text.split("\n"))
            elif "\n" in text:
                _write_paragraphs(cell_tf, text.split("\n"))
            else:
                _write_paragraphs(cell_tf, [text])
        # Clear any trailing cells on this row that aren't covered by columns —
        # templates sometimes have more physical cells than configured columns,
        # and leaving them with example text would leak.
        for extra_col_index in range(len(table_config.columns), len(target_row.cells)):
            _clear_cell(target_row.cells[extra_col_index])

    # Clear every excess data row the record didn't cover. Without this, a
    # record with fewer entries than the template's capacity leaks the
    # duplicated example rows into the output.
    excess_row_count = len(data_rows) - len(rows)
    if excess_row_count > 0:
        warnings.append(
            f"table {shape.name!r} has {len(data_rows)} data rows but only "
            f"{len(rows)} record(s) provided; clearing {excess_row_count} "
            "unused row(s)"
        )
        for target_row in data_rows[len(rows) :]:
            for cell in target_row.cells:
                _clear_cell(cell)

    return warnings


def _clear_cell(cell: Any) -> None:
    """Clear all text from a table cell while preserving paragraph and run formatting.

    Captures the first paragraph's pPr (bullet style, indent, alignment) and
    the first run's rPr (font, size, weight, color) before clearing, then
    re-injects them into the resulting empty paragraph. This is the same
    capture/restore pattern used by ``_write_paragraphs`` — without it,
    ``cell.text_frame.text = ""`` wipes all formatting and cleared cells
    revert to application defaults if later hand-edited in PowerPoint.

    Spanned (merged non-origin) cells have no accessible text frame — python-
    pptx raises on `.text_frame` for them. Skip silently: the merge origin's
    text is cleared when that cell's own iteration visits it, so the visual
    region still renders empty.
    """
    if getattr(cell, "is_spanned", False):
        return

    text_frame = cell.text_frame
    paragraphs = text_frame.paragraphs

    # Capture formatting before clear() wipes it.
    first_pPr_copy = None
    first_rPr_copy = None
    if paragraphs:
        first_p_xml = paragraphs[0]._p
        first_pPr = first_p_xml.find(qn("a:pPr"))
        if first_pPr is not None:
            first_pPr_copy = deepcopy(first_pPr)
        first_r = first_p_xml.find(qn("a:r"))
        if first_r is not None:
            first_rPr = first_r.find(qn("a:rPr"))
            if first_rPr is not None:
                first_rPr_copy = deepcopy(first_rPr)

    text_frame.text = ""

    # Re-inject preserved formatting into the empty paragraph.
    # After `text = ""`, the paragraph may have no <a:r> element. If we
    # captured an rPr, create an empty run so _apply_preserved_format has
    # somewhere to inject it.
    empty_p = text_frame.paragraphs[0]
    if first_rPr_copy is not None and empty_p._p.find(qn("a:r")) is None:
        run_el = etree.SubElement(empty_p._p, qn("a:r"))
        etree.SubElement(run_el, qn("a:t"))
    _apply_preserved_format(empty_p, first_pPr_copy, first_rPr_copy)


def remove_slide(presentation: Any, slide: Slide) -> None:
    """Remove a slide from the presentation.

    python-pptx has no native `slides.remove()` — this walks the slide id
    list on the presentation part, finds the relationship whose target
    matches the target slide, and drops both the rel and the id list entry.

    Ref: https://github.com/scanny/python-pptx/issues/67
    """
    slide_id = slide.slide_id
    xml_slides = presentation.slides._sldIdLst

    for sl in xml_slides:
        rid = sl.get(f"{{{_R_NS}}}id")
        if rid is None:
            continue
        slide_from_rel = presentation.part.related_part(rid).slide
        if slide_from_rel.slide_id == slide_id:
            presentation.part.drop_rel(rid)
            xml_slides.remove(sl)
            return


def _value_char_length(value: Any) -> int:
    """Estimate the text-frame character length a value would occupy."""
    if value is None or value == "":
        return 0
    if isinstance(value, list):
        items = [str(item) for item in value if item is not None and item != ""]
        if not items:
            return 0
        return sum(len(item) for item in items) + (len(items) - 1)  # inter-para newlines
    return len(str(value))


def _capture_baseline_lengths(
    source_slide: Slide, placeholder_map: dict[str, str]
) -> dict[str, int]:
    """Record the text length of each configured placeholder on the source slide.

    Returned dict maps data field name to the source's current character count.
    Shapes without a text frame are skipped. These baselines are compared
    against per-record values to flag probable overflow.
    """
    baselines: dict[str, int] = {}
    for field_name, shape_name in placeholder_map.items():
        shape = find_shape_by_name(source_slide, shape_name)
        if shape is None or not getattr(shape, "has_text_frame", False):
            continue
        baselines[field_name] = len(shape.text_frame.text or "")
    return baselines


def generate_deck(
    config: TemplateConfig,
    records: list[dict[str, Any]],
    output_path: Path | str,
) -> dict[str, Any]:
    """Generate an output pptx from a template + records.

    Returns a dict summary with counts and any warnings. Warnings include
    missing shapes, missing fields, set-value exceptions, and — when
    `config.overflow_ratio > 0` — probable overflow cases where a record's
    field is substantially larger than the source-slide baseline.
    """
    output_path = Path(output_path).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(config.template))

    total_slides = len(presentation.slides)
    if config.source_slide_index < 1 or config.source_slide_index > total_slides:
        raise ValueError(
            f"source_slide_index={config.source_slide_index} out of range "
            f"(template has {total_slides} slide(s))"
        )

    # 1-based -> 0-based
    source_slide = presentation.slides[config.source_slide_index - 1]

    # Capture baseline text lengths BEFORE generation so overflow comparison
    # is against the original example, not the most recently populated slide.
    baselines: dict[str, int] = (
        _capture_baseline_lengths(source_slide, config.placeholders)
        if config.overflow_ratio > 0
        else {}
    )

    warnings: list[str] = []
    generated_count = 0

    for record_index, record in enumerate(records, start=1):
        record_id = record.get("id") or record.get("name") or f"record_{record_index}"
        new_slide = duplicate_slide(presentation, source_slide)

        for field_name, shape_name in config.placeholders.items():
            shape = find_shape_by_name(new_slide, shape_name)
            if shape is None:
                warnings.append(
                    f"record {record_id!r}: shape {shape_name!r} not found on new slide"
                )
                continue
            if field_name not in record:
                warnings.append(f"record {record_id!r}: no value for field {field_name!r}")
                continue

            value = record[field_name]

            # Picture placeholder: route to insert_picture instead of setting text
            if is_picture_placeholder(shape):
                try:
                    record_source = record.get("_recombinase_record_dir")
                    base_dir = Path(record_source) if isinstance(record_source, str) else None
                    set_picture(shape, value, base_dir=base_dir)
                except FileNotFoundError as exc:
                    warnings.append(
                        f"record {record_id!r}: picture for {field_name!r} not found: {exc}"
                    )
                except Exception as exc:
                    warnings.append(
                        f"record {record_id!r}: failed to insert picture "
                        f"{shape_name!r} ({field_name!r}): {exc}"
                    )
                continue

            # Text placeholder: existing path
            try:
                set_shape_value(shape, value)
            except Exception as exc:
                warnings.append(
                    f"record {record_id!r}: failed to set {shape_name!r} ({field_name!r}): {exc}"
                )
                continue

            # Overflow heuristic: compare new text length against source baseline.
            baseline = baselines.get(field_name, 0)
            if baseline > 0:
                new_length = _value_char_length(value)
                if new_length > baseline * config.overflow_ratio:
                    ratio = new_length / baseline
                    warnings.append(
                        f"record {record_id!r}: field {field_name!r} is "
                        f"{ratio:.1f}x the source baseline ({new_length} vs "
                        f"{baseline} chars) — may overflow shape {shape_name!r}"
                    )

        # Tables: populate from config.tables using the record's list-of-dicts data
        for table_field_name, table_config in config.tables.items():
            table_shape = find_shape_by_name(new_slide, table_config.shape)
            if table_shape is None:
                warnings.append(
                    f"record {record_id!r}: table shape {table_config.shape!r} "
                    "not found on new slide"
                )
                continue
            if table_field_name not in record:
                warnings.append(f"record {record_id!r}: no value for table {table_field_name!r}")
                continue
            table_value = record[table_field_name]
            if not isinstance(table_value, list):
                warnings.append(
                    f"record {record_id!r}: table {table_field_name!r} expects "
                    f"a list of row dicts, got {type(table_value).__name__}"
                )
                continue
            try:
                table_warnings = populate_table(table_shape, table_config, table_value)
            except Exception as exc:
                warnings.append(
                    f"record {record_id!r}: failed to populate table "
                    f"{table_config.shape!r} ({table_field_name!r}): {exc}"
                )
                continue
            warnings.extend(f"record {record_id!r}: {tw}" for tw in table_warnings)

        # Sections: populate from config.sections using the record's list-of-
        # {header, items} dicts. A "section" is a header paragraph followed
        # by N bullet paragraphs, and N sections can stack in one text frame.
        for section_field_name, section_config in config.sections.items():
            section_shape = find_shape_by_name(new_slide, section_config.shape)
            if section_shape is None:
                warnings.append(
                    f"record {record_id!r}: section shape "
                    f"{section_config.shape!r} not found on new slide"
                )
                continue
            if section_field_name not in record:
                warnings.append(
                    f"record {record_id!r}: no value for section {section_field_name!r}"
                )
                continue
            section_value = record[section_field_name]
            if not isinstance(section_value, list):
                warnings.append(
                    f"record {record_id!r}: section {section_field_name!r} "
                    f"expects a list of {{header, items}} dicts, got "
                    f"{type(section_value).__name__}"
                )
                continue
            try:
                section_warnings = populate_sections(section_shape, section_config, section_value)
            except Exception as exc:
                warnings.append(
                    f"record {record_id!r}: failed to populate sections "
                    f"{section_config.shape!r} ({section_field_name!r}): {exc}"
                )
                continue
            warnings.extend(f"record {record_id!r}: {sw}" for sw in section_warnings)

        generated_count += 1

    if generated_count == 0 and config.clear_source_slide:
        warnings.append(
            "0 records generated with clear_source_slide=True — output deck will be empty"
        )

    if config.clear_source_slide:
        remove_slide(presentation, source_slide)

    presentation.save(str(output_path))

    return {
        "output": str(output_path),
        "records_generated": generated_count,
        "warnings": warnings,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
    }
