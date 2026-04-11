"""CLI entry point for recombinase, built on Typer for rich human UX.

Subcommands:
- new      : scaffold a project directory (template/, cv-data/, output/)
- inspect  : print structural metadata of a pptx template
- init     : write a scaffold config YAML from a template's shape names
- generate : populate a template from YAML records and write to an output pptx
"""

from __future__ import annotations

import os
import re
import traceback
from collections.abc import Callable
from pathlib import Path

import typer
import yaml
from pptx.exc import PackageNotFoundError

from recombinase import __version__
from recombinase.config import load_config, write_scaffold_config
from recombinase.generate import generate_deck, load_records
from recombinase.inspect import (
    format_template_info,
    inspect_template,
    shape_names_from_slide,
    shape_types_from_slide,
)


def _default_project_dir() -> Path:
    """Pick a sensible default directory for `recombinase new`.

    On Windows with OneDrive configured, use `$env:OneDrive\\cv`. On other
    platforms (or if OneDrive isn't set), fall back to `~/cv`. Users can
    always override by passing an explicit path. Windows environment variables
    are case-insensitive, so using uppercase keys here works on all platforms.
    """
    onedrive = os.environ.get("ONEDRIVE") or os.environ.get("ONEDRIVECOMMERCIAL")
    if onedrive:
        return Path(onedrive) / "cv"
    return Path.home() / "cv"


def _find_template_in_cwd() -> Path | None:
    """Look for a single .pptx/.pptm file in conventional locations.

    Searches (in order): `./template/*.pptm`, `./template/*.pptx`,
    `./*.pptm`, `./*.pptx`. Returns the unique match if there is one;
    returns `None` if no template is found or if more than one candidate
    exists in any directory (ambiguous — the user must name it explicitly).
    """
    cwd = Path.cwd()
    search_dirs = (cwd / "template", cwd)
    for search_dir in search_dirs:
        if not search_dir.is_dir():
            continue
        candidates = sorted([*search_dir.glob("*.pptm"), *search_dir.glob("*.pptx")])
        if len(candidates) == 1:
            return candidates[0]
        if len(candidates) > 1:
            return None  # ambiguous — refuse to guess
    return None


def _resolve_template_arg(template: Path | None) -> Path:
    """Return an explicit template path, or auto-detect it, or exit with an error."""
    if template is not None:
        return template
    found = _find_template_in_cwd()
    if found is None:
        typer.secho(
            "No template specified and none auto-detected. Pass a path "
            "explicitly, or `cd` into a scaffolded project folder (with a "
            "single .pptx/.pptm inside `template/` or the current directory).",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)
    typer.secho(f"Auto-detected template: '{found}'", fg=typer.colors.CYAN)
    return found


_WORKFLOW_EPILOG = """\
[bold]Typical workflow[/bold] (run in order):

  1. [cyan]recombinase new[/cyan]                           scaffold a project folder
  2. [cyan]recombinase inspect[/cyan] TEMPLATE              see template shape names
  3. [cyan]recombinase init[/cyan] TEMPLATE                 write a config to edit
     [dim](then edit the config to map field names to shape names)[/dim]
  4. [cyan]recombinase generate[/cyan] -c CONFIG -d DATA -o OUT.pptx

Run [cyan]recombinase <command> --help[/cyan] for details on a single command.
"""


app = typer.Typer(
    name="recombinase",
    help=(
        "Template-guided pptx synthesis: inspect templates, scaffold configs, "
        "and generate populated decks from structured YAML data."
    ),
    epilog=_WORKFLOW_EPILOG,
    no_args_is_help=True,
    add_completion=False,
    rich_markup_mode="rich",
)


def _version_callback(value: bool) -> None:
    if value:
        typer.echo(f"recombinase {__version__}")
        raise typer.Exit


@app.callback()
def _root(
    version: bool = typer.Option(
        False,
        "--version",
        "-V",
        callback=_version_callback,
        is_eager=True,
        help="Show version and exit.",
    ),
) -> None:
    """Recombinase — template-guided pptx synthesis."""


@app.command("new")
def cmd_new(
    project_dir: Path = typer.Argument(
        None,
        help=(
            "Path to the project directory to create. If omitted, defaults to "
            "$env:OneDrive/cv on Windows (when OneDrive is configured), "
            "otherwise ~/cv."
        ),
        resolve_path=True,
    ),
    force: bool = typer.Option(
        False,
        "--force",
        "-f",
        help="Proceed even if the target directory already exists and is non-empty.",
    ),
) -> None:
    """Scaffold a new project directory with template/, cv-data/, and output/ subfolders.

    Creates a conventional folder layout under the given path:

        <project-dir>/
          template/   — put your .pptx/.pptm template here
          cv-data/    — put your per-record YAML files here
          output/     — generated decks will land here

    Safe to run in OneDrive — it's a plain mkdir + README write, no sync surprises.
    """
    if project_dir is None:
        project_dir = _default_project_dir()
        typer.secho(
            f"No path given — defaulting to '{project_dir}'",
            fg=typer.colors.CYAN,
        )

    if project_dir.exists() and any(project_dir.iterdir()) and not force:
        typer.secho(
            f"Directory already exists and is not empty: '{project_dir}' (use --force to proceed)",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    subfolders = ("template", "cv-data", "output")
    for sub in subfolders:
        (project_dir / sub).mkdir(parents=True, exist_ok=True)

    readme_path = project_dir / "README.md"
    if not readme_path.exists():
        readme_path.write_text(
            "# Recombinase project\n\n"
            "This folder was scaffolded by `recombinase new`.\n\n"
            "## Layout\n\n"
            "- `template/` — place the source .pptx/.pptm template file here\n"
            "- `cv-data/` — one YAML file per record (consultant, use case, etc.)\n"
            "- `output/` — generated decks land here\n\n"
            "## Typical workflow\n\n"
            "```\n"
            "recombinase inspect template/YOUR_TEMPLATE.pptm\n"
            "recombinase init template/YOUR_TEMPLATE.pptm -o template/config.yaml\n"
            "# edit template/config.yaml to map field names to shape names\n"
            "# write one .yaml file per record in cv-data/\n"
            "recombinase generate -c template/config.yaml -d cv-data/ "
            "-o output/deck.pptx\n"
            "```\n",
            encoding="utf-8",
        )

    typer.secho(f"Created project: '{project_dir}'", fg=typer.colors.GREEN)
    for sub in subfolders:
        typer.echo(f"  {project_dir / sub}")
    typer.echo("\nNext: place your template file in the template/ folder, then run:")
    typer.echo(f'  recombinase inspect "{project_dir / "template" / "YOUR_TEMPLATE.pptm"}"')


@app.command("inspect")
def cmd_inspect(
    template: Path = typer.Argument(
        None,
        help=(
            "Path to a .pptx/.pptm template file. If omitted, recombinase "
            "looks for a single pptx/pptm in ./template/ or the current dir."
        ),
        resolve_path=True,
    ),
) -> None:
    """Print structural metadata of a pptx template.

    Shape names, types, placeholder info, and text character counts only —
    no actual slide text is read or printed, so output is safe to share.
    Run this first on any new template to discover the shape names you'll
    reference in the config's `placeholders:` section.

    With no argument, auto-detects a single pptx/pptm in `./template/` or
    the current directory — matches the layout created by `recombinase new`.
    """
    template = _resolve_template_arg(template)
    info = inspect_template(template)
    typer.echo(format_template_info(info))


@app.command("init")
def cmd_init(
    template: Path = typer.Argument(
        None,
        help=(
            "Path to a .pptx/.pptm template file. If omitted, recombinase "
            "looks for a single pptx/pptm in ./template/ or the current dir."
        ),
        resolve_path=True,
    ),
    source_slide_index: int = typer.Option(
        1,
        "--source-slide-index",
        "-s",
        min=1,
        help="1-based index of the slide to read shape names from.",
    ),
    output: Path = typer.Option(
        None,
        "--output",
        "-o",
        help=(
            "Path to write the scaffold config YAML. Defaults to "
            "`./template/config.yaml` if a `template/` folder exists "
            "(scaffolded layout), otherwise `./template-config.yaml`."
        ),
    ),
    force: bool = typer.Option(
        False,
        "--force",
        "-f",
        help="Overwrite an existing output file.",
    ),
) -> None:
    """Write a scaffold config YAML from a template's shape names.

    With no arguments, auto-detects the template in `./template/` or the
    current directory and writes the config next to it. Matches the layout
    created by `recombinase new`.
    """
    template = _resolve_template_arg(template)
    if output is None:
        template_dir = Path.cwd() / "template"
        if template_dir.is_dir():
            output = template_dir / "config.yaml"
        else:
            output = Path("template-config.yaml")
    info = inspect_template(template)
    shape_names = shape_names_from_slide(info, source_slide_index)

    if not shape_names:
        typer.secho(
            f"No shapes found on slide {source_slide_index} of '{template}'",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    if output.exists() and not force:
        typer.secho(
            f"Config file already exists: '{output}' (use --force to overwrite)",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    output.parent.mkdir(parents=True, exist_ok=True)
    write_scaffold_config(template, shape_names, output)
    typer.secho(f"Wrote scaffold config: '{output}'", fg=typer.colors.GREEN)
    typer.echo(f"Found {len(shape_names)} shape(s) on slide {source_slide_index}.")
    typer.echo(
        f'Next: edit "{output}" to map your data field names (left) to the shape names (right).'
    )


@app.command("validate")
def cmd_validate(
    config: Path = typer.Option(
        None,
        "--config",
        "-c",
        resolve_path=True,
        help=(
            "Path to the template config YAML. Defaults to "
            "`./template/config.yaml` when run inside a scaffolded project."
        ),
    ),
    strict: bool = typer.Option(
        False,
        "--strict",
        help="Exit non-zero if unused shapes are present in the template.",
    ),
) -> None:
    """Validate a config against its template — pre-flight check before generate.

    Verifies that:
    - the config loads (YAML is well-formed, required fields present)
    - the template file exists and is readable
    - every shape name in `placeholders` exists on the configured source slide
    - reports which template shapes are NOT mapped (unused)

    Catches typos (`Consultant_Name` vs `Consultant Name`) and stale configs
    (shapes renamed in the template) before you generate 15 broken slides.
    Exit 0 on success, 1 on missing shapes, 2 on unused shapes if --strict.
    """
    cwd = Path.cwd()
    if config is None:
        candidate = cwd / "template" / "config.yaml"
        if not candidate.is_file():
            typer.secho(
                "No --config specified and ./template/config.yaml not found. "
                "Run `recombinase init` first or pass --config explicitly.",
                fg=typer.colors.RED,
                err=True,
            )
            raise typer.Exit(code=1)
        config = candidate.resolve()

    typer.secho(f"Loading config: '{config}'", fg=typer.colors.CYAN)
    cfg = load_config(config)

    typer.secho(f"Inspecting template: '{cfg.template}'", fg=typer.colors.CYAN)
    info = inspect_template(cfg.template)

    template_shape_names = set(shape_names_from_slide(info, cfg.source_slide_index))
    template_shape_has_table = shape_types_from_slide(info, cfg.source_slide_index)
    placeholder_shape_names = set(cfg.placeholders.values())
    table_shape_names = {table.shape for table in cfg.tables.values()}
    config_shape_names = placeholder_shape_names | table_shape_names

    missing_shapes = config_shape_names - template_shape_names
    # Filter default/decorative shape names out of the unused set — they
    # are almost always structural (background rectangles, decorative
    # connectors, logos) that legitimately don't need a data mapping.
    # Without this filter `--strict` is unusable on any real template.
    _decorative_patterns = re.compile(
        r"^(Rectangle|Oval|Freeform|Straight Connector|Straight Arrow Connector|"
        r"TextBox|Title|Subtitle|Content Placeholder|Picture|Group|Table|"
        r"Chart|Diagram|Line|Arrow) \d+$"
    )
    structural_shapes = {name for name in template_shape_names if _decorative_patterns.match(name)}
    unused_shapes = (template_shape_names - config_shape_names) - structural_shapes

    # Type mismatch: placeholders pointed at table shapes, or tables
    # pointed at non-table shapes. Catches a silent failure where the
    # shape name matches but the shape kind is wrong for its role.
    placeholder_on_table: list[tuple[str, str]] = []  # (field, shape)
    table_on_non_table: list[tuple[str, str]] = []  # (field, shape)
    for field_name, shape_name in cfg.placeholders.items():
        if template_shape_has_table.get(shape_name, False):
            placeholder_on_table.append((field_name, shape_name))
    for field_name, table in cfg.tables.items():
        shape_name = table.shape
        if shape_name in template_shape_has_table and not template_shape_has_table[shape_name]:
            table_on_non_table.append((field_name, shape_name))

    # Surface matched shapes for context
    matched = config_shape_names & template_shape_names
    typer.secho(
        f"Matched {len(matched)}/{len(config_shape_names)} configured shapes.",
        fg=typer.colors.GREEN if not missing_shapes else typer.colors.YELLOW,
    )

    if missing_shapes:
        typer.secho(
            f"\nMissing shapes ({len(missing_shapes)}): configured in the "
            "placeholders or tables map but NOT found on the template's source slide.",
            fg=typer.colors.RED,
            err=True,
        )
        for shape_name in sorted(missing_shapes):
            placeholder_fields = [
                field for field, shape in cfg.placeholders.items() if shape == shape_name
            ]
            table_fields = [
                field for field, table in cfg.tables.items() if table.shape == shape_name
            ]
            provenance_bits: list[str] = []
            if placeholder_fields:
                provenance_bits.append(f"placeholder field: {', '.join(placeholder_fields)}")
            if table_fields:
                provenance_bits.append(f"table field: {', '.join(table_fields)}")
            provenance = "; ".join(provenance_bits) if provenance_bits else "unknown source"
            typer.secho(
                f"  - '{shape_name}' ({provenance})",
                fg=typer.colors.RED,
                err=True,
            )
        typer.secho(
            "\nFix: rename the shape in the template OR update the config "
            "right-hand side to match. Run `recombinase inspect` to see "
            "the actual shape names.",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    if placeholder_on_table or table_on_non_table:
        typer.secho(
            "\nShape type mismatch: the shape name matches the template but "
            "the shape kind is wrong for its config role.",
            fg=typer.colors.RED,
            err=True,
        )
        for field_name, shape_name in sorted(placeholder_on_table):
            typer.secho(
                f"  - placeholder field '{field_name}' -> '{shape_name}' "
                "(this is a TABLE shape; move it under `tables:` instead)",
                fg=typer.colors.RED,
                err=True,
            )
        for field_name, shape_name in sorted(table_on_non_table):
            typer.secho(
                f"  - table field '{field_name}' -> '{shape_name}' "
                "(this is NOT a table shape; move it under `placeholders:` instead)",
                fg=typer.colors.RED,
                err=True,
            )
        raise typer.Exit(code=1)

    if unused_shapes:
        typer.secho(
            f"\nUnused shapes ({len(unused_shapes)}): present on the template "
            "but not mapped in the config. These will be inherited as-is "
            "(example text stays in place unless you add them to placeholders).",
            fg=typer.colors.YELLOW,
            err=True,
        )
        for shape_name in sorted(unused_shapes):
            typer.secho(f"  - '{shape_name}'", fg=typer.colors.YELLOW, err=True)

        if strict:
            raise typer.Exit(code=2)

    typer.secho("\nConfig is valid.", fg=typer.colors.GREEN)


@app.command("generate")
def cmd_generate(
    config: Path = typer.Option(
        None,
        "--config",
        "-c",
        resolve_path=True,
        help=(
            "Path to the template config YAML. Defaults to "
            "`./template/config.yaml` when run inside a scaffolded project."
        ),
    ),
    data_dir: Path = typer.Option(
        None,
        "--data-dir",
        "-d",
        resolve_path=True,
        help=(
            "Directory containing per-record YAML files. Defaults to "
            "`./cv-data` when run inside a scaffolded project."
        ),
    ),
    output: Path = typer.Option(
        None,
        "--output",
        "-o",
        resolve_path=True,
        help=(
            "Path to write the generated pptx. Defaults to "
            "`./output/deck.pptx` when run inside a scaffolded project."
        ),
    ),
    force: bool = typer.Option(
        False,
        "--force",
        "-f",
        help="Overwrite an existing output file without prompting.",
    ),
    strict: bool = typer.Option(
        False,
        "--strict",
        help=(
            "Exit non-zero if any record had missing-field or missing-shape "
            "warnings. Default: warnings print but exit 0."
        ),
    ),
    dry_run: bool = typer.Option(
        False,
        "--dry-run",
        help=(
            "Simulate the full pipeline (load config, load records, "
            "validate shape mappings, report warnings) but do NOT write "
            "the output pptx. Use this to preview warnings before a "
            "real run."
        ),
    ),
) -> None:
    """Generate a populated pptx deck from a template + YAML data directory.

    With no arguments, resolves defaults against the `recombinase new`
    scaffolded layout (./template/config.yaml, ./cv-data/, ./output/deck.pptx)
    so you can `cd` into the project and just run `recombinase generate`.

    Pass `--dry-run` to preview without writing.
    """
    cwd = Path.cwd()

    # Auto-default branches: friendly error when the scaffolded path is
    # missing. Explicit paths passed by the user are not validated here —
    # load_config and load_records raise their own errors which main()
    # traps into clean messages.
    if config is None:
        candidate = cwd / "template" / "config.yaml"
        if not candidate.is_file():
            typer.secho(
                "No --config specified and ./template/config.yaml not found. "
                "Run `recombinase init` first or pass --config explicitly.",
                fg=typer.colors.RED,
                err=True,
            )
            raise typer.Exit(code=1)
        config = candidate.resolve()

    if data_dir is None:
        candidate = cwd / "cv-data"
        if not candidate.is_dir():
            typer.secho(
                "No --data-dir specified and ./cv-data/ not found. "
                "Create it or pass --data-dir explicitly.",
                fg=typer.colors.RED,
                err=True,
            )
            raise typer.Exit(code=1)
        data_dir = candidate.resolve()

    if output is None:
        output = (cwd / "output" / "deck.pptx").resolve()

    if output.suffix.lower() != ".pptx":
        typer.secho(
            f"Warning: output path '{output}' does not end in .pptx — "
            "PowerPoint may not open it correctly.",
            fg=typer.colors.YELLOW,
            err=True,
        )

    # Overwrite guard only applies to real runs, not dry-runs.
    if not dry_run and output.exists() and not force:
        typer.secho(
            f"Output file already exists: '{output}' (use --force to overwrite)",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    if not dry_run:
        output.parent.mkdir(parents=True, exist_ok=True)

    typer.secho(f"Loading config: '{config}'", fg=typer.colors.CYAN)
    cfg = load_config(config)

    typer.secho(f"Loading records from: '{data_dir}'", fg=typer.colors.CYAN)
    records = load_records(data_dir)

    if not records:
        typer.secho(
            f"No YAML records found in '{data_dir}'",
            fg=typer.colors.RED,
            err=True,
        )
        raise typer.Exit(code=1)

    if dry_run:
        typer.secho(
            f"DRY RUN: would generate {len(records)} slide(s) from "
            f"'{cfg.template.name}' -> '{output}'",
            fg=typer.colors.CYAN,
        )
        # Pipe through a tmp path so generate_deck's save doesn't touch the
        # real output; caller still gets the same warning list. Use a context
        # manager so the temp directory is cleaned up even if generate_deck
        # raises mid-run.
        import tempfile

        with tempfile.TemporaryDirectory() as tmp_root:
            tmp_output = Path(tmp_root) / "dry-run.pptx"
            result = generate_deck(cfg, records, tmp_output)

        typer.secho(
            f"DRY RUN complete. Would have written: '{output}'",
            fg=typer.colors.CYAN,
        )
        typer.echo(f"Records that would be generated: {result['records_generated']}")
    else:
        typer.secho(
            f"Generating {len(records)} slide(s) from '{cfg.template.name}'...",
            fg=typer.colors.CYAN,
        )
        result = generate_deck(cfg, records, output)

        typer.secho(f"Generated: '{result['output']}'", fg=typer.colors.GREEN)
        typer.echo(f"Records: {result['records_generated']}")

    if result["warnings"]:
        typer.secho(
            f"Warnings ({len(result['warnings'])}):",
            fg=typer.colors.YELLOW,
            err=True,
        )
        for warning in result["warnings"]:
            typer.secho(f"  - {warning}", fg=typer.colors.YELLOW, err=True)
        if strict:
            raise typer.Exit(code=2)


def _print_error(message: str) -> None:
    typer.secho(f"Error: {message}", fg=typer.colors.RED, err=True)


def _format_permission_error(exc: BaseException) -> str:
    target = getattr(exc, "filename", None)
    if target:
        return f"Cannot write '{target}': the file may be open in PowerPoint. Close it and re-run."
    return f"Permission denied: {exc}"


def _format_package_not_found(exc: BaseException) -> str:
    """Format a PackageNotFoundError with PowerPoint-lock detection.

    python-pptx raises `PackageNotFoundError` when it can't parse a file as
    an OOXML zip. The most common root cause on Windows is that the file
    is currently open in PowerPoint, which holds an exclusive lock AND
    creates a `~$<filename>.pptx` lock-file marker in the same directory.

    If we can extract the path from the exception message and find a
    matching `~$` lock file alongside, we know it's the PowerPoint-open
    case and can emit a precise, actionable error. Otherwise we fall back
    to a generic "not a valid pptx" message that still mentions the
    PowerPoint possibility as the most likely cause.
    """
    # python-pptx's message format: "Package not found at 'path/to/file.pptx'"
    path_match = re.search(r"'([^']+)'", str(exc))
    if path_match:
        path = Path(path_match.group(1))
        lock_file = path.parent / f"~${path.name}"
        if lock_file.exists():
            return (
                f"Template '{path}' is currently open in PowerPoint "
                f"(lock file '{lock_file.name}' detected). "
                "Close the file in PowerPoint and re-run."
            )

    return (
        f"Not a valid pptx file: {exc}. Most likely cause: the file is "
        "currently open in PowerPoint — close it and re-run. Otherwise "
        "the file may be corrupt, empty, or in a format python-pptx "
        "can't read."
    )


# Dispatch table mapping exception classes to (formatter, exit_code) tuples.
# Order matters — more specific classes must come before their superclasses.
# Adding a new trap is a one-line addition to this table.
_Formatter = Callable[[BaseException], str]

_ERROR_HANDLERS: tuple[tuple[type[BaseException], _Formatter, int], ...] = (
    (PermissionError, _format_permission_error, 1),
    (PackageNotFoundError, _format_package_not_found, 1),
    (yaml.YAMLError, lambda exc: f"Invalid YAML: {exc}", 1),
    (FileNotFoundError, lambda exc: str(exc), 1),
    (NotADirectoryError, lambda exc: str(exc), 1),
    (ValueError, lambda exc: str(exc), 1),
)


def main(argv: list[str] | None = None) -> int:
    """Entry point. Accepts optional argv list for programmatic / test use.

    Traps known failure modes via `_ERROR_HANDLERS` and prints a clean error
    message instead of a traceback. Set `RECOMBINASE_DEBUG=1` to restore
    the traceback for debugging.
    """
    debug = os.environ.get("RECOMBINASE_DEBUG") == "1"
    try:
        app(args=argv, standalone_mode=False)
    except typer.Exit as exc:
        return exc.exit_code
    except BaseException as exc:
        for exc_type, formatter, exit_code in _ERROR_HANDLERS:
            if isinstance(exc, exc_type):
                _print_error(formatter(exc))
                if debug:
                    traceback.print_exc()
                return exit_code
        # Unclassified: last-resort guard so a human-facing CLI never shows
        # a raw Python traceback on an unexpected exception class.
        _print_error(f"Unexpected {type(exc).__name__}: {exc}")
        if debug:
            traceback.print_exc()
        else:
            typer.secho(
                "  (set RECOMBINASE_DEBUG=1 to see the full traceback)",
                fg=typer.colors.RED,
                err=True,
            )
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
