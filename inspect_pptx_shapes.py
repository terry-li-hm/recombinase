"""Inspect a pptx template: print every slide, layout, and shape name.

Prints only structural metadata — shape names, types, placeholder info,
and text character counts. Does NOT print actual text content, so it's
safe to share the output.

Usage:
    python inspect-pptx-shapes.py "C:\\path\\to\\your template.pptx"
"""
import sys
from pptx import Presentation


def main():
    if len(sys.argv) < 2:
        print("Usage: python inspect-pptx-shapes.py <path-to-pptx>")
        sys.exit(1)

    path = sys.argv[1]
    prs = Presentation(path)

    print(f"File: {path}")
    print(f"Slide count: {len(prs.slides)}")
    print(f"Slide width x height: {prs.slide_width} x {prs.slide_height} EMU")
    print()

    for slide_index, slide in enumerate(prs.slides, start=1):
        print(f"=== Slide {slide_index} (layout: {slide.slide_layout.name!r}) ===")
        for shape in slide.shapes:
            bits = [f"{shape.name!r}"]
            bits.append(f"type={shape.shape_type}")

            if shape.is_placeholder:
                ph = shape.placeholder_format
                bits.append(f"placeholder(type={ph.type}, idx={ph.idx})")

            if shape.has_text_frame:
                text_len = len(shape.text_frame.text or "")
                bits.append(f"text_chars={text_len}")
                # Report paragraph and run counts — helps diagnose rich-text risk
                para_count = len(shape.text_frame.paragraphs)
                run_count = sum(len(p.runs) for p in shape.text_frame.paragraphs)
                bits.append(f"paras={para_count}, runs={run_count}")

            print("  - " + " | ".join(bits))
        print()

    print("=== Slide Layouts available on master ===")
    for layout_index, layout in enumerate(prs.slide_layouts):
        print(f"  [{layout_index}] {layout.name!r}")


if __name__ == "__main__":
    main()
