"""End-to-end test: create a synthetic template, run inspect/init/generate.

This test verifies the full pipeline without needing any external pptx file.
It builds a small template using python-pptx itself, runs recombinase against
it, and asserts the output pptx has the expected populated slides.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from recombinase.config import TemplateConfig, load_config, write_scaffold_config
from recombinase.generate import (
    duplicate_slide,
    find_shape_by_name,
    generate_deck,
    load_records,
    set_shape_value,
)
from recombinase.inspect import (
    format_template_info,
    inspect_template,
    shape_names_from_slide,
)


def _build_sample_template(path: Path) -> None:
    """Build a small sample pptx template with named shapes."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add a few named text boxes
    def add_textbox(name: str, left: float, top: float, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(5), Inches(1))
        box.name = name
        tf = box.text_frame
        tf.text = text
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)

    add_textbox("Consultant_Name", 0.5, 0.5, "EXAMPLE NAME")
    add_textbox("Role_Title", 0.5, 1.3, "EXAMPLE ROLE")
    add_textbox("Summary_Body", 0.5, 2.1, "Example summary paragraph.")
    add_textbox("Background_Bullets", 0.5, 3.0, "Example bullet")
    add_textbox("Key_Skills", 0.5, 4.2, "Example skill")

    prs.save(str(path))


def _build_sample_data(data_dir: Path) -> None:
    """Write two sample YAML records."""
    data_dir.mkdir(parents=True, exist_ok=True)
    (data_dir / "alpha.yaml").write_text(
        """id: alpha
name: Alpha Jones
role: Senior Consultant
summary: Twelve years across banking analytics and governance.
background:
  - Bank A (2010-2015)
  - Bank B (2015-2020)
  - Bank C (2020-present)
key_skills:
  - Risk modelling
  - Governance
  - Data architecture
""",
        encoding="utf-8",
    )
    (data_dir / "bravo.yaml").write_text(
        """id: bravo
name: Bravo Lee
role: Principal Consultant
summary: AI strategy lead for regulated industries.
background:
  - Consulting firm A (2012-2018)
  - Consulting firm B (2018-present)
key_skills:
  - AI strategy
  - Regulation
""",
        encoding="utf-8",
    )


def _build_sample_config(
    template_path: Path, config_path: Path, output_slide_cleared: bool
) -> None:
    # overflow_ratio: 0 disables the v0.1.6 overflow heuristic so these
    # legacy tests keep their strict `warnings == []` assertions. The
    # overflow behavior itself is covered by test_regressions.py.
    config_yaml = f"""template: {template_path}
source_slide_index: 1
clear_source_slide: {str(output_slide_cleared).lower()}
overflow_ratio: 0
placeholders:
  name: Consultant_Name
  role: Role_Title
  summary: Summary_Body
  background: Background_Bullets
  key_skills: Key_Skills
"""
    config_path.write_text(config_yaml, encoding="utf-8")


def test_inspect_sample_template(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    info = inspect_template(template_path)

    assert info.slide_count == 1
    assert len(info.slides) == 1
    slide = info.slides[0]
    shape_names = [s.name for s in slide.shapes]
    assert "Consultant_Name" in shape_names
    assert "Role_Title" in shape_names
    assert "Background_Bullets" in shape_names


def test_format_template_info_no_text_content(tmp_path: Path) -> None:
    """format_template_info must not leak the actual text of shapes."""
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    info = inspect_template(template_path)
    formatted = format_template_info(info)

    # Structural metadata should be present
    assert "Consultant_Name" in formatted
    assert "text_chars=" in formatted
    # But the actual example text must NOT be in the output
    assert "EXAMPLE NAME" not in formatted
    assert "Example bullet" not in formatted
    assert "Example summary paragraph." not in formatted


def test_write_scaffold_config(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    info = inspect_template(template_path)
    names = shape_names_from_slide(info, 1)

    out = tmp_path / "scaffold.yaml"
    write_scaffold_config(template_path, names, out)

    text = out.read_text(encoding="utf-8")
    assert "placeholders:" in text
    assert "Consultant_Name" in text
    assert str(template_path) in text


def test_load_config_happy_path(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    config_path = tmp_path / "config.yaml"
    _build_sample_config(template_path, config_path, output_slide_cleared=True)

    config = load_config(config_path)

    assert isinstance(config, TemplateConfig)
    assert config.template == template_path.resolve()
    assert config.source_slide_index == 1
    assert config.clear_source_slide is True
    assert "name" in config.placeholders
    assert config.placeholders["background"] == "Background_Bullets"


def test_load_config_missing_template_errors(tmp_path: Path) -> None:
    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        """template: /nonexistent/path.pptx
placeholders:
  name: Consultant_Name
""",
        encoding="utf-8",
    )
    with pytest.raises(ValueError, match="Template file not found"):
        load_config(config_path)


def test_load_records_happy_path(tmp_path: Path) -> None:
    data_dir = tmp_path / "data"
    _build_sample_data(data_dir)

    records = load_records(data_dir)

    assert len(records) == 2
    # Sorted by filename: alpha, bravo
    assert records[0]["id"] == "alpha"
    assert records[1]["id"] == "bravo"
    assert isinstance(records[0]["background"], list)
    assert len(records[0]["background"]) == 3


def test_set_shape_value_scalar(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shape = find_shape_by_name(slide, "Consultant_Name")
    assert shape is not None

    set_shape_value(shape, "Real Name Here")
    assert shape.text_frame.text == "Real Name Here"


def test_set_shape_value_list_becomes_paragraphs(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shape = find_shape_by_name(slide, "Background_Bullets")
    assert shape is not None

    set_shape_value(shape, ["Bullet one", "Bullet two", "Bullet three"])
    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 3
    assert paragraphs[0].text == "Bullet one"
    assert paragraphs[1].text == "Bullet two"
    assert paragraphs[2].text == "Bullet three"


def test_set_shape_value_newline_string_becomes_paragraphs(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shape = find_shape_by_name(slide, "Background_Bullets")
    assert shape is not None

    set_shape_value(shape, "Line A\nLine B")
    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 2
    assert paragraphs[0].text == "Line A"
    assert paragraphs[1].text == "Line B"


def test_duplicate_slide_preserves_shapes(tmp_path: Path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    prs = Presentation(str(template_path))
    source = prs.slides[0]
    original_shape_count = len(source.shapes)
    original_names = sorted(s.name for s in source.shapes)

    new_slide = duplicate_slide(prs, source)

    new_names = sorted(s.name for s in new_slide.shapes)
    assert len(new_slide.shapes) == original_shape_count
    assert new_names == original_names
    # Source slide is untouched
    assert len(prs.slides) == 2


def test_generate_deck_end_to_end(tmp_path: Path) -> None:
    """Full pipeline: template + config + data → output pptx."""
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    data_dir = tmp_path / "data"
    _build_sample_data(data_dir)

    config_path = tmp_path / "config.yaml"
    _build_sample_config(template_path, config_path, output_slide_cleared=True)

    output_path = tmp_path / "output.pptx"

    config = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(config, records, output_path)

    assert output_path.exists()
    assert result["records_generated"] == 2
    # No warnings expected — all placeholders should match
    assert result["warnings"] == []

    # Re-open the generated deck and verify the content
    out_prs = Presentation(str(output_path))
    # Source slide was cleared, so only 2 generated slides remain
    assert len(out_prs.slides) == 2

    slide0 = out_prs.slides[0]
    name_shape = find_shape_by_name(slide0, "Consultant_Name")
    role_shape = find_shape_by_name(slide0, "Role_Title")
    background_shape = find_shape_by_name(slide0, "Background_Bullets")

    assert name_shape is not None
    assert role_shape is not None
    assert background_shape is not None

    assert name_shape.text_frame.text == "Alpha Jones"
    assert role_shape.text_frame.text == "Senior Consultant"
    # Background should have 3 bullet paragraphs
    bg_paragraphs = background_shape.text_frame.paragraphs
    assert len(bg_paragraphs) == 3
    assert bg_paragraphs[0].text == "Bank A (2010-2015)"


def test_generate_deck_preserves_source_slide_when_configured(
    tmp_path: Path,
) -> None:
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    data_dir = tmp_path / "data"
    _build_sample_data(data_dir)

    config_path = tmp_path / "config.yaml"
    _build_sample_config(template_path, config_path, output_slide_cleared=False)

    output_path = tmp_path / "output.pptx"

    config = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(config, records, output_path)

    assert result["records_generated"] == 2

    out_prs = Presentation(str(output_path))
    # Source slide kept + 2 generated = 3 total
    assert len(out_prs.slides) == 3


def test_generate_deck_missing_placeholder_warns(tmp_path: Path) -> None:
    """If the config references a shape name that doesn't exist, it should warn."""
    template_path = tmp_path / "template.pptx"
    _build_sample_template(template_path)

    data_dir = tmp_path / "data"
    _build_sample_data(data_dir)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        f"""template: {template_path}
source_slide_index: 1
clear_source_slide: true
placeholders:
  name: Consultant_Name
  nonexistent_field: Nonexistent_Shape
""",
        encoding="utf-8",
    )

    output_path = tmp_path / "output.pptx"

    config = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(config, records, output_path)

    assert result["records_generated"] == 2
    assert len(result["warnings"]) > 0
    assert any("Nonexistent_Shape" in w for w in result["warnings"])
