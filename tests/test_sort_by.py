"""Tests for the sort_by feature — config-driven record ordering."""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml

from recombinase.config import load_config
from recombinase.generate import load_records


def _write_yaml(path: Path, data: dict) -> None:
    path.write_text(yaml.dump(data, sort_keys=False), encoding="utf-8")


def _make_data_dir(tmp_path: Path, records: list[tuple[str, dict]]) -> Path:
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    for filename, record in records:
        _write_yaml(data_dir / filename, record)
    return data_dir


# ---------------------------------------------------------------------------
# load_records with sort_by
# ---------------------------------------------------------------------------


class TestLoadRecordsSortBy:
    """Verify sort_by parameter in load_records."""

    def test_sort_by_numeric_field(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("c-charlie.yaml", {"name": "Charlie", "rank": 3}),
                ("a-alice.yaml", {"name": "Alice", "rank": 1}),
                ("b-bob.yaml", {"name": "Bob", "rank": 2}),
            ],
        )
        records = load_records(data_dir, sort_by="rank")
        names = [rec["name"] for rec in records]
        assert names == ["Alice", "Bob", "Charlie"]

    def test_sort_by_string_field(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("03.yaml", {"name": "Charlie", "tier": "c-consultant"}),
                ("01.yaml", {"name": "Alice", "tier": "a-partner"}),
                ("02.yaml", {"name": "Bob", "tier": "b-principal"}),
            ],
        )
        records = load_records(data_dir, sort_by="tier")
        names = [rec["name"] for rec in records]
        assert names == ["Alice", "Bob", "Charlie"]

    def test_missing_field_sorts_last(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("a.yaml", {"name": "Alice", "rank": 2}),
                ("b.yaml", {"name": "Bob"}),  # no rank
                ("c.yaml", {"name": "Charlie", "rank": 1}),
            ],
        )
        records = load_records(data_dir, sort_by="rank")
        names = [rec["name"] for rec in records]
        assert names == ["Charlie", "Alice", "Bob"]

    def test_equal_values_preserve_filename_order(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("b-bob.yaml", {"name": "Bob", "rank": 1}),
                ("a-alice.yaml", {"name": "Alice", "rank": 1}),
                ("c-charlie.yaml", {"name": "Charlie", "rank": 1}),
            ],
        )
        records = load_records(data_dir, sort_by="rank")
        names = [rec["name"] for rec in records]
        # Filename order: a-alice, b-bob, c-charlie (all rank 1)
        assert names == ["Alice", "Bob", "Charlie"]

    def test_no_sort_by_uses_filename_order(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("c.yaml", {"name": "Charlie", "rank": 1}),
                ("a.yaml", {"name": "Alice", "rank": 3}),
                ("b.yaml", {"name": "Bob", "rank": 2}),
            ],
        )
        records = load_records(data_dir, sort_by=None)
        names = [rec["name"] for rec in records]
        assert names == ["Alice", "Bob", "Charlie"]  # filename order

    def test_sort_by_float_field(self, tmp_path: Path) -> None:
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("a.yaml", {"name": "Alice", "priority": 2.5}),
                ("b.yaml", {"name": "Bob", "priority": 1.0}),
            ],
        )
        records = load_records(data_dir, sort_by="priority")
        names = [rec["name"] for rec in records]
        assert names == ["Bob", "Alice"]

    def test_sort_by_mixed_present_absent(self, tmp_path: Path) -> None:
        """Multiple records missing the field all sort last, in filename order."""
        data_dir = _make_data_dir(
            tmp_path,
            [
                ("c.yaml", {"name": "Charlie"}),
                ("a.yaml", {"name": "Alice", "rank": 1}),
                ("b.yaml", {"name": "Bob"}),
            ],
        )
        records = load_records(data_dir, sort_by="rank")
        names = [rec["name"] for rec in records]
        assert names == ["Alice", "Bob", "Charlie"]


# ---------------------------------------------------------------------------
# Config parsing
# ---------------------------------------------------------------------------


class TestConfigSortBy:
    """Verify sort_by is parsed from config YAML."""

    def _write_config(self, tmp_path: Path, extra: dict | None = None) -> Path:
        template = tmp_path / "template.pptx"
        # Create a minimal valid pptx so config validation passes
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(template))

        config_data: dict = {
            "template": "./template.pptx",
            "placeholders": {"name": "name"},
        }
        if extra:
            config_data.update(extra)
        config_path = tmp_path / "config.yaml"
        _write_yaml(config_path, config_data)
        return config_path

    def test_sort_by_parsed(self, tmp_path: Path) -> None:
        config_path = self._write_config(tmp_path, {"sort_by": "rank"})
        cfg = load_config(config_path)
        assert cfg.sort_by == "rank"

    def test_sort_by_default_none(self, tmp_path: Path) -> None:
        config_path = self._write_config(tmp_path)
        cfg = load_config(config_path)
        assert cfg.sort_by is None

    def test_sort_by_non_string_raises(self, tmp_path: Path) -> None:
        config_path = self._write_config(tmp_path, {"sort_by": 42})
        with pytest.raises(ValueError, match=r"sort_by.*string"):
            load_config(config_path)
