"""v0.1.9 robustness tests — close the remaining coverage gaps from the v0.1.5 review.

Each test here exists to catch a specific regression that would otherwise
slip through. These are paths the existing suite reached only indirectly
or not at all.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches
from typer.testing import CliRunner

from recombinase.cli import main
from recombinase.config import load_config
from recombinase.generate import (
    duplicate_slide,
    find_shape_by_name,
    load_records,
    remove_slide,
    set_shape_value,
)

runner = CliRunner()


# -- load_records error branches ------------------------------------------


def test_load_records_raises_when_dir_missing(tmp_path: Path) -> None:
    with pytest.raises(FileNotFoundError, match="Data directory not found"):
        load_records(tmp_path / "nope")


def test_load_records_raises_when_path_is_file(tmp_path: Path) -> None:
    file_path = tmp_path / "notadir.txt"
    file_path.write_text("hello", encoding="utf-8")
    with pytest.raises(NotADirectoryError, match="Not a directory"):
        load_records(file_path)


def test_load_records_raises_on_non_dict_yaml_record(tmp_path: Path) -> None:
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "bad.yaml").write_text("- just\n- a list\n", encoding="utf-8")
    with pytest.raises(ValueError, match="expected top-level mapping"):
        load_records(data_dir)


def test_load_records_silently_skips_empty_yaml(tmp_path: Path) -> None:
    """Empty YAML files (comment-only, whitespace-only) should be skipped, not errored."""
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "empty.yaml").write_text("", encoding="utf-8")
    (data_dir / "comment-only.yaml").write_text("# just a comment\n", encoding="utf-8")
    (data_dir / "real.yaml").write_text(
        yaml.safe_dump({"name": "only real record"}), encoding="utf-8"
    )

    records = load_records(data_dir)
    assert len(records) == 1
    assert records[0]["name"] == "only real record"


def test_load_records_picks_up_yml_extension(tmp_path: Path) -> None:
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "a.yaml").write_text(yaml.safe_dump({"id": "a"}), encoding="utf-8")
    (data_dir / "b.yml").write_text(yaml.safe_dump({"id": "b"}), encoding="utf-8")

    records = load_records(data_dir)
    ids = sorted(r["id"] for r in records)
    assert ids == ["a", "b"]


# -- remove_slide direct test ---------------------------------------------


def test_remove_slide_drops_slide_and_relationship(simple_template: Path) -> None:
    """Directly verify remove_slide drops both the sldId entry and the rel."""
    prs = Presentation(str(simple_template))

    # Add a second slide so we can remove one safely
    new_slide = duplicate_slide(prs, prs.slides[0])
    assert len(prs.slides) == 2

    target_slide_id = prs.slides[0].slide_id

    # Count relationships on the presentation part before
    pres_rels_before = len(prs.part.rels)

    remove_slide(prs, prs.slides[0])

    # Slide count dropped
    assert len(prs.slides) == 1
    # The remaining slide is the one we duplicated
    assert prs.slides[0].slide_id == new_slide.slide_id
    # Removed slide's slide_id is not present anywhere in the sldIdLst
    remaining_ids = [s.slide_id for s in prs.slides]
    assert target_slide_id not in remaining_ids
    # Presentation rels count dropped (at least the removed slide's rel)
    assert len(prs.part.rels) < pres_rels_before


# -- duplicate_slide: external hyperlink rel path -------------------------


def test_duplicate_slide_preserves_external_hyperlink(tmp_path: Path) -> None:
    """A hyperlink on the source slide is an external rel — must be preserved."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "LinkedText"
    tb.text_frame.text = "Portfolio"
    # Add a hyperlink to the first run
    tb.text_frame.paragraphs[0].runs[0].hyperlink.address = "https://example.com/portfolio"

    source_path = tmp_path / "with_link.pptx"
    prs.save(str(source_path))

    # Now duplicate the slide and round-trip
    prs = Presentation(str(source_path))
    source = prs.slides[0]
    source_ext_rels = [rel for rel in source.part.rels.values() if rel.is_external]
    assert source_ext_rels, "source fixture should have an external rel"

    new_slide = duplicate_slide(prs, source)
    new_ext_rels = [rel for rel in new_slide.part.rels.values() if rel.is_external]
    assert len(new_ext_rels) == len(source_ext_rels)
    assert {r.target_ref for r in new_ext_rels} == {r.target_ref for r in source_ext_rels}

    # Round-trip: save and verify the external rel survives
    output = tmp_path / "roundtrip.pptx"
    prs.save(str(output))
    reopened = Presentation(str(output))
    dup_slide = reopened.slides[1]
    reopened_ext = [rel for rel in dup_slide.part.rels.values() if rel.is_external]
    assert any("example.com/portfolio" in r.target_ref for r in reopened_ext)


# -- duplicate_slide: notesSlide skip -------------------------------------


def test_duplicate_slide_does_not_copy_notes(tmp_path: Path) -> None:
    """Source-slide notes should NOT follow the slide into its duplicate."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"

    # Add notes to the source slide
    slide.notes_slide.notes_text_frame.text = "Private presenter notes"

    path = tmp_path / "with_notes.pptx"
    prs.save(str(path))

    prs = Presentation(str(path))
    source = prs.slides[0]
    # Confirm fixture has notesSlide rel
    source_has_notes_rel = any("notesSlide" in rel.reltype for rel in source.part.rels.values())
    assert source_has_notes_rel, "source should have notes rel"

    new_slide = duplicate_slide(prs, source)
    new_has_notes_rel = any("notesSlide" in rel.reltype for rel in new_slide.part.rels.values())
    assert not new_has_notes_rel, "duplicated slide should not inherit source's notes"


# -- _walk_shapes: nested groups ------------------------------------------


def test_walk_shapes_recurses_through_nested_groups(tmp_path: Path) -> None:
    """Group inside a group — _walk_shapes recursion should reach every level."""
    from copy import deepcopy

    from lxml import etree
    from pptx.oxml.ns import qn

    from recombinase.generate import _walk_shapes

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    inner_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    inner_tb.name = "Deeply_Nested"
    inner_tb.text_frame.text = "EXAMPLE"

    # Build outer grpSp → inner grpSp → inner_tb
    spTree = slide.shapes._spTree

    def make_grp(name: str, parent):
        grp = etree.SubElement(parent, qn("p:grpSp"))
        nv = etree.SubElement(grp, qn("p:nvGrpSpPr"))
        cNv = etree.SubElement(nv, qn("p:cNvPr"))
        cNv.set("id", "777")
        cNv.set("name", name)
        etree.SubElement(nv, qn("p:cNvGrpSpPr"))
        etree.SubElement(nv, qn("p:nvPr"))
        grpPr = etree.SubElement(grp, qn("p:grpSpPr"))
        xfrm = etree.SubElement(grpPr, qn("a:xfrm"))
        for tag in ("a:off", "a:ext", "a:chOff", "a:chExt"):
            el = etree.SubElement(xfrm, qn(tag))
            if "ext" in tag.lower() or tag.split(":")[1] == "chExt":
                el.set("cx", "914400")
                el.set("cy", "685800")
            else:
                el.set("x", "0")
                el.set("y", "0")
        return grp

    outer = make_grp("Outer_Group", spTree)
    inner = make_grp("Inner_Group", outer)

    # Move inner_tb into the inner group
    sp = inner_tb._element
    parent = sp.getparent()
    sp_copy = deepcopy(sp)
    parent.remove(sp)
    inner.append(sp_copy)

    path = tmp_path / "nested_group.pptx"
    prs.save(str(path))

    prs = Presentation(str(path))
    slide = prs.slides[0]
    names = [s.name for s in _walk_shapes(slide.shapes)]

    assert "Outer_Group" in names
    assert "Inner_Group" in names
    assert "Deeply_Nested" in names

    # And find_shape_by_name should reach the deeply-nested one
    assert find_shape_by_name(slide, "Deeply_Nested") is not None


# -- set_shape_value: bool type ------------------------------------------


@pytest.mark.parametrize(
    ("value", "expected"),
    [
        (True, "True"),
        (False, ""),  # False is falsy in `value == ""` path? No — but value is not "" or None
    ],
)
def test_set_shape_value_handles_bool(simple_template: Path, value: bool, expected: str) -> None:
    """bool is an int subclass in Python — verify it stringifies predictably."""
    prs = Presentation(str(simple_template))
    shape = find_shape_by_name(prs.slides[0], "Field")
    assert shape is not None

    set_shape_value(shape, value)
    # Current implementation: `if value is None or value == ""` catches None and "".
    # False doesn't match either (False != "" and False is not None), so it
    # falls through to `str(value)` → "False". True → "True".
    assert shape.text_frame.text == ("True" if value else "False")


# -- load_config: relative template path ---------------------------------


def test_load_config_resolves_relative_template_path(simple_template: Path, tmp_path: Path) -> None:
    """Relative template path in config should resolve against config file's directory."""
    # Put the config in tmp_path, template in tmp_path/template/
    template_dir = tmp_path / "template"
    template_dir.mkdir()
    relative_template = template_dir / "t.pptx"
    shutil.copy(simple_template, relative_template)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        "template: template/t.pptx\nplaceholders:\n  name: Field\n",
        encoding="utf-8",
    )

    cfg = load_config(config_path)
    assert cfg.template == relative_template.resolve()
    assert cfg.template.exists()


# -- main() exception handler: PermissionError ---------------------------


def test_main_handles_permission_error_with_hint(
    tmp_path: Path, simple_template: Path, write_config, monkeypatch, capsys
) -> None:
    """Force a PermissionError during save to exercise the PowerPoint-lock message."""
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=0)
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r.yaml").write_text(yaml.safe_dump({"name": "A"}), encoding="utf-8")
    output = tmp_path / "locked.pptx"

    # Monkeypatch the Presentation class's save method. `from pptx import
    # Presentation` gives the factory function; the real class with `.save`
    # lives in pptx.presentation.
    from pptx.presentation import Presentation as _PresClass

    def fake_save(self, path: str) -> None:
        raise PermissionError(13, "locked", str(path))

    monkeypatch.setattr(_PresClass, "save", fake_save)

    rc = main(
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
            "--force",
        ]
    )

    assert rc == 1
    captured = capsys.readouterr()
    combined = captured.out + captured.err
    assert "Error:" in combined
    # PowerPoint-lock hint text
    assert "PowerPoint" in combined or "Cannot write" in combined


# -- main() exception handler: yaml.YAMLError ----------------------------


def test_main_handles_yaml_error_cleanly(simple_template: Path, tmp_path: Path, capsys) -> None:
    """Malformed YAML config should produce a clean 'Invalid YAML' message."""
    config_path = tmp_path / "broken.yaml"
    config_path.write_text("template: [unclosed\nplaceholders:\n  name: Field\n", encoding="utf-8")

    rc = main(
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(tmp_path),
            "-o",
            str(tmp_path / "out.pptx"),
        ]
    )
    # Depending on where it trips: YAML parse error becomes ValueError via
    # load_config wrapping, or yaml.YAMLError surfaces directly. Either
    # routes through main()'s error dispatch.
    assert rc == 1
    captured = capsys.readouterr()
    combined = captured.out + captured.err
    assert "Error:" in combined
    assert "Traceback" not in combined
