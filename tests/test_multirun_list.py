"""Tests for list-value → per-run replacement in multi-run shapes.

When a shape has a single paragraph with multiple runs (mixed formatting),
a list data value maps each item to the corresponding run, preserving
each run's rPr (font, size, bold, color, italic).
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from recombinase.generate import set_shape_value


def _make_multirun_shape(prs, runs_spec):
    """Create a textbox with multiple runs in a single paragraph.

    runs_spec: list of (text, bold, rgb_hex) tuples.
    Returns the shape.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
    tf = shape.text_frame
    paragraph = tf.paragraphs[0]
    # First run uses the existing paragraph
    for idx, (text, bold, rgb_hex) in enumerate(runs_spec):
        if idx == 0:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = text
        else:
            run = paragraph.add_run()
            run.text = text
        run.font.bold = bold
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor.from_string(rgb_hex)
    return shape


def test_list_value_replaces_per_run():
    """List with same length as runs replaces text per-run."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    shape = _make_multirun_shape(prs, [
        ("Bold part. ", True, "1A1A1A"),
        ("Grey part.", False, "6B6B6B"),
    ])

    set_shape_value(shape, ["New bold text. ", "New grey text."])

    paragraph = shape.text_frame.paragraphs[0]
    runs = paragraph.runs
    assert len(runs) == 2
    assert runs[0].text == "New bold text. "
    assert runs[0].font.bold is True
    assert str(runs[0].font.color.rgb) == "1A1A1A"
    assert runs[1].text == "New grey text."
    assert runs[1].font.bold is False
    assert str(runs[1].font.color.rgb) == "6B6B6B"


def test_list_fewer_items_than_runs():
    """Fewer items than runs removes excess runs."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    shape = _make_multirun_shape(prs, [
        ("Part one. ", True, "1A1A1A"),
        ("Part two. ", False, "6B6B6B"),
        ("Part three.", False, "999999"),
    ])

    set_shape_value(shape, ["Only this."])

    paragraph = shape.text_frame.paragraphs[0]
    runs = paragraph.runs
    assert len(runs) == 1
    assert runs[0].text == "Only this."
    assert runs[0].font.bold is True


def test_list_more_items_than_runs():
    """More items than runs appends new runs cloning last run's format."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    shape = _make_multirun_shape(prs, [
        ("Bold. ", True, "1A1A1A"),
        ("Grey.", False, "6B6B6B"),
    ])

    set_shape_value(shape, ["One. ", "Two. ", "Three."])

    paragraph = shape.text_frame.paragraphs[0]
    runs = paragraph.runs
    assert len(runs) == 3
    assert runs[0].text == "One. "
    assert runs[1].text == "Two. "
    assert runs[2].text == "Three."
    # Extra run inherits last source run's format
    assert runs[2].font.bold is False


def test_scalar_value_still_flattens():
    """Scalar value on a multi-run shape still works (existing behavior)."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    shape = _make_multirun_shape(prs, [
        ("Bold. ", True, "1A1A1A"),
        ("Grey.", False, "6B6B6B"),
    ])

    set_shape_value(shape, "Single string replaces all")

    paragraph = shape.text_frame.paragraphs[0]
    assert paragraph.text == "Single string replaces all"


def test_single_run_list_still_makes_paragraphs():
    """List value on a single-run shape still creates paragraphs (bullets)."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
    shape.text_frame.paragraphs[0].text = "Single run"

    set_shape_value(shape, ["Bullet one", "Bullet two", "Bullet three"])

    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 3
    assert paragraphs[0].text == "Bullet one"
    assert paragraphs[1].text == "Bullet two"
    assert paragraphs[2].text == "Bullet three"
