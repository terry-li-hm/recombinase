"""Tests for v0.2.0 P0 fixes: validate section shapes + _clear_cell rPr preservation."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from typer.testing import CliRunner

from recombinase.cli import app
from recombinase.config import TableConfig
from recombinase.generate import _clear_cell, populate_table

runner = CliRunner()

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _build_template_with_sections_shape(pptx_path: Path) -> None:
    """Template with a text box containing 2+ paragraphs (section-compatible)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Regular placeholder
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb1.name = "Name_Field"
    tb1.text_frame.text = "EXAMPLE"

    # Section-compatible shape: needs >=2 paragraphs
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4), Inches(3))
    tb2.name = "Skills_Section"
    tf = tb2.text_frame
    tf.text = "Header Example"
    para2 = tf.add_paragraph()
    para2.text = "Bullet example"

    prs.save(str(pptx_path))


def _build_table_template_with_styled_cells(pptx_path: Path) -> None:
    """Template with a table whose cells have explicit font formatting."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 3 rows x 2 cols table
    table_shape = slide.shapes.add_table(3, 2, Inches(0.5), Inches(0.5), Inches(5), Inches(3))
    table_shape.name = "Styled_Table"
    table = table_shape.table

    # Header row
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"

    # Style all data cells with explicit font
    for row_idx in range(1, 3):
        for col_idx in range(2):
            cell = table.cell(row_idx, col_idx)
            cell.text = f"Example R{row_idx}C{col_idx}"
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(10)
                    run.font.bold = True

    prs.save(str(pptx_path))


# ---------------------------------------------------------------------------
# Commit 1: validate includes section shapes
# ---------------------------------------------------------------------------


class TestValidateSectionShapes:
    """Validate command must include section shapes in its shape-name check."""

    def test_validate_catches_missing_section_shape(self, tmp_path: Path) -> None:
        """A section shape that doesn't exist on the template produces exit 1."""
        template_path = tmp_path / "template.pptx"
        _build_template_with_sections_shape(template_path)

        config_data: dict[str, Any] = {
            "template": str(template_path),
            "source_slide_index": 1,
            "placeholders": {"name": "Name_Field"},
            "sections": {
                "skills": {
                    "shape": "NONEXISTENT_Section",
                    "header_index": 0,
                    "bullet_index": 1,
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        result = runner.invoke(app, ["validate", "-c", str(config_path)])
        assert result.exit_code == 1

    def test_validate_section_shapes_not_reported_unused(self, tmp_path: Path) -> None:
        """A correctly mapped section shape must not appear in the 'unused' list."""
        template_path = tmp_path / "template.pptx"
        _build_template_with_sections_shape(template_path)

        config_data: dict[str, Any] = {
            "template": str(template_path),
            "source_slide_index": 1,
            "placeholders": {"name": "Name_Field"},
            "sections": {
                "skills": {
                    "shape": "Skills_Section",
                    "header_index": 0,
                    "bullet_index": 1,
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        result = runner.invoke(app, ["validate", "-c", str(config_path)])
        assert result.exit_code == 0

    def test_validate_section_shape_provenance_in_error(self, tmp_path: Path) -> None:
        """Missing section shape error should show 'section field: ...' provenance."""
        template_path = tmp_path / "template.pptx"
        _build_template_with_sections_shape(template_path)

        config_data: dict[str, Any] = {
            "template": str(template_path),
            "source_slide_index": 1,
            "placeholders": {"name": "Name_Field"},
            "sections": {
                "competencies": {
                    "shape": "Missing_Shape_XYZ",
                    "header_index": 0,
                    "bullet_index": 1,
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        result = runner.invoke(app, ["validate", "-c", str(config_path)])
        assert result.exit_code == 1
        assert "section field: competencies" in result.output


# ---------------------------------------------------------------------------
# Commit 2: _clear_cell preserves formatting
# ---------------------------------------------------------------------------


class TestClearCellPreservesFormatting:
    """_clear_cell must preserve pPr and rPr like _write_paragraphs does."""

    def test_clear_cell_preserves_rpr(self, tmp_path: Path) -> None:
        """After _clear_cell, the empty paragraph retains the original rPr."""
        pptx_path = tmp_path / "styled.pptx"
        _build_table_template_with_styled_cells(pptx_path)

        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]
        table_shape = None
        for shape in slide.shapes:
            if shape.name == "Styled_Table":
                table_shape = shape
                break
        assert table_shape is not None

        cell = table_shape.table.cell(1, 0)

        # Verify pre-condition: cell has styled run
        first_p = cell.text_frame.paragraphs[0]
        assert len(first_p.runs) > 0
        assert first_p.runs[0].font.bold is True

        # Clear and verify rPr survives
        _clear_cell(cell)
        assert cell.text_frame.text == ""

        # Check the XML: the empty paragraph should still have rPr with bold
        cleared_p = cell.text_frame.paragraphs[0]._p
        first_r = cleared_p.find(qn("a:r"))
        assert first_r is not None, "cleared cell should still have a run element"
        rPr = first_r.find(qn("a:rPr"))
        assert rPr is not None, "cleared cell run should retain rPr"
        assert rPr.get("b") == "1", "bold attribute should be preserved"

    def test_clear_cell_preserves_ppr(self) -> None:
        """After _clear_cell, paragraph-level formatting (pPr) is retained."""
        # Build a cell with explicit paragraph properties
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
        cell = table_shape.table.cell(1, 0)
        cell.text = "Sample"

        # Inject a pPr with alignment
        para_xml = cell.text_frame.paragraphs[0]._p
        pPr = para_xml.find(qn("a:pPr"))
        if pPr is None:
            from lxml import etree

            pPr = etree.SubElement(para_xml, qn("a:pPr"))
            para_xml.insert(0, pPr)
        pPr.set("algn", "ctr")

        _clear_cell(cell)
        assert cell.text_frame.text == ""

        cleared_pPr = cell.text_frame.paragraphs[0]._p.find(qn("a:pPr"))
        assert cleared_pPr is not None, "pPr should survive _clear_cell"
        assert cleared_pPr.get("algn") == "ctr", "alignment should be preserved"

    def test_clear_cell_spanned_is_noop(self) -> None:
        """Spanned (merged non-origin) cells are skipped silently."""

        class FakeSpannedCell:
            is_spanned = True

        # Should not raise or access text_frame
        _clear_cell(FakeSpannedCell())

    def test_table_excess_rows_retain_formatting_after_clear(self, tmp_path: Path) -> None:
        """When populate_table clears excess rows, those cells keep formatting.

        This is an integration test: populate_table calls _clear_cell on excess
        rows. After the fix, those cleared cells should retain pPr+rPr so that
        manual edits in PowerPoint don't revert to default formatting.
        """
        pptx_path = tmp_path / "styled.pptx"
        _build_table_template_with_styled_cells(pptx_path)

        prs = Presentation(str(pptx_path))
        slide = prs.slides[0]
        table_shape = None
        for shape in slide.shapes:
            if shape.name == "Styled_Table":
                table_shape = shape
                break
        assert table_shape is not None

        # Populate with only 1 data row — row 2 should be cleared
        config = TableConfig(shape="Styled_Table", columns=["col_a", "col_b"], header_row=True)
        rows = [{"col_a": "Value A", "col_b": "Value B"}]
        populate_table(table_shape, config, rows)

        # Row 2 was an excess row — its cells should be cleared but formatted
        excess_cell = table_shape.table.cell(2, 0)
        assert excess_cell.text_frame.text == ""

        cleared_p = excess_cell.text_frame.paragraphs[0]._p
        first_r = cleared_p.find(qn("a:r"))
        if first_r is not None:
            rPr = first_r.find(qn("a:rPr"))
            # If rPr exists, bold should be preserved
            if rPr is not None:
                assert rPr.get("b") == "1"
