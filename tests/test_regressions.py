"""Regression tests for v0.1.5 review-driven fixes.

Covers the bugs surfaced by the multi-reviewer code review:
- duplicate_slide must preserve shape relationships (rIds) so pictures,
  hyperlinks, and embedded charts survive the copy
- find_shape_by_name must recurse into group shapes
- set_shape_value must handle all scalar types and edge cases
- load_config must raise clean ValueError on malformed YAML / wrong types
- CLI generate must create parent directories and refuse to clobber
- CLI main() must trap known exception classes with clean error messages
"""

from __future__ import annotations

import struct
import zlib
from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from typer.testing import CliRunner

from recombinase.cli import app, main
from recombinase.config import load_config
from recombinase.generate import (
    _walk_shapes,
    duplicate_slide,
    find_shape_by_name,
    generate_deck,
    load_records,
    set_shape_value,
)
from recombinase.inspect import format_template_info, inspect_template

runner = CliRunner()


# -- test fixture builders -------------------------------------------------


def _tiny_png_bytes() -> bytes:
    """A 1x1 red PNG for test fixtures, self-contained so tests need no assets."""
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = b"\x00\x00\x00\rIHDR" + struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = struct.pack(">I", zlib.crc32(ihdr[4:]) & 0xFFFFFFFF)
    idat_data = zlib.compress(b"\x00\xff\x00\x00")
    idat = struct.pack(">I", len(idat_data)) + b"IDAT" + idat_data
    idat_crc = struct.pack(">I", zlib.crc32(b"IDAT" + idat_data) & 0xFFFFFFFF)
    iend = b"\x00\x00\x00\x00IEND" + struct.pack(">I", zlib.crc32(b"IEND") & 0xFFFFFFFF)
    return sig + ihdr + ihdr_crc + idat + idat_crc + iend


def _build_template_with_picture(pptx_path: Path, png_path: Path) -> None:
    """Template with a text box and a picture (to exercise rels duplication)."""
    png_path.write_bytes(_tiny_png_bytes())
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "Consultant_Name"
    tb.text_frame.text = "EXAMPLE"
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(14)

    pic = slide.shapes.add_picture(str(png_path), Inches(5), Inches(0.5), Inches(1), Inches(1))
    pic.name = "Headshot"

    prs.save(str(pptx_path))


def _build_template_with_group(pptx_path: Path) -> None:
    """Template with two text boxes that will be grouped via direct XML manipulation.

    python-pptx doesn't expose a public "group shapes" API, so we build a
    minimal grpSp element and reparent two existing shapes into it. This
    exercises the group-recursion code path without depending on PowerPoint.
    """
    from copy import deepcopy

    from lxml import etree
    from pptx.oxml.ns import qn

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb1.name = "Grouped_Name"
    tb1.text_frame.text = "EXAMPLE"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(1))
    tb2.name = "Grouped_Role"
    tb2.text_frame.text = "EXAMPLE"

    # Build a <p:grpSp> element and move tb1/tb2 into it.
    spTree = slide.shapes._spTree
    grpSp = etree.SubElement(spTree, qn("p:grpSp"))
    nvGrpSpPr = etree.SubElement(grpSp, qn("p:nvGrpSpPr"))
    cNvPr = etree.SubElement(nvGrpSpPr, qn("p:cNvPr"))
    cNvPr.set("id", "999")
    cNvPr.set("name", "Consultant_Card")
    etree.SubElement(nvGrpSpPr, qn("p:cNvGrpSpPr"))
    etree.SubElement(nvGrpSpPr, qn("p:nvPr"))
    grpSpPr = etree.SubElement(grpSp, qn("p:grpSpPr"))
    xfrm = etree.SubElement(grpSpPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", "0")
    off.set("y", "0")
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", "9144000")
    ext.set("cy", "6858000")
    chOff = etree.SubElement(xfrm, qn("a:chOff"))
    chOff.set("x", "0")
    chOff.set("y", "0")
    chExt = etree.SubElement(xfrm, qn("a:chExt"))
    chExt.set("cx", "9144000")
    chExt.set("cy", "6858000")

    # Move tb1 and tb2 into the group.
    for tb in (tb1, tb2):
        sp = tb._element
        parent = sp.getparent()
        sp_copy = deepcopy(sp)
        parent.remove(sp)
        grpSp.append(sp_copy)

    prs.save(str(pptx_path))


def _simple_template(pptx_path: Path) -> None:
    """Plain text-only template used by set_shape_value tests."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(2))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"
    prs.save(str(pptx_path))


def _write_config_yaml(
    config_path: Path, template_path: Path, placeholders: dict[str, str]
) -> None:
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "placeholders": placeholders,
            }
        ),
        encoding="utf-8",
    )


# -- duplicate_slide rels preservation (P0) --------------------------------


def test_duplicate_slide_copies_image_relationship(tmp_path: Path) -> None:
    """A picture on the source slide must survive duplication with its image intact."""
    template_path = tmp_path / "with_pic.pptx"
    png_path = tmp_path / "dot.png"
    _build_template_with_picture(template_path, png_path)

    prs = Presentation(str(template_path))
    source = prs.slides[0]

    # Source has a slideLayout rel + an image rel. After duplication the new
    # slide's rels must include both reltypes, and the image rel must target
    # the same image part.
    source_image_targets = {
        rel.target_part.partname for rel in source.part.rels.values() if "image" in rel.reltype
    }
    assert source_image_targets, "fixture should have at least one image rel"

    new_slide = duplicate_slide(prs, source)

    new_image_targets = {
        rel.target_part.partname for rel in new_slide.part.rels.values() if "image" in rel.reltype
    }
    assert new_image_targets == source_image_targets, (
        "duplicated slide should reference the same image parts as the source"
    )

    # Round-trip through save/reopen to confirm the output file is valid.
    output_path = tmp_path / "dup_output.pptx"
    prs.save(str(output_path))
    reopened = Presentation(str(output_path))
    assert len(reopened.slides) == 2

    # Every Picture shape on the duplicated slide should still resolve its
    # image part via the rels map. If the rel was dropped, accessing .image
    # would raise KeyError.
    from pptx.shapes.picture import Picture

    dup = reopened.slides[1]
    pic_count = 0
    for shape in dup.shapes:
        if isinstance(shape, Picture):
            pic_count += 1
            # .image access forces the rel lookup
            assert shape.image.content_type.startswith("image/"), (
                "duplicated picture should still resolve its image rel"
            )
    assert pic_count >= 1, "duplicated slide should contain at least one picture"


# -- group shape recursion (P0) --------------------------------------------


def test_walk_shapes_yields_grouped_children(tmp_path: Path) -> None:
    template_path = tmp_path / "with_group.pptx"
    _build_template_with_group(template_path)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]

    names = [shape.name for shape in _walk_shapes(slide.shapes)]
    # Should yield the group itself AND both children, regardless of order
    assert "Consultant_Card" in names
    assert "Grouped_Name" in names
    assert "Grouped_Role" in names


def test_find_shape_by_name_finds_shape_inside_group(tmp_path: Path) -> None:
    template_path = tmp_path / "with_group.pptx"
    _build_template_with_group(template_path)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]

    # Without group recursion, find_shape_by_name would return None for
    # grouped children because slide.shapes only yields top-level shapes.
    found = find_shape_by_name(slide, "Grouped_Name")
    assert found is not None
    assert found.name == "Grouped_Name"


def test_inspect_reports_group_and_children(tmp_path: Path) -> None:
    template_path = tmp_path / "with_group.pptx"
    _build_template_with_group(template_path)

    info = inspect_template(template_path)
    shape_names = [s.name for s in info.slides[0].shapes]
    assert "Consultant_Card" in shape_names
    assert "Grouped_Name" in shape_names
    assert "Grouped_Role" in shape_names

    # The formatter should indent grouped children visibly
    formatted = format_template_info(info)
    assert "Consultant_Card" in formatted
    assert "Grouped_Name" in formatted


# -- set_shape_value edge cases (P1) ---------------------------------------


@pytest.mark.parametrize(
    ("value", "expected"),
    [
        (42, "42"),
        (3.14, "3.14"),
        (0, "0"),
        ("", ""),
        (None, ""),
        ([], ""),
    ],
)
def test_set_shape_value_handles_edge_cases(tmp_path: Path, value: object, expected: str) -> None:
    template_path = tmp_path / "simple.pptx"
    _simple_template(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "Field")
    assert shape is not None

    set_shape_value(shape, value)
    assert shape.text_frame.text == expected


def test_set_shape_value_skips_none_items_in_list(tmp_path: Path) -> None:
    template_path = tmp_path / "simple.pptx"
    _simple_template(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "Field")
    assert shape is not None

    set_shape_value(shape, ["one", None, "two", "", "three"])
    paras = [p.text for p in shape.text_frame.paragraphs]
    assert paras == ["one", "two", "three"]


def test_set_shape_value_noop_on_shape_without_text_frame() -> None:
    class Stub:
        has_text_frame = False

    # Should not raise even though Stub has no text_frame attribute
    set_shape_value(Stub(), "anything")


# -- config validation (P1) ------------------------------------------------


def test_load_config_rejects_non_dict_yaml(tmp_path: Path) -> None:
    config_path = tmp_path / "bad.yaml"
    config_path.write_text("- just\n- a list\n", encoding="utf-8")

    with pytest.raises(ValueError, match="expected top-level mapping"):
        load_config(config_path)


def test_load_config_rejects_empty_yaml(tmp_path: Path) -> None:
    config_path = tmp_path / "empty.yaml"
    config_path.write_text("", encoding="utf-8")

    with pytest.raises(ValueError, match="empty"):
        load_config(config_path)


def test_load_config_rejects_invalid_yaml_syntax(tmp_path: Path) -> None:
    config_path = tmp_path / "broken.yaml"
    config_path.write_text("template: [unclosed\n", encoding="utf-8")

    with pytest.raises(ValueError, match="invalid YAML"):
        load_config(config_path)


def test_load_config_rejects_wrong_type_for_source_slide_index(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "bad.yaml"
    config_path.write_text(
        f"template: {template_path}\nsource_slide_index: two\nplaceholders:\n  name: Field\n",
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match=r"source_slide_index.*integer"):
        load_config(config_path)


def test_load_config_rejects_non_string_template(tmp_path: Path) -> None:
    config_path = tmp_path / "bad.yaml"
    config_path.write_text("template: 42\nplaceholders:\n  name: Field\n", encoding="utf-8")

    with pytest.raises(ValueError, match="'template' must be a string"):
        load_config(config_path)


def test_load_config_rejects_non_dict_placeholders(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "bad.yaml"
    config_path.write_text(
        f"template: {template_path}\nplaceholders:\n  - a\n  - b\n",
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match="'placeholders' must be a mapping"):
        load_config(config_path)


def test_load_config_rejects_non_bool_clear_source_slide(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "bad.yaml"
    config_path.write_text(
        f"template: {template_path}\nclear_source_slide: maybe\nplaceholders:\n  name: Field\n",
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match=r"clear_source_slide.*boolean"):
        load_config(config_path)


# -- generate_deck edge cases (P1) -----------------------------------------


def test_generate_deck_rejects_out_of_range_source_slide_index(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "config.yaml"
    _write_config_yaml(config_path, template_path, {"name": "Field"})

    cfg = load_config(config_path)
    cfg.source_slide_index = 99

    (tmp_path / "data").mkdir()
    (tmp_path / "data" / "r.yaml").write_text("name: X", encoding="utf-8")
    records = load_records(tmp_path / "data")

    with pytest.raises(ValueError, match="out of range"):
        generate_deck(cfg, records, tmp_path / "out.pptx")


def test_generate_deck_creates_parent_dir_automatically(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "config.yaml"
    _write_config_yaml(config_path, template_path, {"name": "Field"})

    (tmp_path / "data").mkdir()
    (tmp_path / "data" / "r.yaml").write_text("name: X", encoding="utf-8")

    # Parent dir `deep/nested/subdir/` does not exist
    output = tmp_path / "deep" / "nested" / "subdir" / "out.pptx"
    cfg = load_config(config_path)
    records = load_records(tmp_path / "data")
    result = generate_deck(cfg, records, output)

    assert output.exists()
    assert result["records_generated"] == 1


# -- CLI generate tests (P0/P1) --------------------------------------------


def _setup_minimal_project(tmp_path: Path) -> tuple[Path, Path, Path]:
    """Return (config_path, data_dir, template_path) for CLI tests."""
    template_path = tmp_path / "template.pptx"
    _simple_template(template_path)
    config_path = tmp_path / "config.yaml"
    _write_config_yaml(config_path, template_path, {"name": "Field"})
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "record.yaml").write_text("name: Test Value", encoding="utf-8")
    return config_path, data_dir, template_path


def test_cli_generate_creates_parent_directory(tmp_path: Path) -> None:
    config_path, data_dir, _ = _setup_minimal_project(tmp_path)
    output = tmp_path / "nested" / "sub" / "deck.pptx"

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
        ],
    )
    assert result.exit_code == 0, result.output
    assert output.exists()


def test_cli_generate_refuses_overwrite_without_force(tmp_path: Path) -> None:
    config_path, data_dir, _ = _setup_minimal_project(tmp_path)
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"not a real pptx but exists")

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
        ],
    )
    assert result.exit_code == 1
    assert output.read_bytes() == b"not a real pptx but exists"  # untouched


def test_cli_generate_overwrites_with_force(tmp_path: Path) -> None:
    config_path, data_dir, _ = _setup_minimal_project(tmp_path)
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"stale content")

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
            "--force",
        ],
    )
    assert result.exit_code == 0, result.output
    assert output.read_bytes() != b"stale content"


def test_cli_generate_warns_on_non_pptx_suffix(tmp_path: Path) -> None:
    config_path, data_dir, _ = _setup_minimal_project(tmp_path)
    output = tmp_path / "deck.ppt"

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
        ],
    )
    # Should succeed (warning not error) but the warning text should appear
    assert result.exit_code == 0, result.output
    combined = result.output + (result.stderr if hasattr(result, "stderr") else "")
    assert "does not end in .pptx" in combined


# -- CLI main() exception handling (P0) ------------------------------------


def test_cli_main_handles_missing_file_cleanly(tmp_path: Path) -> None:
    rc = main(
        [
            "generate",
            "-c",
            str(tmp_path / "nope.yaml"),
            "-d",
            str(tmp_path),
            "-o",
            str(tmp_path / "out.pptx"),
        ]
    )
    assert rc == 1


def test_cli_main_handles_corrupt_pptx_cleanly(tmp_path: Path) -> None:
    # Create a file that looks like pptx but isn't
    fake_pptx = tmp_path / "corrupt.pptx"
    fake_pptx.write_bytes(b"not actually a zip file")
    rc = main(["inspect", str(fake_pptx)])
    assert rc == 1


# -- workflow epilog visible in top-level help (P1) ------------------------


def test_top_level_help_mentions_workflow_order() -> None:
    result = runner.invoke(app, ["--help"])
    assert result.exit_code == 0
    # Epilog should teach the new→inspect→init→generate sequence
    assert "Typical workflow" in result.output
    assert "new" in result.output
    assert "inspect" in result.output
    assert "init" in result.output
    assert "generate" in result.output


def test_inspect_help_mentions_safe_to_share() -> None:
    result = runner.invoke(app, ["inspect", "--help"])
    assert result.exit_code == 0
    assert "safe to share" in result.output.lower()
