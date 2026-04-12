"""Tests for validate --data-dir: pre-flight record validation against config."""

from __future__ import annotations

from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches
from typer.testing import CliRunner

from recombinase.cli import app

runner = CliRunner()


def _make_template(path: Path, shape_name: str = "Consultant_Name") -> None:
    """Minimal template with one named textbox."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = shape_name
    tb.text_frame.text = "EXAMPLE"
    prs.save(str(path))


def _make_template_with_table(path: Path, text_shape: str, table_shape: str) -> None:
    """Template with one textbox and one table shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = text_shape
    tb.text_frame.text = "EXAMPLE"
    tbl = slide.shapes.add_table(2, 2, Inches(0.5), Inches(3), Inches(4), Inches(2))
    tbl.name = table_shape
    prs.save(str(path))


def _write_config(
    config_path: Path,
    template_path: Path,
    placeholders: dict[str, str],
    tables: dict[str, dict] | None = None,
) -> None:
    data: dict = {
        "template": str(template_path),
        "source_slide_index": 1,
        "clear_source_slide": True,
        "placeholders": placeholders,
    }
    if tables:
        data["tables"] = tables
    config_path.write_text(yaml.safe_dump(data), encoding="utf-8")


# ---------------------------------------------------------------------------
# Happy path: all record fields match config
# ---------------------------------------------------------------------------


def test_validate_data_dir_happy_path(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _make_template(template_path, "Consultant_Name")
    config_path = tmp_path / "config.yaml"
    _write_config(config_path, template_path, {"name": "Consultant_Name"})

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump({"id": "alpha", "name": "Alpha Jones"}),
        encoding="utf-8",
    )

    result = runner.invoke(app, ["validate", "-c", str(config_path), "-d", str(data_dir)])
    assert result.exit_code == 0, result.output
    assert "Data files valid" in result.output


# ---------------------------------------------------------------------------
# Missing placeholder field warns
# ---------------------------------------------------------------------------


def test_validate_data_dir_missing_field_warns(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _make_template(template_path, "Consultant_Name")
    config_path = tmp_path / "config.yaml"
    _write_config(config_path, template_path, {"name": "Consultant_Name"})

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    # Record deliberately omits the 'name' field
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump({"id": "alpha"}),
        encoding="utf-8",
    )

    result = runner.invoke(app, ["validate", "-c", str(config_path), "-d", str(data_dir)])
    assert result.exit_code == 0, result.output
    combined = result.output + (result.stderr if hasattr(result, "stderr") else "")
    assert "missing placeholder field 'name'" in combined


# ---------------------------------------------------------------------------
# Table field has scalar instead of list warns
# ---------------------------------------------------------------------------


def test_validate_data_dir_table_not_list_warns(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _make_template_with_table(template_path, "Consultant_Name", "Skills_Table")
    config_path = tmp_path / "config.yaml"
    _write_config(
        config_path,
        template_path,
        {"name": "Consultant_Name"},
        tables={"skills": {"shape": "Skills_Table", "columns": ["skill"]}},
    )

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    # 'skills' should be a list of dicts but is given as a plain string
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump({"id": "alpha", "name": "Alpha", "skills": "not-a-list"}),
        encoding="utf-8",
    )

    result = runner.invoke(app, ["validate", "-c", str(config_path), "-d", str(data_dir)])
    assert result.exit_code == 0, result.output
    combined = result.output + (result.stderr if hasattr(result, "stderr") else "")
    assert "expected list" in combined
    assert "skills" in combined


# ---------------------------------------------------------------------------
# Extra fields in record not in config warns
# ---------------------------------------------------------------------------


def test_validate_data_dir_extra_fields_warns(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _make_template(template_path, "Consultant_Name")
    config_path = tmp_path / "config.yaml"
    _write_config(config_path, template_path, {"name": "Consultant_Name"})

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump({"id": "alpha", "name": "Alpha", "ghost_field": "unexpected"}),
        encoding="utf-8",
    )

    result = runner.invoke(app, ["validate", "-c", str(config_path), "-d", str(data_dir)])
    assert result.exit_code == 0, result.output
    combined = result.output + (result.stderr if hasattr(result, "stderr") else "")
    assert "unused field 'ghost_field'" in combined


# ---------------------------------------------------------------------------
# --strict with data warnings exits 2
# ---------------------------------------------------------------------------


def test_validate_data_dir_strict_exits_2(tmp_path: Path) -> None:
    template_path = tmp_path / "t.pptx"
    _make_template(template_path, "Consultant_Name")
    config_path = tmp_path / "config.yaml"
    _write_config(config_path, template_path, {"name": "Consultant_Name"})

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    # Missing required 'name' field triggers a data warning
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump({"id": "alpha"}),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        ["validate", "-c", str(config_path), "-d", str(data_dir), "--strict"],
    )
    assert result.exit_code == 2
