"""v0.1.12: table cell population + picture placeholder insertion.

Two new features covering the remaining real-template gaps:
- populate_table walks a table.rows x cells grid driven by a TableConfig,
  with header-row skipping and list-as-multi-line-cell support
- set_picture calls PicturePlaceholder.insert_picture on picture
  placeholder shapes, with optional relative-path resolution
- generate_deck routes placeholder shapes to set_shape_value OR
  set_picture based on runtime type, and processes config.tables after
  the main placeholders loop
"""

from __future__ import annotations

import struct
import zlib
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt

from recombinase.config import SectionConfig, TableConfig, load_config
from recombinase.generate import (
    find_shape_by_name,
    generate_deck,
    is_picture_placeholder,
    load_records,
    populate_sections,
    populate_table,
    set_picture,
)

# -- populate_table: basic population --------------------------------------


def _build_template_with_table(path: Path) -> None:
    """Build a template with a named 3x3 table (header row + 2 data rows)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Add a 3-row x 3-col table (1 header row + 2 body rows)
    table_shape = slide.shapes.add_table(
        rows=3, cols=3, left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(2)
    )
    table_shape.name = "recent_projects"
    # Populate header row
    table = table_shape.table
    table.cell(0, 0).text = "Client & Project"
    table.cell(0, 1).text = "Role"
    table.cell(0, 2).text = "Achievements"
    # Pre-fill body rows with example content to simulate a real template
    table.cell(1, 0).text = "Example Client A — Example Project"
    table.cell(1, 1).text = "Example Role"
    table.cell(1, 2).text = "Example achievement"
    table.cell(2, 0).text = "Example Client B — Example Project"
    table.cell(2, 1).text = "Example Role"
    table.cell(2, 2).text = "Example achievement"
    prs.save(str(path))


def test_populate_table_skips_header_and_writes_data_rows(tmp_path: Path) -> None:
    template_path = tmp_path / "table.pptx"
    _build_template_with_table(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "recent_projects")
    assert shape is not None

    config = TableConfig(
        shape="recent_projects",
        columns=["client_project", "role", "achievements"],
        header_row=True,
    )

    rows = [
        {
            "client_project": "Bank A — Data governance",
            "role": "Lead consultant",
            "achievements": "Delivered framework",
        },
        {
            "client_project": "Bank B — AI tiering",
            "role": "Advisor",
            "achievements": "Designed taxonomy",
        },
    ]

    warnings = populate_table(shape, config, rows)
    assert warnings == []

    table = shape.table
    # Header row unchanged
    assert table.cell(0, 0).text == "Client & Project"
    assert table.cell(0, 1).text == "Role"
    assert table.cell(0, 2).text == "Achievements"
    # Data row 1
    assert table.cell(1, 0).text == "Bank A — Data governance"
    assert table.cell(1, 1).text == "Lead consultant"
    assert table.cell(1, 2).text == "Delivered framework"
    # Data row 2
    assert table.cell(2, 0).text == "Bank B — AI tiering"
    assert table.cell(2, 1).text == "Advisor"
    assert table.cell(2, 2).text == "Designed taxonomy"


def test_populate_table_joins_list_values_with_newlines(tmp_path: Path) -> None:
    template_path = tmp_path / "table.pptx"
    _build_template_with_table(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "recent_projects")
    assert shape is not None

    config = TableConfig(
        shape="recent_projects",
        columns=["client_project", "role", "achievements"],
        header_row=True,
    )

    rows = [
        {
            "client_project": "Bank A — Data governance",
            "role": "Lead",
            "achievements": [
                "Delivered framework",
                "Completed audit cycle",
                "Trained 20 reviewers",
            ],
        },
    ]

    warnings = populate_table(shape, config, rows)
    # One record against a 2-data-row template: expect a single
    # "clearing unused row(s)" warning for the untouched row.
    assert len(warnings) == 1
    assert "clearing 1 unused row" in warnings[0]

    table = shape.table
    cell_text = table.cell(1, 2).text
    assert "Delivered framework" in cell_text
    assert "Completed audit cycle" in cell_text
    assert "Trained 20 reviewers" in cell_text
    # Three lines separated by newlines
    assert cell_text.count("\n") == 2
    # Excess row 2 was cleared, not left with example text
    assert table.cell(2, 0).text == ""
    assert table.cell(2, 1).text == ""
    assert table.cell(2, 2).text == ""


def test_populate_table_truncates_when_rows_exceed_capacity(tmp_path: Path) -> None:
    template_path = tmp_path / "table.pptx"
    _build_template_with_table(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "recent_projects")
    assert shape is not None

    config = TableConfig(
        shape="recent_projects",
        columns=["client_project", "role", "achievements"],
        header_row=True,
    )

    # Template has 2 data rows (plus header). Provide 4 records.
    rows = [
        {"client_project": f"Client {i}", "role": "Role", "achievements": "Ach"} for i in range(4)
    ]

    warnings = populate_table(shape, config, rows)
    assert len(warnings) == 1
    assert "truncating" in warnings[0].lower()
    assert "2 data rows" in warnings[0]

    table = shape.table
    # Only the first 2 records wrote
    assert table.cell(1, 0).text == "Client 0"
    assert table.cell(2, 0).text == "Client 1"


def test_populate_table_no_header_row_writes_into_row_0(tmp_path: Path) -> None:
    template_path = tmp_path / "table.pptx"
    _build_template_with_table(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "recent_projects")
    assert shape is not None

    config = TableConfig(
        shape="recent_projects",
        columns=["client_project", "role", "achievements"],
        header_row=False,
    )

    rows = [
        {"client_project": "New Client", "role": "New Role", "achievements": "New Ach"},
    ]

    warnings = populate_table(shape, config, rows)
    # header_row=False means all 3 rows are data rows. One record leaves 2
    # excess rows, which are now cleared and surfaced as a warning.
    assert len(warnings) == 1
    assert "clearing 2 unused row" in warnings[0]

    table = shape.table
    # Row 0 overwritten (header treated as data)
    assert table.cell(0, 0).text == "New Client"
    # Excess rows 1 + 2 cleared — no example text leakage
    assert table.cell(1, 0).text == ""
    assert table.cell(1, 1).text == ""
    assert table.cell(2, 0).text == ""


def test_populate_table_warns_on_non_table_shape(tmp_path: Path) -> None:
    """If the configured shape isn't a table, emit a warning and skip."""
    from recombinase.generate import _write_paragraphs  # noqa: F401 (keep import side-effect)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "not_a_table"
    tb.text_frame.text = "I am a text box"

    path = tmp_path / "textbox.pptx"
    prs.save(str(path))
    prs = Presentation(str(path))
    shape = find_shape_by_name(prs.slides[0], "not_a_table")
    assert shape is not None

    config = TableConfig(shape="not_a_table", columns=["a"])
    warnings = populate_table(shape, config, [{"a": "x"}])
    assert len(warnings) == 1
    assert "not a table shape" in warnings[0]


# -- set_picture: picture placeholder insertion ----------------------------


def _tiny_png(path: Path) -> None:
    """Write a minimal 1x1 red PNG to disk."""
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = b"\x00\x00\x00\rIHDR" + struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = struct.pack(">I", zlib.crc32(ihdr[4:]) & 0xFFFFFFFF)
    idat_data = zlib.compress(b"\x00\xff\x00\x00")
    idat = struct.pack(">I", len(idat_data)) + b"IDAT" + idat_data
    idat_crc = struct.pack(">I", zlib.crc32(b"IDAT" + idat_data) & 0xFFFFFFFF)
    iend = b"\x00\x00\x00\x00IEND" + struct.pack(">I", zlib.crc32(b"IEND") & 0xFFFFFFFF)
    path.write_bytes(sig + ihdr + ihdr_crc + idat + idat_crc + iend)


def _build_template_with_picture_placeholder(pptx_path: Path) -> None:
    """Build a template whose layout includes a picture placeholder.

    python-pptx doesn't expose a public API for adding a picture placeholder
    to a layout, so we use the layout index from the default pptx template
    that already includes one: layout 8 ("Picture with Caption") has a
    picture placeholder at idx 1.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    # Rename the picture placeholder to our test name
    for shape in slide.shapes:
        if is_picture_placeholder(shape):
            shape.name = "photo"
            break
    prs.save(str(pptx_path))


def test_is_picture_placeholder_detects_picture_placeholders(tmp_path: Path) -> None:
    template_path = tmp_path / "pic.pptx"
    _build_template_with_picture_placeholder(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "photo")
    assert shape is not None
    assert is_picture_placeholder(shape) is True


def test_is_picture_placeholder_false_for_regular_text_box(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "not_a_picture"
    tb.text_frame.text = "text"
    assert is_picture_placeholder(tb) is False


def test_set_picture_inserts_image_into_placeholder(tmp_path: Path) -> None:
    template_path = tmp_path / "pic.pptx"
    _build_template_with_picture_placeholder(template_path)

    png_path = tmp_path / "headshot.png"
    _tiny_png(png_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "photo")
    assert shape is not None

    set_picture(shape, str(png_path))

    # After insert_picture, the shape should still have a matching name and
    # the save-reopen cycle should still find it.
    output = tmp_path / "out.pptx"
    prs.save(str(output))
    reopened = Presentation(str(output))
    new_shape = find_shape_by_name(reopened.slides[0], "photo")
    assert new_shape is not None


def test_set_picture_raises_on_missing_file(tmp_path: Path) -> None:
    import pytest

    template_path = tmp_path / "pic.pptx"
    _build_template_with_picture_placeholder(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "photo")
    assert shape is not None

    with pytest.raises(FileNotFoundError, match="picture file not found"):
        set_picture(shape, str(tmp_path / "nonexistent.jpg"))


def test_set_picture_resolves_relative_path_against_base_dir(tmp_path: Path) -> None:
    template_path = tmp_path / "pic.pptx"
    _build_template_with_picture_placeholder(template_path)

    # Create a png in a sub directory
    base_dir = tmp_path / "data"
    base_dir.mkdir()
    png_path = base_dir / "headshot.png"
    _tiny_png(png_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "photo")
    assert shape is not None

    # Pass a relative path plus the base_dir
    set_picture(shape, "headshot.png", base_dir=base_dir)

    # If it worked, there's no exception and the shape remains
    output = tmp_path / "out.pptx"
    prs.save(str(output))


def test_set_picture_noop_on_empty_value(tmp_path: Path) -> None:
    template_path = tmp_path / "pic.pptx"
    _build_template_with_picture_placeholder(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "photo")
    assert shape is not None

    # Should not raise
    set_picture(shape, "")
    set_picture(shape, None)


# -- end-to-end: generate_deck with table config + placeholders ------------


def test_generate_deck_populates_table_from_record(tmp_path: Path, write_config) -> None:
    """Full pipeline: template with a table shape, config has `tables:`,
    records have list-of-dicts, generated output has the rows populated."""
    template_path = tmp_path / "with_table.pptx"
    _build_template_with_table(template_path)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "overflow_ratio": 0,  # disable overflow for this test
                "placeholders": {},
                "tables": {
                    "recent_projects": {
                        "shape": "recent_projects",
                        "header_row": True,
                        "columns": ["client_project", "role", "achievements"],
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r1.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "alpha",
                "recent_projects": [
                    {
                        "client_project": "Bank A — Framework",
                        "role": "Lead",
                        "achievements": ["One", "Two"],
                    },
                    {
                        "client_project": "Bank B — AI",
                        "role": "Advisor",
                        "achievements": ["Three"],
                    },
                ],
            }
        ),
        encoding="utf-8",
    )

    config = load_config(config_path)
    records = load_records(data_dir)
    output = tmp_path / "out.pptx"
    result = generate_deck(config, records, output)

    assert result["records_generated"] == 1
    assert result["warnings"] == []

    # Re-open and verify the table was populated on the first generated slide
    prs = Presentation(str(output))
    shape = find_shape_by_name(prs.slides[0], "recent_projects")
    assert shape is not None
    assert shape.has_table

    table = shape.table
    # Header row preserved
    assert table.cell(0, 0).text == "Client & Project"
    # Data rows populated
    assert "Bank A — Framework" in table.cell(1, 0).text
    assert table.cell(1, 1).text == "Lead"
    assert "One" in table.cell(1, 2).text
    assert "Two" in table.cell(1, 2).text
    assert "Bank B — AI" in table.cell(2, 0).text


def test_generate_deck_warns_when_table_record_is_not_list(tmp_path: Path, write_config) -> None:
    template_path = tmp_path / "with_table.pptx"
    _build_template_with_table(template_path)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "overflow_ratio": 0,
                "placeholders": {},
                "tables": {
                    "recent_projects": {
                        "shape": "recent_projects",
                        "columns": ["client_project"],
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r1.yaml").write_text(
        yaml.safe_dump({"id": "x", "recent_projects": "not a list"}),
        encoding="utf-8",
    )

    config = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(config, records, tmp_path / "out.pptx")

    assert result["records_generated"] == 1
    assert any("expects a list" in w for w in result["warnings"])


def test_load_config_parses_tables_section(tmp_path: Path, simple_template: Path) -> None:
    config_path = tmp_path / "c.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(simple_template),
                "placeholders": {"name": "Field"},
                "tables": {
                    "projects": {
                        "shape": "ProjTable",
                        "columns": ["a", "b"],
                        "header_row": False,
                        "list_joiner": " | ",
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    cfg = load_config(config_path)
    assert "projects" in cfg.tables
    tc = cfg.tables["projects"]
    assert tc.shape == "ProjTable"
    assert tc.columns == ["a", "b"]
    assert tc.header_row is False
    assert tc.list_joiner == " | "


def test_load_config_rejects_malformed_tables_section(
    tmp_path: Path, simple_template: Path
) -> None:
    import pytest

    config_path = tmp_path / "c.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(simple_template),
                "placeholders": {"name": "Field"},
                "tables": "not a dict",
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match="'tables' must be a mapping"):
        load_config(config_path)


# -- populate_table: run-level rPr preservation (v0.1.13 regression) -------


def test_populate_table_preserves_cell_run_font_size(tmp_path: Path) -> None:
    """Single-line cell values must preserve the template's run-level rPr.

    client CV regression: the role column rendered at default 18pt instead
    of the template's 7pt because `populate_table` took a fast path via
    `cell.text_frame.text = text` for single-line values — the same bug
    `set_shape_value` was fixed for in v0.1.10. python-pptx's `.text`
    setter wipes pPr AND rPr, so run-level font size is silently lost.

    This test builds a template with body cells styled at 7pt, populates
    them with single-line values, and asserts the output cell still
    reports 7pt. Without the fix, the populated cell reports no sz or
    the default, and the visible deck shows 18pt for that column.
    """
    template_path = tmp_path / "table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        rows=3, cols=3, left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(2)
    )
    table_shape.name = "recent_projects"
    table = table_shape.table
    table.cell(0, 0).text = "Client & Project"
    table.cell(0, 1).text = "Role"
    table.cell(0, 2).text = "Achievements"
    # Pre-fill body rows with example content at an explicit 7pt font so
    # the template carries real run-level rPr that populate_table must
    # preserve on overwrite.
    for row_idx in (1, 2):
        for col_idx in range(3):
            cell = table.cell(row_idx, col_idx)
            cell.text = "Example content"
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(7)
    prs.save(str(template_path))

    # Sanity check: template's role cell reports 7pt before populate
    prs_check = Presentation(str(template_path))
    role_cell_before = prs_check.slides[0].shapes[0].table.cell(1, 1)
    assert role_cell_before.text_frame.paragraphs[0].runs[0].font.size == Pt(7)

    prs2 = Presentation(str(template_path))
    shape = find_shape_by_name(prs2.slides[0], "recent_projects")
    assert shape is not None
    config = TableConfig(
        shape="recent_projects",
        columns=["client_project", "role", "achievements"],
        header_row=True,
    )
    rows = [
        {
            "client_project": "Bank A — Governance",
            "role": "Project lead",
            "achievements": "Shipped framework",
        },
        {
            "client_project": "Bank B — Tiering",
            "role": "Advisor",
            "achievements": "Built taxonomy",
        },
    ]
    warnings = populate_table(shape, config, rows)
    assert warnings == []

    # Assert the ROLE column preserved its 7pt font size after populate.
    role_cell_after = shape.table.cell(1, 1)
    assert role_cell_after.text == "Project lead"
    run = role_cell_after.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(7), (
        f"role column lost font size on populate: got {run.font.size}, expected {Pt(7)}"
    )
    # And the second data row too
    role_cell_row2 = shape.table.cell(2, 1)
    assert role_cell_row2.text == "Advisor"
    run2 = role_cell_row2.text_frame.paragraphs[0].runs[0]
    assert run2.font.size == Pt(7)


# -- populate_sections: sectioned-list shapes (v0.1.14) --------------------


def _build_template_with_sectioned_shape(path: Path) -> None:
    """Build a template with a text box whose first two paragraphs carry
    distinct header vs bullet styles (bold 11pt header, normal 9pt bullet).

    Mirrors the real-world "Key Competencies" CV cell: four headers
    (FS Industry, Functional, Technical, Methodical), each followed by
    a variable number of bulleted items. We seed paragraph 0 with the
    header profile and paragraph 1 with the bullet profile so
    `populate_sections` can capture both.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txbox = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(5)
    )
    txbox.name = "Key_Competencies"
    tf = txbox.text_frame
    # Paragraph 0: header profile (bold, 11pt)
    tf.text = "FS Industry"
    header_run = tf.paragraphs[0].runs[0]
    header_run.font.size = Pt(11)
    header_run.font.bold = True
    # Paragraph 1: bullet profile (9pt, not bold)
    bullet_p = tf.add_paragraph()
    bullet_p.text = "Wealth Management"
    bullet_run = bullet_p.runs[0]
    bullet_run.font.size = Pt(9)
    bullet_run.font.bold = False
    prs.save(str(path))


def test_populate_sections_renders_header_and_bullet_profiles(tmp_path: Path) -> None:
    """Four sections, each with 2-3 bullets, render with correct profiles.

    Asserts that every header paragraph carries the captured 11pt bold
    profile and every bullet paragraph carries the captured 9pt non-bold
    profile, for the whole sectioned-list emission.
    """
    template_path = tmp_path / "sections.pptx"
    _build_template_with_sectioned_shape(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "Key_Competencies")
    assert shape is not None

    section_config = SectionConfig(shape="Key_Competencies")
    sections_data = [
        {"header": "FS Industry", "items": ["Wealth Management", "Retail Banking"]},
        {"header": "Functional", "items": ["Risk Management", "Compliance", "Audit"]},
        {"header": "Technical", "items": ["Python", "PySpark"]},
        {"header": "Methodical", "items": ["Agile", "Waterfall"]},
    ]

    warnings = populate_sections(shape, section_config, sections_data)
    assert warnings == []

    paragraphs = list(shape.text_frame.paragraphs)
    # Expected emission: 4 headers + 2+3+2+2 bullets = 13 paragraphs
    assert len(paragraphs) == 13

    expected = [
        ("FS Industry", True),
        ("Wealth Management", False),
        ("Retail Banking", False),
        ("Functional", True),
        ("Risk Management", False),
        ("Compliance", False),
        ("Audit", False),
        ("Technical", True),
        ("Python", False),
        ("PySpark", False),
        ("Methodical", True),
        ("Agile", False),
        ("Waterfall", False),
    ]
    for paragraph, (text, is_header) in zip(paragraphs, expected, strict=True):
        assert paragraph.text == text
        run = paragraph.runs[0]
        if is_header:
            assert run.font.size == Pt(11), f"header {text!r} lost font size: got {run.font.size}"
            assert run.font.bold is True, f"header {text!r} lost bold"
        else:
            assert run.font.size == Pt(9), f"bullet {text!r} lost font size: got {run.font.size}"
            assert run.font.bold is False, f"bullet {text!r} gained bold"


def test_populate_sections_warns_on_index_out_of_range(tmp_path: Path) -> None:
    """header_index / bullet_index past the template's paragraph count warn."""
    template_path = tmp_path / "sections_short.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txbox = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(1)
    )
    txbox.name = "Short"
    txbox.text_frame.text = "Only one paragraph"
    prs.save(str(template_path))

    prs2 = Presentation(str(template_path))
    shape = find_shape_by_name(prs2.slides[0], "Short")
    assert shape is not None

    # bullet_index=1 but template has only paragraph 0
    section_config = SectionConfig(shape="Short", header_index=0, bullet_index=1)
    warnings = populate_sections(shape, section_config, [{"header": "A", "items": ["x"]}])
    assert any("bullet_index=1" in w for w in warnings)


def test_populate_sections_skips_malformed_sections(tmp_path: Path) -> None:
    """Non-dict sections, missing headers, and non-list items are warned and
    handled without killing the whole populate."""
    template_path = tmp_path / "sections.pptx"
    _build_template_with_sectioned_shape(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "Key_Competencies")
    assert shape is not None

    sections_data = [
        {"header": "Good", "items": ["one"]},
        "not a dict",  # -> warn + skip
        {"items": ["orphan"]},  # missing header -> warn + skip
        {"header": "Also Good", "items": "not a list"},  # -> warn, render header only
    ]
    warnings = populate_sections(shape, SectionConfig(shape="Key_Competencies"), sections_data)
    assert any("not a dict" in w or "section 1" in w for w in warnings)
    assert any("header" in w for w in warnings)
    assert any("'items' must be a list" in w for w in warnings)

    paragraphs = list(shape.text_frame.paragraphs)
    # Emission: "Good", "one", "Also Good" (3 paragraphs)
    assert [p.text for p in paragraphs] == ["Good", "one", "Also Good"]


def test_load_config_parses_sections_block(tmp_path: Path) -> None:
    """`sections:` block round-trips through load_config with defaults."""
    template_path = tmp_path / "t.pptx"
    Presentation().save(str(template_path))

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "sections": {
                    "key_competencies": {
                        "shape": "Key_Competencies",
                    },
                    "key_competencies_custom": {
                        "shape": "KC2",
                        "header_index": 2,
                        "bullet_index": 5,
                    },
                },
            }
        ),
        encoding="utf-8",
    )
    config = load_config(config_path)
    assert "key_competencies" in config.sections
    assert config.sections["key_competencies"].shape == "Key_Competencies"
    assert config.sections["key_competencies"].header_index == 0
    assert config.sections["key_competencies"].bullet_index == 1
    assert config.sections["key_competencies_custom"].header_index == 2
    assert config.sections["key_competencies_custom"].bullet_index == 5


def test_load_config_rejects_equal_header_and_bullet_index(tmp_path: Path) -> None:
    """header_index == bullet_index fails validation with a clear message."""
    import pytest

    template_path = tmp_path / "t.pptx"
    Presentation().save(str(template_path))

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "sections": {
                    "kc": {
                        "shape": "KC",
                        "header_index": 3,
                        "bullet_index": 3,
                    }
                },
            }
        ),
        encoding="utf-8",
    )
    with pytest.raises(ValueError, match="header_index and bullet_index"):
        load_config(config_path)


def test_generate_deck_populates_sections_end_to_end(tmp_path: Path) -> None:
    """End-to-end: template + record with sections renders into the output."""
    template_path = tmp_path / "deck.pptx"
    _build_template_with_sectioned_shape(template_path)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "source_slide_index": 1,
                "sections": {
                    "key_competencies": {"shape": "Key_Competencies"},
                },
                "clear_source_slide": False,
            }
        ),
        encoding="utf-8",
    )
    config = load_config(config_path)

    record = {
        "id": "alice",
        "key_competencies": [
            {"header": "FS Industry", "items": ["Retail"]},
            {"header": "Functional", "items": ["Risk", "Audit"]},
        ],
    }
    output_path = tmp_path / "out.pptx"
    result = generate_deck(config, [record], output_path)
    assert result["records_generated"] == 1
    assert not any("failed" in w or "not found" in w for w in result["warnings"])

    prs_out = Presentation(str(output_path))
    # Source slide at index 0, generated at index 1 (clear_source_slide=False)
    new_slide = prs_out.slides[1]
    shape = find_shape_by_name(new_slide, "Key_Competencies")
    assert shape is not None
    texts = [p.text for p in shape.text_frame.paragraphs]
    assert texts == ["FS Industry", "Retail", "Functional", "Risk", "Audit"]


# -- populate_table / set_shape_value: multi-run-br preservation (v0.1.15) -


def _build_cell_with_multirun_br(cell: object) -> None:
    """Seed a table cell with a single paragraph containing:

        <a:r b="1">Example Primary, Firm</a:r>
        <a:br/>
        <a:r i="1">(N years)</a:r>

    This is the CV template idiom `_write_multirun_br` preserves:
    bold run + soft break + italic run in a single paragraph.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    tf = cell.text_frame
    tf.text = ""  # normalise to single empty paragraph
    p_xml = tf.paragraphs[0]._p
    # Drop the empty run the normalisation left behind
    for existing_run in p_xml.findall(qn("a:r")):
        p_xml.remove(existing_run)

    # Run 0: bold primary
    r0 = etree.SubElement(p_xml, qn("a:r"))
    rPr0 = etree.SubElement(r0, qn("a:rPr"))
    rPr0.set("b", "1")
    t0 = etree.SubElement(r0, qn("a:t"))
    t0.text = "Example Primary, Firm"
    # Soft break
    etree.SubElement(p_xml, qn("a:br"))
    # Run 1: italic secondary
    r1 = etree.SubElement(p_xml, qn("a:r"))
    rPr1 = etree.SubElement(r1, qn("a:rPr"))
    rPr1.set("i", "1")
    t1 = etree.SubElement(r1, qn("a:t"))
    t1.text = "(5 years)"


def test_populate_table_preserves_multirun_br_on_scalar_with_newline(
    tmp_path: Path,
) -> None:
    """Scalar values containing `\\n` route to the multirun-br writer when
    the source cell has the dual-run-br idiom, preserving per-run rPr
    and the soft break. Regression for the client career/recent_projects
    column where "Role\\n(duration)" was flattening into a single run."""
    template_path = tmp_path / "table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        rows=3, cols=2, left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(2)
    )
    table_shape.name = "career"
    table = table_shape.table
    table.cell(0, 0).text = "Role"
    table.cell(0, 1).text = "Organisation"
    # Body rows: column 0 has the multi-run-br idiom
    for row_idx in (1, 2):
        _build_cell_with_multirun_br(table.cell(row_idx, 0))
        table.cell(row_idx, 1).text = "Example Firm"
    prs.save(str(template_path))

    prs2 = Presentation(str(template_path))
    shape = find_shape_by_name(prs2.slides[0], "career")
    assert shape is not None
    config = TableConfig(
        shape="career",
        columns=["role", "organisation"],
        header_row=True,
    )
    rows = [
        {"role": "Management Principal, client\n(3 years)", "organisation": "client"},
        {"role": "AGM, China CITIC Bank International\n(3.5 years)", "organisation": "CNCBI"},
    ]
    warnings = populate_table(shape, config, rows)
    assert warnings == []

    # Row 1: two runs preserved, first bold + new primary text, second italic + new years
    cell_r1 = shape.table.cell(1, 0)
    paragraphs_r1 = cell_r1.text_frame.paragraphs
    assert len(paragraphs_r1) == 1, "dual-run-br cell must stay single-paragraph"
    runs_r1 = paragraphs_r1[0].runs
    assert len(runs_r1) == 2, f"expected 2 runs, got {len(runs_r1)}"
    assert runs_r1[0].text == "Management Principal, client"
    assert runs_r1[0].font.bold is True, "primary run lost bold rPr"
    assert runs_r1[1].text == "(3 years)"
    assert runs_r1[1].font.italic is True, "secondary run lost italic rPr"
    # Soft break still in the XML
    from pptx.oxml.ns import qn

    br_elements = paragraphs_r1[0]._p.findall(qn("a:br"))
    assert len(br_elements) == 1, "soft break <a:br/> was dropped"

    # Row 2: same structure, different text
    cell_r2 = shape.table.cell(2, 0)
    runs_r2 = cell_r2.text_frame.paragraphs[0].runs
    assert runs_r2[0].text == "AGM, China CITIC Bank International"
    assert runs_r2[0].font.bold is True
    assert runs_r2[1].text == "(3.5 years)"
    assert runs_r2[1].font.italic is True


def test_set_shape_value_preserves_multirun_br_on_placeholder() -> None:
    """Same multirun-br preservation applies to a standalone text box
    populated via set_shape_value, not just table cells. In-memory — text
    boxes don't survive save/reload through direct lxml manipulation (the
    pptx serialisation pipeline rebuilds text body from cached state),
    but the writer itself operates on live lxml elements so we can
    unit-test it without a round-trip."""
    from recombinase.generate import set_shape_value

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txbox = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(2)
    )
    txbox.name = "role_title"
    _build_cell_with_multirun_br(txbox)

    set_shape_value(txbox, "Principal Consultant, AI Solution Lead\n(10 years)")

    paragraphs = txbox.text_frame.paragraphs
    assert len(paragraphs) == 1
    runs = paragraphs[0].runs
    assert len(runs) == 2
    assert runs[0].text == "Principal Consultant, AI Solution Lead"
    assert runs[0].font.bold is True
    assert runs[1].text == "(10 years)"
    assert runs[1].font.italic is True


def test_populate_table_falls_back_on_multi_paragraph_cell(tmp_path: Path) -> None:
    """When the source cell is multi-paragraph (not multi-run-br), values
    with newlines still route through _write_paragraphs so each line
    becomes a bulleted paragraph — the achievements column idiom."""
    template_path = tmp_path / "fallback.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        rows=2, cols=1, left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(2)
    )
    table_shape.name = "achievements_tbl"
    table_shape.table.cell(0, 0).text = "Header"
    # Body: plain single-paragraph single-run cell (no dual-run-br idiom)
    table_shape.table.cell(1, 0).text = "Example achievement"
    prs.save(str(template_path))

    prs2 = Presentation(str(template_path))
    shape = find_shape_by_name(prs2.slides[0], "achievements_tbl")
    assert shape is not None
    config = TableConfig(shape="achievements_tbl", columns=["ach"], header_row=True)
    rows = [{"ach": ["First win", "Second win", "Third win"]}]
    warnings = populate_table(shape, config, rows)
    assert warnings == []

    cell = shape.table.cell(1, 0)
    texts = [p.text for p in cell.text_frame.paragraphs]
    assert texts == ["First win", "Second win", "Third win"], (
        "non-multirun-br cell must still use paragraph-based write"
    )


def test_write_multirun_br_handles_unequal_parts_gracefully() -> None:
    """Fewer parts than runs → trailing runs cleared but rPr preserved.
    More parts than runs → excess parts merged into last run's text.
    In-memory; no save/reload (text boxes don't survive round-trip with
    manually-built runs, see note in the set_shape_value test above)."""
    from pptx.oxml.ns import qn

    from recombinase.generate import _write_multirun_br

    # Case 1: one part, two runs → run 0 gets text, run 1 cleared
    prs1 = Presentation()
    slide1 = prs1.slides.add_slide(prs1.slide_layouts[5])
    txbox1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(2))
    _build_cell_with_multirun_br(txbox1)
    tf1 = txbox1.text_frame
    _write_multirun_br(tf1, ["Only one line"])
    # Read via lxml to bypass any python-pptx run-filtering on empty text
    runs1_xml = tf1.paragraphs[0]._p.findall(qn("a:r"))
    assert len(runs1_xml) == 2, f"expected 2 runs preserved, got {len(runs1_xml)}"
    t0 = runs1_xml[0].find(qn("a:t"))
    t1 = runs1_xml[1].find(qn("a:t"))
    assert t0 is not None
    assert t0.text == "Only one line"
    assert t1 is not None
    assert t1.text in ("", None)
    # rPr preserved even when the run text is empty
    assert runs1_xml[0].find(qn("a:rPr")).get("b") == "1"
    assert runs1_xml[1].find(qn("a:rPr")).get("i") == "1"

    # Case 2: three parts, two runs → last two merge into run 1
    prs2 = Presentation()
    slide2 = prs2.slides.add_slide(prs2.slide_layouts[5])
    txbox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(2))
    _build_cell_with_multirun_br(txbox2)
    tf2 = txbox2.text_frame
    _write_multirun_br(tf2, ["Primary", "Secondary", "Tertiary"])
    runs2_xml = tf2.paragraphs[0]._p.findall(qn("a:r"))
    assert len(runs2_xml) == 2
    assert runs2_xml[0].find(qn("a:t")).text == "Primary"
    assert runs2_xml[1].find(qn("a:t")).text == "Secondary Tertiary"
