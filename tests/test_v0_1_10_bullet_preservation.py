"""v0.1.10: verify bullet/paragraph formatting is preserved across list expansion.

The bug this closes: v0.1.9's `_write_paragraphs` called `text_frame.clear()`
which wiped all paragraph-level `<a:pPr>` elements. The first new paragraph
inherited bullet styling from the master/layout via level-0 defaults, but
subsequent paragraphs added via `add_paragraph()` did not. Result: a 3-item
list on a bulleted template placeholder ended up with one bullet followed by
two un-bulleted lines.

v0.1.10 fix: capture the first existing paragraph's pPr BEFORE clearing,
re-inject it into every output paragraph. Also preserves the first run's
rPr for font/size/color consistency.
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from recombinase.generate import find_shape_by_name, set_shape_value


def _make_template_with_styled_paragraph(path: Path) -> None:
    """Build a template where the text box has a paragraph with a non-default pPr.

    We can't easily hand-build a real bullet pPr without a master/layout change,
    but we can inject an indent marker (`marL`) into the pPr that's trivially
    detectable post-write. If pPr preservation works, the marker survives
    across all generated paragraphs.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(3))
    tb.name = "BulletField"
    tb.text_frame.text = "EXAMPLE"

    # Add a custom pPr with a recognizable marL (left margin) marker
    first_p = tb.text_frame.paragraphs[0]._p
    pPr = etree.SubElement(first_p, qn("a:pPr"))
    pPr.set("marL", "457200")  # 0.5 inch in EMU — our preservation marker
    pPr.set("indent", "-457200")
    # pPr must come BEFORE any runs in OOXML
    first_p.insert(0, pPr)
    first_p.remove(pPr)  # pull out the auto-appended copy
    first_p.insert(0, pPr)

    # Add a run-level marker too (font size)
    first_r = first_p.find(qn("a:r"))
    if first_r is not None:
        rPr = first_r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(first_r, qn("a:rPr"))
            first_r.insert(0, rPr)
            first_r.remove(rPr)
            first_r.insert(0, rPr)
        rPr.set("sz", "2400")  # 24pt — recognizable marker

    prs.save(str(path))


def test_write_paragraphs_preserves_pPr_across_all_items(tmp_path: Path) -> None:
    """Every output paragraph should carry the marL marker from the template."""
    template_path = tmp_path / "bullet.pptx"
    _make_template_with_styled_paragraph(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "BulletField")
    assert shape is not None

    set_shape_value(
        shape,
        [
            "Company A — Role (2010-2015)",
            "Company B — Role (2015-2020)",
            "Company C — Role (2020-present)",
        ],
    )

    # Verify all three paragraphs have the preserved marL
    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 3

    for i, p in enumerate(paragraphs):
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None, f"paragraph {i} lost its pPr after write"
        assert pPr.get("marL") == "457200", (
            f"paragraph {i} pPr marL was not preserved: got {pPr.get('marL')!r}"
        )


def test_write_paragraphs_preserves_rPr_across_all_items(tmp_path: Path) -> None:
    """Every output paragraph's first run should carry the preserved rPr."""
    template_path = tmp_path / "bullet.pptx"
    _make_template_with_styled_paragraph(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "BulletField")
    assert shape is not None

    set_shape_value(shape, ["First item", "Second item", "Third item"])

    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 3

    for i, p in enumerate(paragraphs):
        first_r = p._p.find(qn("a:r"))
        assert first_r is not None, f"paragraph {i} has no run"
        rPr = first_r.find(qn("a:rPr"))
        assert rPr is not None, f"paragraph {i} first run lost its rPr"
        assert rPr.get("sz") == "2400", (
            f"paragraph {i} rPr sz was not preserved: got {rPr.get('sz')!r}"
        )


def test_write_paragraphs_round_trips_through_save_reopen(tmp_path: Path) -> None:
    """After save/reopen the preserved styling must still be on every paragraph."""
    template_path = tmp_path / "bullet.pptx"
    _make_template_with_styled_paragraph(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "BulletField")
    assert shape is not None
    set_shape_value(shape, ["one", "two", "three", "four"])

    output_path = tmp_path / "out.pptx"
    prs.save(str(output_path))

    # Reopen and re-check
    reopened = Presentation(str(output_path))
    shape = find_shape_by_name(reopened.slides[0], "BulletField")
    assert shape is not None
    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 4

    for i, p in enumerate(paragraphs):
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None, f"paragraph {i} lost pPr across round-trip"
        assert pPr.get("marL") == "457200", f"paragraph {i} marL not preserved across round-trip"


def test_write_paragraphs_single_item_still_works(tmp_path: Path) -> None:
    """Regression: single-item list should still produce one paragraph with the style."""
    template_path = tmp_path / "bullet.pptx"
    _make_template_with_styled_paragraph(template_path)

    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "BulletField")
    assert shape is not None

    set_shape_value(shape, ["only item"])

    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 1
    pPr = paragraphs[0]._p.find(qn("a:pPr"))
    assert pPr is not None
    assert pPr.get("marL") == "457200"


# Avoid unused import warnings
_ = deepcopy
