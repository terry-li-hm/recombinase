"""Expand test coverage for exception handlers, edge cases, and fallbacks.

Tests in this file exercise code paths that had 0 coverage before v0.2.0:
- generate_deck exception-handler branches for set_shape_value, populate_table,
  and populate_sections
- populate_table non-dict rows, column-exceeds-cell-count, trailing-cell clearing
- populate_sections shape-with-no-text-frame, empty data, non-dict entries
- find_shape_by_name duplicate-name resolution
- record_id fallback chain (id → name → record_N)
- write_scaffold_config slug collision deduplication
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from unittest.mock import patch

import yaml
from pptx import Presentation
from pptx.util import Inches

from recombinase.config import SectionConfig, TableConfig, TemplateConfig, write_scaffold_config
from recombinase.generate import (
    find_shape_by_name,
    generate_deck,
    populate_sections,
    populate_table,
)

# ---------------------------------------------------------------------------
# Internal helpers for building synthetic templates
# ---------------------------------------------------------------------------


def _build_simple_textbox_template_bytes() -> bytes:
    """One named textbox on a blank layout — minimal generate_deck target."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_textbox_template_to_path(path: Path) -> None:
    path.write_bytes(_build_simple_textbox_template_bytes())


def _build_table_textbox_template(path: Path) -> None:
    """Template with a textbox AND a 2-row x 2-col table for exception tests."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"

    tbl_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2), Inches(5), Inches(2))
    tbl_shape.name = "My_Table"
    tbl_shape.table.cell(0, 0).text = "Col A"
    tbl_shape.table.cell(0, 1).text = "Col B"
    tbl_shape.table.cell(1, 0).text = "ex a"
    tbl_shape.table.cell(1, 1).text = "ex b"

    prs.save(str(path))


def _build_sections_template(path: Path) -> None:
    """Template with a textbox AND a text frame with 2 paragraphs for sections."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"

    sec_tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5), Inches(3))
    sec_tb.name = "Sections_Shape"
    sec_tf = sec_tb.text_frame
    sec_tf.text = "Header Example"
    sec_tf.add_paragraph().text = "Bullet Example"

    prs.save(str(path))


# ---------------------------------------------------------------------------
# 1. generate_deck: set_shape_value raises → warning contains "failed to set"
# ---------------------------------------------------------------------------


def test_generate_deck_set_shape_value_exception_warns(tmp_path: Path) -> None:
    """When set_shape_value raises, generate_deck emits a warning with 'failed to set'."""
    template_path = tmp_path / "template.pptx"
    _build_textbox_template_to_path(template_path)
    output_path = tmp_path / "out.pptx"

    config = TemplateConfig(
        template=template_path,
        source_slide_index=1,
        placeholders={"field": "Field"},
        clear_source_slide=False,
        overflow_ratio=0,
    )
    records = [{"field": "value"}]

    with patch("recombinase.generate.set_shape_value", side_effect=RuntimeError("boom")):
        result = generate_deck(config, records, output_path)

    assert any("failed to set" in w and "boom" in w for w in result["warnings"]), (
        f"Expected 'failed to set ... boom' warning; got: {result['warnings']}"
    )


# ---------------------------------------------------------------------------
# 2. generate_deck: populate_table raises → warning contains "failed to populate table"
# ---------------------------------------------------------------------------


def test_generate_deck_populate_table_exception_warns(tmp_path: Path) -> None:
    """When populate_table raises, generate_deck emits a warning about it."""
    template_path = tmp_path / "template.pptx"
    _build_table_textbox_template(template_path)
    output_path = tmp_path / "out.pptx"

    config = TemplateConfig(
        template=template_path,
        source_slide_index=1,
        placeholders={},
        tables={
            "my_table": TableConfig(
                shape="My_Table",
                columns=["Col A", "Col B"],
                header_row=True,
            )
        },
        clear_source_slide=False,
        overflow_ratio=0,
    )
    records = [{"my_table": [{"Col A": "a", "Col B": "b"}]}]

    with patch("recombinase.generate.populate_table", side_effect=RuntimeError("table boom")):
        result = generate_deck(config, records, output_path)

    assert any("failed to populate table" in w for w in result["warnings"]), (
        f"Expected 'failed to populate table' warning; got: {result['warnings']}"
    )


# ---------------------------------------------------------------------------
# 3. generate_deck: populate_sections raises → warning
# ---------------------------------------------------------------------------


def test_generate_deck_populate_sections_exception_warns(tmp_path: Path) -> None:
    """When populate_sections raises, generate_deck emits a warning about it."""
    template_path = tmp_path / "template.pptx"
    _build_sections_template(template_path)
    output_path = tmp_path / "out.pptx"

    config = TemplateConfig(
        template=template_path,
        source_slide_index=1,
        placeholders={},
        sections={
            "skills": SectionConfig(
                shape="Sections_Shape",
                header_index=0,
                bullet_index=1,
            )
        },
        clear_source_slide=False,
        overflow_ratio=0,
    )
    records = [{"skills": [{"header": "Tech", "items": ["Python"]}]}]

    with patch("recombinase.generate.populate_sections", side_effect=RuntimeError("sections boom")):
        result = generate_deck(config, records, output_path)

    assert any("failed to populate sections" in w for w in result["warnings"]), (
        f"Expected 'failed to populate sections' warning; got: {result['warnings']}"
    )


# ---------------------------------------------------------------------------
# 4. populate_table: non-dict row warns about "expected a dict"
# ---------------------------------------------------------------------------


def test_populate_table_non_dict_row_warns(tmp_path: Path) -> None:
    """A string row in populate_table triggers a warning containing 'expected a dict'."""
    pptx_path = tmp_path / "template.pptx"
    _build_table_textbox_template(pptx_path)

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    table_shape = next(s for s in slide.shapes if s.name == "My_Table")

    table_config = TableConfig(shape="My_Table", columns=["Col A", "Col B"], header_row=True)
    rows = ["not a dict"]
    warnings = populate_table(table_shape, table_config, rows)

    assert any("expected a dict" in w for w in warnings), (
        f"Expected 'expected a dict' warning; got: {warnings}"
    )


# ---------------------------------------------------------------------------
# 5. populate_table: column index exceeds cell count → warning
# ---------------------------------------------------------------------------


def test_populate_table_column_exceeds_cell_count(tmp_path: Path) -> None:
    """Config with more columns than the table has cells emits an 'exceeds' warning."""
    # Build a template with a 1-column table (plus header)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl_shape = slide.shapes.add_table(2, 1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
    tbl_shape.name = "Narrow_Table"
    tbl_shape.table.cell(0, 0).text = "Col A"
    tbl_shape.table.cell(1, 0).text = "ex"
    pptx_path = tmp_path / "narrow.pptx"
    prs.save(str(pptx_path))

    prs2 = Presentation(str(pptx_path))
    slide2 = prs2.slides[0]
    narrow_shape = next(s for s in slide2.shapes if s.name == "Narrow_Table")

    table_config = TableConfig(
        shape="Narrow_Table",
        columns=["Col A", "Col B"],  # 2 columns but table only has 1
        header_row=True,
    )
    rows = [{"Col A": "value a", "Col B": "value b"}]
    warnings = populate_table(narrow_shape, table_config, rows)

    assert any("exceeds" in w or "exceed" in w for w in warnings), (
        f"Expected column-exceeds warning; got: {warnings}"
    )


# ---------------------------------------------------------------------------
# 6. populate_table: extra trailing cells cleared when config has fewer columns
# ---------------------------------------------------------------------------


def test_populate_table_extra_trailing_cells_cleared(tmp_path: Path) -> None:
    """Cells beyond config columns are cleared (no example text leaks)."""
    # Build template with 3 columns, config uses 2
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl_shape = slide.shapes.add_table(2, 3, Inches(0.5), Inches(0.5), Inches(6), Inches(2))
    tbl_shape.name = "Wide_Table"
    tbl_shape.table.cell(0, 0).text = "Col A"
    tbl_shape.table.cell(0, 1).text = "Col B"
    tbl_shape.table.cell(0, 2).text = "Col C"
    tbl_shape.table.cell(1, 0).text = "ex a"
    tbl_shape.table.cell(1, 1).text = "ex b"
    tbl_shape.table.cell(1, 2).text = "ex c"  # this should be cleared
    pptx_path = tmp_path / "wide.pptx"
    prs.save(str(pptx_path))

    prs2 = Presentation(str(pptx_path))
    slide2 = prs2.slides[0]
    wide_shape = next(s for s in slide2.shapes if s.name == "Wide_Table")

    table_config = TableConfig(
        shape="Wide_Table",
        columns=["Col A", "Col B"],  # only 2 of 3 columns configured
        header_row=True,
    )
    rows = [{"Col A": "value a", "Col B": "value b"}]
    populate_table(wide_shape, table_config, rows)

    # Third cell in data row should be cleared
    data_row = wide_shape.table.rows[1]
    third_cell = data_row.cells[2]
    assert third_cell.text_frame.text == "", (
        f"Expected cleared trailing cell, got: {third_cell.text_frame.text!r}"
    )


# ---------------------------------------------------------------------------
# 7. populate_sections: shape with no text frame → warning
# ---------------------------------------------------------------------------


def test_populate_sections_shape_no_text_frame(tmp_path: Path) -> None:
    """A shape without a text frame triggers a warning from populate_sections."""
    from unittest.mock import MagicMock

    # Build a mock shape that reports has_text_frame = False
    mock_shape = MagicMock()
    mock_shape.has_text_frame = False
    mock_shape.name = "NoFrame_Shape"

    section_config = SectionConfig(shape="NoFrame_Shape", header_index=0, bullet_index=1)
    section_data = [{"header": "Tech", "items": ["Python"]}]

    warnings = populate_sections(mock_shape, section_config, section_data)

    assert any("no text frame" in w for w in warnings), (
        f"Expected 'no text frame' warning; got: {warnings}"
    )


# ---------------------------------------------------------------------------
# 8. populate_sections: empty data clears text frame
# ---------------------------------------------------------------------------


def test_populate_sections_empty_data_clears(tmp_path: Path) -> None:
    """populate_sections with an empty list clears the text frame."""
    pptx_path = tmp_path / "template.pptx"
    _build_sections_template(pptx_path)

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    sec_shape = next(s for s in slide.shapes if s.name == "Sections_Shape")

    # Pre-condition: shape has content
    assert sec_shape.text_frame.text != ""

    section_config = SectionConfig(shape="Sections_Shape", header_index=0, bullet_index=1)
    populate_sections(sec_shape, section_config, [])

    assert sec_shape.text_frame.text == "", (
        f"Expected empty text frame after empty data; got: {sec_shape.text_frame.text!r}"
    )


# ---------------------------------------------------------------------------
# 9. populate_sections: non-dict entry warns
# ---------------------------------------------------------------------------


def test_populate_sections_non_dict_entry(tmp_path: Path) -> None:
    """Non-dict entries in sections_data emit a warning."""
    pptx_path = tmp_path / "template.pptx"
    _build_sections_template(pptx_path)

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    sec_shape = next(s for s in slide.shapes if s.name == "Sections_Shape")

    section_config = SectionConfig(shape="Sections_Shape", header_index=0, bullet_index=1)
    warnings = populate_sections(sec_shape, section_config, [42, "not a dict"])

    assert any(
        "expected a dict" in w.lower() or "not a dict" in w.lower() or "int" in w for w in warnings
    ), f"Expected warning about non-dict entry; got: {warnings}"


# ---------------------------------------------------------------------------
# 10. find_shape_by_name: returns first shape on duplicate names
# ---------------------------------------------------------------------------


def test_find_shape_by_name_returns_first_on_duplicate(tmp_path: Path) -> None:
    """When two shapes share the same name, find_shape_by_name returns the first."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add two textboxes with the same name but distinguishable text
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
    tb1.name = "Dup"
    tb1.text_frame.text = "FIRST"

    tb2 = slide.shapes.add_textbox(Inches(4.0), Inches(0.5), Inches(3), Inches(1))
    tb2.name = "Dup"
    tb2.text_frame.text = "SECOND"

    pptx_path = tmp_path / "dup.pptx"
    prs.save(str(pptx_path))

    prs2 = Presentation(str(pptx_path))
    slide2 = prs2.slides[0]

    found = find_shape_by_name(slide2, "Dup")
    assert found is not None
    assert found.text_frame.text == "FIRST", (
        f"Expected first shape 'FIRST'; got: {found.text_frame.text!r}"
    )


# ---------------------------------------------------------------------------
# 11. record_id fallback chain
# ---------------------------------------------------------------------------


def test_record_id_fallback_chain(tmp_path: Path) -> None:
    """Warnings reference the right identifier: id > name > record_N fallback."""
    template_path = tmp_path / "template.pptx"
    _build_textbox_template_to_path(template_path)
    output_path = tmp_path / "out.pptx"

    config = TemplateConfig(
        template=template_path,
        source_slide_index=1,
        placeholders={"field": "Field"},
        clear_source_slide=False,
        overflow_ratio=0,
    )

    # All three records are missing the "field" key → triggers warning that
    # includes the record_id in the message.
    records = [
        {"id": "rec-001"},  # should appear as 'rec-001'
        {"name": "Alice"},  # should appear as 'Alice'
        {},  # no id or name → should appear as 'record_3'
    ]

    result = generate_deck(config, records, output_path)
    warnings = result["warnings"]
    warning_text = "\n".join(warnings)

    assert "rec-001" in warning_text, f"Expected record id 'rec-001' in warnings; got: {warnings}"
    assert "Alice" in warning_text, f"Expected record name 'Alice' in warnings; got: {warnings}"
    assert "record_3" in warning_text, f"Expected fallback 'record_3' in warnings; got: {warnings}"


# ---------------------------------------------------------------------------
# 12. write_scaffold_config: slug collision deduplication
# ---------------------------------------------------------------------------


def test_slug_collision_deduplication(tmp_path: Path) -> None:
    """Shape names that slug to the same key get distinct keys (_2, _3, ...)."""
    template_path = tmp_path / "template.pptx"
    _build_textbox_template_to_path(template_path)
    output_path = tmp_path / "config.yaml"

    # "My Field" and "My-Field" both slug to "my_field"
    shape_names = ["My Field", "My-Field", "Another Shape"]
    write_scaffold_config(template_path, shape_names, output_path)

    content = output_path.read_text(encoding="utf-8")
    config = yaml.safe_load(content)
    placeholders = config["placeholders"]

    keys = list(placeholders.keys())
    assert "my_field" in keys, f"Expected 'my_field' in {keys}"
    assert "my_field_2" in keys, f"Expected 'my_field_2' in {keys}"
    assert len(set(keys)) == len(keys), f"Duplicate keys found: {keys}"
    assert len(keys) == 3, f"Expected 3 keys, got: {keys}"
