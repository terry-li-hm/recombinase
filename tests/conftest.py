"""Shared pytest fixtures for the recombinase test suite.

Keeps fixture builders in one place so test files don't redefine them.
"""

from __future__ import annotations

import struct
import zlib
from pathlib import Path
from typing import Any

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt


def _tiny_png_bytes() -> bytes:
    """A 1x1 red PNG for fixtures — self-contained, no test assets needed."""
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = b"\x00\x00\x00\rIHDR" + struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = struct.pack(">I", zlib.crc32(ihdr[4:]) & 0xFFFFFFFF)
    idat_data = zlib.compress(b"\x00\xff\x00\x00")
    idat = struct.pack(">I", len(idat_data)) + b"IDAT" + idat_data
    idat_crc = struct.pack(">I", zlib.crc32(b"IDAT" + idat_data) & 0xFFFFFFFF)
    iend = b"\x00\x00\x00\x00IEND" + struct.pack(">I", zlib.crc32(b"IEND") & 0xFFFFFFFF)
    return sig + ihdr + ihdr_crc + idat + idat_crc + iend


@pytest.fixture
def tiny_png(tmp_path: Path) -> Path:
    """A 1x1 PNG written to tmp_path, returned as its path."""
    path = tmp_path / "dot.png"
    path.write_bytes(_tiny_png_bytes())
    return path


def _build_plain_template(path: Path, baseline_text: str = "EXAMPLE") -> None:
    """Template with a single named text box on slide 1."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(2))
    tb.name = "Field"
    tb.text_frame.text = baseline_text
    prs.save(str(path))


@pytest.fixture
def simple_template(tmp_path: Path) -> Path:
    """A tiny pptx template with one named text box. Baseline = 'EXAMPLE'."""
    path = tmp_path / "simple.pptx"
    _build_plain_template(path)
    return path


@pytest.fixture
def simple_template_with_baseline(tmp_path: Path):
    """Factory fixture: build a simple template with custom baseline text."""

    def _factory(baseline: str = "EXAMPLE") -> Path:
        path = tmp_path / f"simple_{len(baseline)}.pptx"
        _build_plain_template(path, baseline)
        return path

    return _factory


def _build_rich_sample_template(path: Path) -> None:
    """Template with multiple named text boxes — mirrors the original
    test_end_to_end fixture. Used by full-pipeline tests."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def add_textbox(name: str, left: float, top: float, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(5), Inches(1))
        box.name = name
        box.text_frame.text = text
        for p in box.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)

    add_textbox("Consultant_Name", 0.5, 0.5, "EXAMPLE NAME")
    add_textbox("Role_Title", 0.5, 1.3, "EXAMPLE ROLE")
    add_textbox("Summary_Body", 0.5, 2.1, "Example summary paragraph.")
    add_textbox("Background_Bullets", 0.5, 3.0, "Example bullet")
    add_textbox("Key_Skills", 0.5, 4.2, "Example skill")

    prs.save(str(path))


@pytest.fixture
def rich_template(tmp_path: Path) -> Path:
    """A template with 5 named text boxes covering the common CV field set."""
    path = tmp_path / "rich.pptx"
    _build_rich_sample_template(path)
    return path


def _build_template_with_picture(pptx_path: Path, png_path: Path) -> None:
    """Template with a text box and a picture (tests rels duplication)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "Consultant_Name"
    tb.text_frame.text = "EXAMPLE"

    pic = slide.shapes.add_picture(str(png_path), Inches(5), Inches(0.5), Inches(1), Inches(1))
    pic.name = "Headshot"

    prs.save(str(pptx_path))


@pytest.fixture
def template_with_picture(tmp_path: Path, tiny_png: Path) -> Path:
    """Template with one text box and one Picture shape (named Headshot)."""
    path = tmp_path / "with_pic.pptx"
    _build_template_with_picture(path, tiny_png)
    return path


def _build_template_with_group(pptx_path: Path) -> None:
    """Template with two text boxes reparented into a `<p:grpSp>` group.

    python-pptx has no public API for grouping shapes, so we hand-build the
    grpSp XML and move the text boxes into it. Verified that python-pptx
    classifies the result as MSO_SHAPE_TYPE.GROUP with accessible .shapes.
    """
    from copy import deepcopy

    from lxml import etree
    from pptx.oxml.ns import qn

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb1.name = "Grouped_Name"
    tb1.text_frame.text = "EXAMPLE"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(1))
    tb2.name = "Grouped_Role"
    tb2.text_frame.text = "EXAMPLE"

    spTree = slide.shapes._spTree
    grpSp = etree.SubElement(spTree, qn("p:grpSp"))
    nvGrpSpPr = etree.SubElement(grpSp, qn("p:nvGrpSpPr"))
    cNvPr = etree.SubElement(nvGrpSpPr, qn("p:cNvPr"))
    cNvPr.set("id", "999")
    cNvPr.set("name", "Consultant_Card")
    etree.SubElement(nvGrpSpPr, qn("p:cNvGrpSpPr"))
    etree.SubElement(nvGrpSpPr, qn("p:nvPr"))
    grpSpPr = etree.SubElement(grpSp, qn("p:grpSpPr"))
    xfrm = etree.SubElement(grpSpPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", "0")
    off.set("y", "0")
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", "9144000")
    ext.set("cy", "6858000")
    chOff = etree.SubElement(xfrm, qn("a:chOff"))
    chOff.set("x", "0")
    chOff.set("y", "0")
    chExt = etree.SubElement(xfrm, qn("a:chExt"))
    chExt.set("cx", "9144000")
    chExt.set("cy", "6858000")

    for tb in (tb1, tb2):
        sp = tb._element
        parent = sp.getparent()
        sp_copy = deepcopy(sp)
        parent.remove(sp)
        grpSp.append(sp_copy)

    prs.save(str(pptx_path))


@pytest.fixture
def template_with_group(tmp_path: Path) -> Path:
    """Template where two text boxes are grouped via a <p:grpSp> element."""
    path = tmp_path / "with_group.pptx"
    _build_template_with_group(path)
    return path


def _write_config(
    config_path: Path,
    template_path: Path,
    placeholders: dict[str, str],
    **extra: Any,
) -> None:
    """Write a TemplateConfig YAML file with the given mapping."""
    data: dict[str, Any] = {
        "template": str(template_path),
        "source_slide_index": 1,
        "clear_source_slide": True,
        "placeholders": placeholders,
    }
    data.update(extra)
    config_path.write_text(yaml.safe_dump(data), encoding="utf-8")


@pytest.fixture
def write_config():
    """Factory fixture that writes a config YAML and returns its path."""
    return _write_config


def _build_sample_data_dir(data_dir: Path) -> None:
    """Two YAML records — alpha and bravo — matching rich_template fields."""
    data_dir.mkdir(parents=True, exist_ok=True)
    (data_dir / "alpha.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "alpha",
                "name": "Alpha Jones",
                "role": "Senior Consultant",
                "summary": "Twelve years across banking analytics and governance.",
                "background": [
                    "Bank A (2010-2015)",
                    "Bank B (2015-2020)",
                    "Bank C (2020-present)",
                ],
                "key_skills": ["Risk modelling", "Governance", "Data architecture"],
            }
        ),
        encoding="utf-8",
    )
    (data_dir / "bravo.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "bravo",
                "name": "Bravo Lee",
                "role": "Principal Consultant",
                "summary": "AI strategy lead for regulated industries.",
                "background": [
                    "Consulting firm A (2012-2018)",
                    "Consulting firm B (2018-present)",
                ],
                "key_skills": ["AI strategy", "Regulation"],
            }
        ),
        encoding="utf-8",
    )


@pytest.fixture
def sample_data_dir(tmp_path: Path) -> Path:
    """A populated data directory with alpha.yaml and bravo.yaml."""
    data_dir = tmp_path / "data"
    _build_sample_data_dir(data_dir)
    return data_dir
