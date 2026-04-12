"""Tests for v0.2.0 edge case hardening: nested dict in table cell, empty deck
warning, and txBody zero-items guard."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from recombinase.config import TableConfig, TemplateConfig
from recombinase.generate import generate_deck, populate_table

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _build_two_col_table_template(pptx_path: Path) -> None:
    """Template with a 2-row x 2-col table (row 0 = header, row 1 = data)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(0.5), Inches(5), Inches(2))
    table_shape.name = "Test_Table"
    tbl = table_shape.table
    tbl.cell(0, 0).text = "Col A"
    tbl.cell(0, 1).text = "Col B"
    tbl.cell(1, 0).text = "example a"
    tbl.cell(1, 1).text = "example b"
    prs.save(str(pptx_path))


def _build_simple_textbox_template(pptx_path: Path) -> None:
    """Template with one named text box — minimal target for generate_deck."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tb.name = "Field"
    tb.text_frame.text = "EXAMPLE"
    prs.save(str(pptx_path))


# ---------------------------------------------------------------------------
# 6a: nested dict in table cell
# ---------------------------------------------------------------------------


def test_nested_dict_in_table_cell_warns_and_clears(tmp_path: Path) -> None:
    """populate_table emits a warning and clears the cell when a value is a dict."""
    pptx_path = tmp_path / "template.pptx"
    _build_two_col_table_template(pptx_path)

    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    table_shape = next(s for s in slide.shapes if s.name == "Test_Table")

    table_config = TableConfig(
        shape="Test_Table",
        columns=["Col A", "Col B"],
        header_row=True,
    )

    rows = [{"Col A": "normal value", "Col B": {"nested": "dict"}}]
    warnings = populate_table(table_shape, table_config, rows)

    assert any("value is a dict" in w for w in warnings), (
        f"Expected 'value is a dict' warning, got: {warnings}"
    )

    # The cell that held the dict value should now be empty
    data_row = table_shape.table.rows[1]
    col_b_cell = data_row.cells[1]
    assert col_b_cell.text_frame.text == "", (
        f"Expected cleared cell, got: {col_b_cell.text_frame.text!r}"
    )


# ---------------------------------------------------------------------------
# 6b: empty deck warning
# ---------------------------------------------------------------------------


def test_empty_deck_warning(tmp_path: Path) -> None:
    """generate_deck warns when 0 records are generated with clear_source_slide=True."""
    pptx_path = tmp_path / "template.pptx"
    _build_simple_textbox_template(pptx_path)
    output_path = tmp_path / "output.pptx"

    config = TemplateConfig(
        template=pptx_path,
        source_slide_index=1,
        clear_source_slide=True,
        placeholders={"field": "Field"},
    )

    result = generate_deck(config, [], output_path)

    assert any("output deck will be empty" in w for w in result["warnings"]), (
        f"Expected 'output deck will be empty' warning, got: {result['warnings']}"
    )
    assert result["records_generated"] == 0
