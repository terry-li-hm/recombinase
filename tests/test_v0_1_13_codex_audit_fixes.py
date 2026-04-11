"""v0.1.13: Codex audit fix-ups.

Covers the seven findings from the 2026-04-11 Codex audit:

HIGH
----
1. populate_table clears cells when a row dict is missing a column key
   (previously the duplicated example text leaked into the output).
2. populate_table clears excess data rows when the record has fewer rows
   than the template table can hold (previously example rows leaked).

MED
---
3. cli.validate honours config.tables — table shape typos are now caught.
4. load_records stamps `_recombinase_record_dir` so relative picture paths
   resolve against the YAML file's directory instead of CWD.
5. write_scaffold_config uses yaml.safe_dump so YAML-significant characters
   in shape names are quoted correctly.

LOW
---
6. load_config distinguishes None from wrong-type empty config sections.
7. cli.generate --dry-run uses TemporaryDirectory so the temp dir is
   cleaned up even when generate_deck raises.
"""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches

from recombinase.config import TableConfig, load_config, write_scaffold_config
from recombinase.generate import (
    find_shape_by_name,
    load_records,
    populate_table,
)

# -- populate_table HIGH fixes ---------------------------------------------


def _build_template_with_table(path: Path, data_rows: int = 2) -> None:
    """Template with a named table: 1 header row + N data rows, all pre-filled."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    total_rows = data_rows + 1
    table_shape = slide.shapes.add_table(
        rows=total_rows,
        cols=3,
        left=Inches(0.5),
        top=Inches(0.5),
        width=Inches(8),
        height=Inches(2),
    )
    table_shape.name = "projects"
    table = table_shape.table
    table.cell(0, 0).text = "Client"
    table.cell(0, 1).text = "Role"
    table.cell(0, 2).text = "Outcome"
    for row_index in range(1, total_rows):
        table.cell(row_index, 0).text = f"EXAMPLE CLIENT {row_index}"
        table.cell(row_index, 1).text = f"EXAMPLE ROLE {row_index}"
        table.cell(row_index, 2).text = f"EXAMPLE OUTCOME {row_index}"
    prs.save(str(path))


def _load_table_shape(template_path: Path):
    prs = Presentation(str(template_path))
    shape = find_shape_by_name(prs.slides[0], "projects")
    assert shape is not None
    return prs, shape


def test_missing_column_key_clears_cell_and_warns(tmp_path: Path) -> None:
    """Finding 1: a row dict missing a column key used to leak example text."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    # Row 0 is missing 'role'; row 1 is complete
    rows = [
        {"client": "Bank A", "outcome": "Delivered"},
        {"client": "Bank B", "role": "Lead", "outcome": "Designed"},
    ]
    warnings = populate_table(shape, config, rows)

    table = shape.table
    # The missing column should be CLEARED, not left with 'EXAMPLE ROLE 1'
    assert table.cell(1, 0).text == "Bank A"
    assert table.cell(1, 1).text == ""
    assert table.cell(1, 2).text == "Delivered"
    # Row 1 is fully populated
    assert table.cell(2, 0).text == "Bank B"
    assert table.cell(2, 1).text == "Lead"
    assert table.cell(2, 2).text == "Designed"
    # And the user is warned about the missing column
    assert any("missing column 'role'" in w for w in warnings)


def test_none_value_clears_cell(tmp_path: Path) -> None:
    """Finding 1b: explicit None should clear, not skip."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    rows = [
        {"client": "Bank A", "role": None, "outcome": "Delivered"},
        {"client": "Bank B", "role": "Lead", "outcome": "Designed"},
    ]
    populate_table(shape, config, rows)
    table = shape.table
    assert table.cell(1, 1).text == ""  # cleared, not "EXAMPLE ROLE 1"


def test_empty_string_value_clears_cell(tmp_path: Path) -> None:
    """Finding 1c: empty string should clear, not skip."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    rows = [
        {"client": "Bank A", "role": "", "outcome": "Delivered"},
        {"client": "Bank B", "role": "Lead", "outcome": "Designed"},
    ]
    populate_table(shape, config, rows)
    table = shape.table
    assert table.cell(1, 1).text == ""


def test_short_record_clears_excess_data_rows(tmp_path: Path) -> None:
    """Finding 2: unused data rows used to leak example text."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=3)  # header + 3 data rows
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    rows = [
        {"client": "Bank A", "role": "Lead", "outcome": "Delivered"},
    ]
    warnings = populate_table(shape, config, rows)

    table = shape.table
    # Row 1 populated
    assert table.cell(1, 0).text == "Bank A"
    # Rows 2 and 3 cleared — no leaked EXAMPLE text
    assert table.cell(2, 0).text == ""
    assert table.cell(2, 1).text == ""
    assert table.cell(2, 2).text == ""
    assert table.cell(3, 0).text == ""
    assert table.cell(3, 1).text == ""
    assert table.cell(3, 2).text == ""
    assert any("clearing 2 unused row" in w for w in warnings)


def test_combined_missing_column_and_short_record(tmp_path: Path) -> None:
    """Finding 1 + 2: both bugs at once."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=3)
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    # One record, missing 'outcome'
    rows = [{"client": "Bank A", "role": "Lead"}]
    warnings = populate_table(shape, config, rows)

    table = shape.table
    assert table.cell(1, 0).text == "Bank A"
    assert table.cell(1, 1).text == "Lead"
    assert table.cell(1, 2).text == ""  # missing column cleared
    # Excess rows cleared
    for row_index in (2, 3):
        for col_index in range(3):
            assert table.cell(row_index, col_index).text == ""
    # Both warnings present
    assert any("missing column 'outcome'" in w for w in warnings)
    assert any("clearing 2 unused row" in w for w in warnings)


def test_merged_cells_do_not_crash_clear(tmp_path: Path) -> None:
    """Gemini finding (post-Codex fix): _clear_cell must skip spanned cells.

    python-pptx raises on `cell.text_frame` for spanned (merged non-origin)
    cells. The Codex audit fix added three unconditional _clear_cell loops —
    any template with merged cells in the data region would crash with
    ValueError before this guard.

    Build a table where two adjacent cells in an excess-row region are merged,
    then provide a short record. The clear-excess-rows path must succeed
    without raising; only the merge origin's text needs to go to empty.
    """
    template = tmp_path / "merged.pptx"
    _build_template_with_table(template, data_rows=3)

    # Merge two cells in row 2 (an excess row we'll clear) BEFORE saving.
    prs = Presentation(str(template))
    shape = find_shape_by_name(prs.slides[0], "projects")
    assert shape is not None
    origin = shape.table.cell(2, 0)
    other = shape.table.cell(2, 1)
    origin.merge(other)
    prs.save(str(template))

    _, shape = _load_table_shape(template)
    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    # One record → rows 2 and 3 are excess → the merged cells get cleared
    rows = [{"client": "Bank A", "role": "Lead", "outcome": "Delivered"}]

    # Must NOT raise. Before the is_spanned guard, this raised ValueError
    # on cell.text_frame access for the spanned non-origin cell.
    warnings = populate_table(shape, config, rows)

    table = shape.table
    # Populated row is intact
    assert table.cell(1, 0).text == "Bank A"
    # Excess row 3 fully cleared (non-merged)
    assert table.cell(3, 0).text == ""
    assert table.cell(3, 1).text == ""
    assert table.cell(3, 2).text == ""
    # Merge origin in row 2 is cleared (spanned sibling skipped silently)
    assert table.cell(2, 0).text == ""
    assert any("clearing 2 unused row" in w for w in warnings)


def test_fully_populated_still_works_no_extra_warnings(tmp_path: Path) -> None:
    """Regression: a fully populated table still succeeds with no warnings."""
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)
    _, shape = _load_table_shape(template)

    config = TableConfig(
        shape="projects",
        columns=["client", "role", "outcome"],
        header_row=True,
    )
    rows = [
        {"client": "Bank A", "role": "Lead", "outcome": "Delivered"},
        {"client": "Bank B", "role": "Advisor", "outcome": "Designed"},
    ]
    warnings = populate_table(shape, config, rows)
    assert warnings == []
    table = shape.table
    assert table.cell(1, 0).text == "Bank A"
    assert table.cell(2, 0).text == "Bank B"


# -- MED: cli validate covers table shapes ---------------------------------


def test_validate_flags_missing_table_shape(tmp_path: Path) -> None:
    """Finding 3: typos in table.shape must be flagged by `recombinase validate`."""
    from typer.testing import CliRunner

    from recombinase.cli import app

    # Build a template with the table named 'projects'
    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)

    # Config references a DIFFERENT table shape name (typo)
    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "placeholders": {},
                "tables": {
                    "projects": {
                        "shape": "projets",  # typo
                        "columns": ["client", "role", "outcome"],
                        "header_row": True,
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    runner = CliRunner()
    result = runner.invoke(app, ["validate", "--config", str(config_path)])
    assert result.exit_code == 1
    assert "projets" in result.output
    assert "table field" in result.output


def test_validate_accepts_matching_table_shape(tmp_path: Path) -> None:
    """Finding 3b: correct table shape names validate cleanly."""
    from typer.testing import CliRunner

    from recombinase.cli import app

    template = tmp_path / "t.pptx"
    _build_template_with_table(template, data_rows=2)

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "placeholders": {},
                "tables": {
                    "projects": {
                        "shape": "projects",
                        "columns": ["client", "role", "outcome"],
                        "header_row": True,
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    runner = CliRunner()
    result = runner.invoke(app, ["validate", "--config", str(config_path)])
    assert result.exit_code == 0
    assert "valid" in result.output.lower()


# -- MED: load_records stamps record dir -----------------------------------


def test_load_records_stamps_record_dir(tmp_path: Path) -> None:
    """Finding 4: relative picture paths must resolve against the YAML dir."""
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r1.yaml").write_text(
        yaml.safe_dump({"id": "r1", "name": "Alpha"}),
        encoding="utf-8",
    )
    records = load_records(data_dir)
    assert len(records) == 1
    record = records[0]
    assert record["_recombinase_record_dir"] == str(data_dir)
    # Sanity: stays attached through a Path() conversion, the same way
    # generate_deck uses it to build base_dir for set_picture.
    base_dir = Path(record["_recombinase_record_dir"])
    assert base_dir == data_dir


# -- MED: write_scaffold_config uses yaml.safe_dump ------------------------


def test_scaffold_config_quotes_special_characters(tmp_path: Path) -> None:
    """Finding 5: shape names with YAML-significant chars used to break the scaffold."""
    template_path = tmp_path / "t.pptx"
    # create a minimal dummy so write_scaffold_config has a real path string
    Presentation().save(str(template_path))
    output_path = tmp_path / "scaffold.yaml"

    # Shape names containing YAML-significant characters
    shape_names = [
        "Name: Title",  # colon-space — yaml mapping indicator
        "Role #1",  # hash — comment indicator
        "Summary {short}",  # braces — flow mapping
        "List [bullets]",  # brackets — flow sequence
    ]
    write_scaffold_config(template_path, shape_names, output_path)

    # The file must parse cleanly through load_config
    cfg = load_config(output_path)
    # Every shape name should be present in the round-tripped placeholders
    round_tripped = set(cfg.placeholders.values())
    for name in shape_names:
        assert name in round_tripped, f"shape {name!r} lost in scaffold round-trip"


# -- LOW: load_config distinguishes None from wrong-type ------------------


def test_load_config_rejects_wrong_type_placeholders(tmp_path: Path) -> None:
    """Finding 6: a list where a mapping is expected used to be silently accepted."""
    template_path = tmp_path / "t.pptx"
    Presentation().save(str(template_path))

    config_path = tmp_path / "c.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "placeholders": [],  # wrong type — used to short-circuit to {}
                "tables": {
                    "projects": {"shape": "projects", "columns": ["a", "b"]},
                },
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match="'placeholders' must be a mapping"):
        load_config(config_path)


def test_load_config_accepts_null_placeholders(tmp_path: Path) -> None:
    """Finding 6b: explicit null should still be treated as empty."""
    template_path = tmp_path / "t.pptx"
    Presentation().save(str(template_path))

    config_path = tmp_path / "c.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "placeholders": None,
                "tables": {
                    "projects": {"shape": "projects", "columns": ["a", "b"]},
                },
            }
        ),
        encoding="utf-8",
    )

    cfg = load_config(config_path)
    assert cfg.placeholders == {}
    assert "projects" in cfg.tables


def test_load_config_rejects_wrong_type_table_columns(tmp_path: Path) -> None:
    """Finding 6c: a dict where a list is expected for columns must be rejected."""
    template_path = tmp_path / "t.pptx"
    Presentation().save(str(template_path))

    config_path = tmp_path / "c.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "placeholders": {"name": "Field"},
                "tables": {
                    "projects": {"shape": "projects", "columns": {"a": 1}},
                },
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match="must be a list of strings"):
        load_config(config_path)


# -- LOW: dry-run tempdir is cleaned up on exception -----------------------


def test_dry_run_uses_context_manager_for_tempdir(monkeypatch, tmp_path: Path) -> None:
    """Finding 7: if generate_deck raises, the tempdir must still be cleaned up.

    We patch generate_deck to raise and confirm TemporaryDirectory semantics —
    namely that after the command returns, no new permanent directory is
    lingering inside the system tempdir. We can't easily observe the exact
    tempdir name, so we rely on the context-manager guarantee: the command's
    own `with` block has exited by the time we inspect state.
    """
    import tempfile

    from typer.testing import CliRunner

    from recombinase import cli as cli_mod
    from recombinase.cli import app

    # Minimal scaffold so the command gets as far as calling generate_deck
    template_path = tmp_path / "t.pptx"
    _build_template_with_table(template_path, data_rows=2)

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r1.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "r1",
                "projects": [
                    {"client": "A", "role": "Lead", "outcome": "Delivered"},
                ],
            }
        ),
        encoding="utf-8",
    )

    config_path = tmp_path / "config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            {
                "template": str(template_path),
                "source_slide_index": 1,
                "clear_source_slide": True,
                "overflow_ratio": 0,
                "placeholders": {},
                "tables": {
                    "projects": {
                        "shape": "projects",
                        "columns": ["client", "role", "outcome"],
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    captured_tempdirs: list[str] = []
    original_temp_dir_cls = tempfile.TemporaryDirectory

    class TrackingTempDir(original_temp_dir_cls):  # type: ignore[misc, valid-type]
        def __enter__(self) -> str:
            path = super().__enter__()
            captured_tempdirs.append(path)
            return path

    monkeypatch.setattr(tempfile, "TemporaryDirectory", TrackingTempDir)

    def boom(*_args: object, **_kwargs: object) -> None:
        raise RuntimeError("synthetic failure mid-generate")

    monkeypatch.setattr(cli_mod, "generate_deck", boom)

    runner = CliRunner()
    # Writing to a temp output path inside tmp_path so the command doesn't
    # bail on "file exists".
    output_path = tmp_path / "out.pptx"
    result = runner.invoke(
        app,
        [
            "generate",
            "--config",
            str(config_path),
            "--data-dir",
            str(data_dir),
            "--output",
            str(output_path),
            "--dry-run",
        ],
    )
    # The command should have propagated the RuntimeError — but the
    # TemporaryDirectory must have been cleaned up regardless.
    assert result.exit_code != 0
    assert captured_tempdirs, "dry-run did not enter a TemporaryDirectory context"
    for tempdir in captured_tempdirs:
        assert not Path(tempdir).exists(), (
            f"dry-run tempdir {tempdir} leaked after generate_deck raised"
        )
