"""Tests for v0.2.0 config validation hardening: unknown keys + duplicate shapes."""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches

from recombinase.config import load_config


def _build_minimal_template(path: Path) -> None:
    """Minimal pptx: one slide with one named textbox and one table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    tb.name = "Name_Field"
    tb.text_frame.text = "EXAMPLE"
    table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2.0), Inches(4), Inches(2))
    table_shape.name = "Data_Table"
    prs.save(str(path))


class TestUnknownTopLevelKeyWarns:
    """load_config warns on unrecognized top-level keys (e.g. typos)."""

    def test_unknown_top_level_key_warns(self, tmp_path: Path) -> None:
        """A YAML with a typo key 'placeholder:' produces a UserWarning."""
        template_path = tmp_path / "template.pptx"
        _build_minimal_template(template_path)

        config_data = {
            "template": str(template_path),
            "source_slide_index": 1,
            "placeholders": {"name": "Name_Field"},
            # 'placeholder' is a common typo for 'placeholders'
            "placeholder": {"name": "Name_Field"},
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        with pytest.warns(UserWarning, match="unrecognized"):
            load_config(config_path)


class TestUnknownTableSubKeyWarns:
    """load_config warns on unrecognized sub-keys inside a table entry."""

    def test_unknown_table_sub_key_warns(self, tmp_path: Path) -> None:
        """A table entry with unknown sub-key 'colour:' produces a UserWarning."""
        template_path = tmp_path / "template.pptx"
        _build_minimal_template(template_path)

        config_data = {
            "template": str(template_path),
            "source_slide_index": 1,
            "tables": {
                "data": {
                    "shape": "Data_Table",
                    "columns": ["col_a", "col_b"],
                    "header_row": True,
                    "colour": "blue",  # unknown key
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        with pytest.warns(UserWarning, match="unrecognized"):
            load_config(config_path)


class TestUnknownSectionSubKeyWarns:
    """load_config warns on unrecognized sub-keys inside a section entry."""

    def test_unknown_section_sub_key_warns(self, tmp_path: Path) -> None:
        """A section entry with unknown sub-key 'style:' produces a UserWarning."""
        template_path = tmp_path / "template.pptx"
        _build_minimal_template(template_path)

        # Need a textbox with >=2 paragraphs for the section shape
        prs = Presentation(str(template_path))
        slide = prs.slides[0]
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4), Inches(2))
        tb.name = "Section_Shape"
        tf = tb.text_frame
        tf.text = "Header"
        tf.add_paragraph().text = "Bullet"
        prs.save(str(template_path))

        config_data = {
            "template": str(template_path),
            "source_slide_index": 1,
            "sections": {
                "skills": {
                    "shape": "Section_Shape",
                    "header_index": 0,
                    "bullet_index": 1,
                    "style": "fancy",  # unknown key
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        with pytest.warns(UserWarning, match="unrecognized"):
            load_config(config_path)


class TestDuplicateShapeAcrossSections:
    """load_config raises ValueError when the same shape name appears in multiple sections."""

    def test_duplicate_shape_across_placeholders_and_tables(self, tmp_path: Path) -> None:
        """Same shape name in placeholders and tables raises ValueError."""
        template_path = tmp_path / "template.pptx"
        _build_minimal_template(template_path)

        config_data = {
            "template": str(template_path),
            "source_slide_index": 1,
            "placeholders": {"name": "Data_Table"},  # reusing shape of the table
            "tables": {
                "data": {
                    "shape": "Data_Table",  # duplicate
                    "columns": ["col_a", "col_b"],
                }
            },
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        with pytest.raises(ValueError, match="appears in multiple config sections"):
            load_config(config_path)


class TestValidConfigNoWarnings:
    """A well-formed config produces no warnings."""

    def test_valid_config_no_warnings(self, tmp_path: Path) -> None:
        """A clean config with known keys only produces no 'unrecognized' warnings."""
        import warnings

        template_path = tmp_path / "template.pptx"
        _build_minimal_template(template_path)

        config_data = {
            "template": str(template_path),
            "source_slide_index": 1,
            "clear_source_slide": True,
            "overflow_ratio": 1.5,
            "placeholders": {"name": "Name_Field"},
        }
        config_path = tmp_path / "config.yaml"
        config_path.write_text(yaml.safe_dump(config_data), encoding="utf-8")

        with warnings.catch_warnings(record=True) as captured:
            warnings.simplefilter("always")
            load_config(config_path)

        # No 'unrecognized' warnings from our validation code
        our_warnings = [w for w in captured if "unrecognized" in str(w.message).lower()]
        assert len(our_warnings) == 0
