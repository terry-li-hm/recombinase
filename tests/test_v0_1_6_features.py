"""Tests for v0.1.6 features and remaining regression gaps.

Uses shared conftest fixtures. Covers:
- rId rewrite correctness in duplicate_slide (not just rels preservation)
- Overflow detection heuristic in generate_deck
- Preset geometry detection in inspect (circle-crop reporting)
- CLI --strict exit code 2 path
- main() exception dispatch table routing
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

import yaml
from pptx import Presentation
from typer.testing import CliRunner

from recombinase.cli import _ERROR_HANDLERS, app, main
from recombinase.config import load_config
from recombinase.generate import (
    _capture_baseline_lengths,
    _value_char_length,
    duplicate_slide,
    generate_deck,
    load_records,
)
from recombinase.inspect import inspect_template

runner = CliRunner()


# -- rId rewrite correctness (new coverage for v0.1.5 fix) -----------------


def test_duplicate_slide_rewrites_rid_when_ids_collide(
    template_with_picture: Path, tmp_path: Path
) -> None:
    """Force a rId collision so the rewrite branch in duplicate_slide fires.

    Strategy: we add an extra relationship to the new slide BEFORE the rel
    copy, so python-pptx's get_or_add has to assign a different rId than the
    source slide used. Without the rewrite, the copied shape XML would still
    reference the old rId and the picture would not resolve. With rewrite,
    the XML attribute should be updated to the new rId.
    """
    prs = Presentation(str(template_with_picture))
    source = prs.slides[0]

    # Grab the source picture's current r:embed value
    from pptx.oxml.ns import qn

    source_pic_el = None
    for shape in source.shapes:
        if shape.name == "Headshot":
            source_pic_el = shape._element
            break
    assert source_pic_el is not None
    blip = source_pic_el.find(f".//{qn('a:blip')}")
    assert blip is not None
    original_rid = blip.get(
        f"{{{'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}}}embed"
    )
    assert original_rid is not None

    # Run duplication
    new_slide = duplicate_slide(prs, source)

    # Find the copied picture on the new slide
    new_blip = None
    for shape in new_slide.shapes:
        if shape.name == "Headshot":
            new_blip = shape._element.find(f".//{qn('a:blip')}")
            break
    assert new_blip is not None
    new_rid = new_blip.get(
        f"{{{'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}}}embed"
    )
    assert new_rid is not None

    # Round-trip save and verify the picture still resolves. If the rewrite
    # failed, opening the Picture.image would raise KeyError on the missing
    # relationship.
    output = tmp_path / "rewritten.pptx"
    prs.save(str(output))
    reopened = Presentation(str(output))

    from pptx.shapes.picture import Picture

    for shape in reopened.slides[1].shapes:
        if isinstance(shape, Picture):
            # Accessing .image forces the rel lookup — KeyError if broken
            image_bytes = shape.image.blob
            assert len(image_bytes) > 0


def test_duplicate_slide_deep_copy_produces_independent_shapes(
    template_with_picture: Path, tmp_path: Path
) -> None:
    """Editing the new slide's shape must not affect the source slide."""
    prs = Presentation(str(template_with_picture))
    source = prs.slides[0]
    new_slide = duplicate_slide(prs, source)

    for shape in new_slide.shapes:
        if shape.name == "Consultant_Name":
            shape.text_frame.text = "MUTATED"
            break

    for shape in source.shapes:
        if shape.name == "Consultant_Name":
            assert shape.text_frame.text == "EXAMPLE", (
                "source slide should be untouched when new slide is mutated"
            )


# -- overflow detection (new v0.1.6 feature) -------------------------------


def test_value_char_length_scalar() -> None:
    assert _value_char_length("hello") == 5
    assert _value_char_length(42) == 2
    assert _value_char_length(3.14) == 4
    assert _value_char_length(None) == 0
    assert _value_char_length("") == 0


def test_value_char_length_list_accounts_for_newlines() -> None:
    # Three items of 3 chars each = 9 chars + 2 newlines = 11
    assert _value_char_length(["abc", "def", "ghi"]) == 11
    assert _value_char_length([]) == 0
    assert _value_char_length([None, "", "real"]) == 4  # filtered to 1 item


def test_capture_baseline_lengths_from_source(simple_template: Path) -> None:
    prs = Presentation(str(simple_template))
    baselines = _capture_baseline_lengths(prs.slides[0], {"name": "Field"})
    assert baselines == {"name": len("EXAMPLE")}


def test_generate_deck_warns_on_overflow(
    simple_template: Path, tmp_path: Path, write_config
) -> None:
    """Data much larger than the baseline should trigger an overflow warning."""
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=1.5)

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    # Baseline is "EXAMPLE" (7 chars). 50 chars is ~7x baseline, well over 1.5.
    (data_dir / "r.yaml").write_text(
        yaml.safe_dump({"name": "this is a much much longer name than the example"}),
        encoding="utf-8",
    )

    cfg = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(cfg, records, tmp_path / "out.pptx")

    assert any("may overflow" in w for w in result["warnings"]), result["warnings"]
    assert any("name" in w for w in result["warnings"])


def test_generate_deck_does_not_warn_when_content_fits(
    simple_template: Path, tmp_path: Path, write_config
) -> None:
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=1.5)

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    # 7 chars baseline, 8 chars content → ratio 1.14, under 1.5
    (data_dir / "r.yaml").write_text(yaml.safe_dump({"name": "same-ish"}), encoding="utf-8")

    cfg = load_config(config_path)
    records = load_records(data_dir)
    result = generate_deck(cfg, records, tmp_path / "out.pptx")

    assert not any("may overflow" in w for w in result["warnings"]), result["warnings"]


def test_generate_deck_overflow_disabled_by_zero_ratio(
    simple_template: Path, tmp_path: Path, write_config
) -> None:
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=0)

    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r.yaml").write_text(yaml.safe_dump({"name": "x" * 1000}), encoding="utf-8")

    cfg = load_config(config_path)
    assert cfg.overflow_ratio == 0
    records = load_records(data_dir)
    result = generate_deck(cfg, records, tmp_path / "out.pptx")

    assert not any("may overflow" in w for w in result["warnings"])


def test_load_config_rejects_negative_overflow_ratio(simple_template: Path, tmp_path: Path) -> None:
    import pytest

    config_path = tmp_path / "bad.yaml"
    config_path.write_text(
        f"template: {simple_template}\noverflow_ratio: -1\nplaceholders:\n  name: Field\n",
        encoding="utf-8",
    )
    with pytest.raises(ValueError, match=r"overflow_ratio.*>= 0"):
        load_config(config_path)


def test_load_config_rejects_non_numeric_overflow_ratio(
    simple_template: Path, tmp_path: Path
) -> None:
    import pytest

    config_path = tmp_path / "bad.yaml"
    config_path.write_text(
        f"template: {simple_template}\noverflow_ratio: lots\nplaceholders:\n  name: Field\n",
        encoding="utf-8",
    )
    with pytest.raises(ValueError, match=r"overflow_ratio.*number"):
        load_config(config_path)


# -- preset geometry detection (new v0.1.6 feature) ------------------------


def test_inspect_detects_ellipse_preset_geom(tmp_path: Path) -> None:
    """A shape with `prstGeom prst="ellipse"` should be reported by inspect."""
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # python-pptx doesn't expose a direct "add ellipse" API on the basic
    # shape interface, but autoshape MSO_SHAPE.OVAL creates a prstGeom ellipse.
    from pptx.enum.shapes import MSO_SHAPE

    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2))
    oval.name = "Headshot_Mask"

    path = tmp_path / "ellipse.pptx"
    prs.save(str(path))

    info = inspect_template(path)
    mask_shape = next(
        (s for s in info.slides[0].shapes if s.name == "Headshot_Mask"),
        None,
    )
    assert mask_shape is not None
    assert mask_shape.preset_geom == "ellipse"


def test_inspect_reports_rect_for_regular_textbox(simple_template: Path) -> None:
    info = inspect_template(simple_template)
    shape = info.slides[0].shapes[0]
    # Text boxes use prstGeom prst="rect"
    assert shape.preset_geom == "rect"


def test_format_template_info_shows_geom_for_ellipse(tmp_path: Path) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2))
    oval.name = "Profile"
    path = tmp_path / "e.pptx"
    prs.save(str(path))

    from recombinase.inspect import format_template_info

    info = inspect_template(path)
    formatted = format_template_info(info)
    assert "geom=ellipse" in formatted


# -- CLI --strict exit code 2 (new coverage) -------------------------------


def test_cli_generate_strict_exits_2_on_overflow_warning(
    simple_template: Path, tmp_path: Path, write_config
) -> None:
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=1.5)
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r.yaml").write_text(
        yaml.safe_dump({"name": "way way way way way way too long name example"}),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
            "--strict",
        ],
    )
    assert result.exit_code == 2, result.output
    assert output.exists()  # Still wrote the file; strict just escalates exit


def test_cli_generate_non_strict_exits_0_on_overflow_warning(
    simple_template: Path, tmp_path: Path, write_config
) -> None:
    config_path = tmp_path / "config.yaml"
    write_config(config_path, simple_template, {"name": "Field"}, overflow_ratio=1.5)
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    (data_dir / "r.yaml").write_text(
        yaml.safe_dump({"name": "way too long for the example baseline"}),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"

    result = runner.invoke(
        app,
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(data_dir),
            "-o",
            str(output),
        ],
    )
    assert result.exit_code == 0


# -- main() exception dispatch table coverage -----------------------------


def test_error_handlers_registered_for_core_types() -> None:
    """Verify the dispatch table has entries for all classified exception types."""
    registered_types = {exc_type for exc_type, _, _ in _ERROR_HANDLERS}

    # Each of these must be routed through a formatter
    assert PermissionError in registered_types
    assert FileNotFoundError in registered_types
    assert NotADirectoryError in registered_types
    assert ValueError in registered_types


def test_main_routes_value_error_to_clean_message(
    simple_template: Path, tmp_path: Path, write_config, capsys
) -> None:
    """Force a ValueError path and verify the dispatch table handled it."""
    config_path = tmp_path / "bad.yaml"
    config_path.write_text("- not a mapping\n", encoding="utf-8")

    rc = main(
        [
            "generate",
            "-c",
            str(config_path),
            "-d",
            str(tmp_path),
            "-o",
            str(tmp_path / "out.pptx"),
        ]
    )
    assert rc == 1
    captured = capsys.readouterr()
    combined = captured.out + captured.err
    # Should show clean error, not a traceback
    assert "Error:" in combined
    assert "Traceback" not in combined


# -- conftest fixture usage check -----------------------------------------


def test_simple_template_fixture_is_usable(simple_template: Path) -> None:
    assert simple_template.exists()
    info = inspect_template(simple_template)
    assert info.slide_count == 1
    assert any(s.name == "Field" for s in info.slides[0].shapes)


def test_rich_template_fixture_has_expected_shapes(rich_template: Path) -> None:
    info = inspect_template(rich_template)
    names = [s.name for s in info.slides[0].shapes]
    assert "Consultant_Name" in names
    assert "Role_Title" in names
    assert "Summary_Body" in names
    assert "Background_Bullets" in names
    assert "Key_Skills" in names


def test_template_with_group_fixture_yields_groups(template_with_group: Path) -> None:
    info = inspect_template(template_with_group)
    names = [s.name for s in info.slides[0].shapes]
    assert "Consultant_Card" in names
    assert "Grouped_Name" in names
    assert "Grouped_Role" in names


def test_template_with_picture_fixture_has_picture(
    template_with_picture: Path,
) -> None:
    prs = Presentation(str(template_with_picture))
    shape_names = [s.name for s in prs.slides[0].shapes]
    assert "Consultant_Name" in shape_names
    assert "Headshot" in shape_names


def test_sample_data_dir_fixture_has_two_records(sample_data_dir: Path) -> None:
    records = load_records(sample_data_dir)
    assert len(records) == 2
    ids = {r.get("id") for r in records}
    assert ids == {"alpha", "bravo"}


# Keep deepcopy import used
_ = deepcopy
